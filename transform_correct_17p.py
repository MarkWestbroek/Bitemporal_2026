#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt

pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(pptx_file)
print(f'Totaal slides: {len(prs.slides)}\n')

# Template uit slide 3
slide3 = prs.slides[2]
title_template = slide3.shapes[0]
content_template = slide3.shapes[1]

print(f'Template: Titel @ ({title_template.left}, {title_template.top})')
print(f'Template: Content @ ({content_template.left}, {content_template.top})\n')

# Transformeer slides 5-17
for slide_idx in range(4, 17):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    # Verzamel alle content
    title_text = None
    all_bullets = []
    intro_line = None
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        
        text = shape.text.strip()
        if not text:
            continue
        
        # Font size bepalen
        font_size = None
        if len(shape.text_frame.paragraphs) > 0:
            para = shape.text_frame.paragraphs[0]
            if len(para.runs) > 0:
                font_size = para.runs[0].font.size
        
        # Titel = grote font (360000+)
        if font_size and font_size > 200000 and not title_text:
            title_text = text
        # Rest = bullets
        elif text != '→' and font_size and font_size < 200000:
            all_bullets.append(text)
    
    if not title_text:
        print(f'Slide {slide_num}: SKIP (no title)')
        continue
    
    # Intro is de eerste bullet
    if all_bullets:
        intro_line = all_bullets[0]
        rest_bullets = all_bullets[1:]
    else:
        intro_line = ''
        rest_bullets = []
    
    # Verwijder alle TEXT_BOX shapes (behoudt PICTURE)
    for shape in list(slide.shapes):
        if hasattr(shape, 'text_frame'):
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Voeg TITEL toe
    title_box = slide.shapes.add_textbox(
        title_template.left,
        title_template.top,
        title_template.width,
        title_template.height
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    
    # Voeg CONTENT toe (intro + bullets)
    content_box = slide.shapes.add_textbox(
        content_template.left,
        content_template.top,
        content_template.width,
        content_template.height
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.clear()
    
    # Eerste line: intro (24pt, bold)
    p = tf.paragraphs[0]
    p.text = intro_line
    p.font.size = Pt(24)  # 203200 EMU
    p.font.bold = True
    
    # Rest: bullets met arrow prefix (12pt, bold)
    for bullet in rest_bullets:
        # Empty line
        p = tf.add_paragraph()
        p.text = ''
        
        # Bullet line
        p = tf.add_paragraph()
        p.text = f'→ {bullet}'
        p.font.size = Pt(12)  # 161925 EMU
        p.font.bold = True
    
    print(f'Slide {slide_num}: ✓ "{title_text}" ({len(all_bullets)} items)')

prs.save(pptx_file)
print(f'\n✓ Klaar!')
