#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt

# Gebruik het herstelde 17p bestand
pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(pptx_file)
print(f'Totaal slides: {len(prs.slides)}\n')

# Analyse slide 3 en 4 zeer gedetailleerd
for slide_idx in [2, 3]:
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    print(f'=== SLIDE {slide_num} ===')
    print(f'Totaal shapes: {len(slide.shapes)}\n')
    
    for i, shape in enumerate(slide.shapes):
        print(f'Shape {i}:')
        print(f'  Type: {shape.shape_type}')
        print(f'  Naam: {shape.name}')
        
        if hasattr(shape, 'text_frame'):
            txt = shape.text
            print(f'  Text length: {len(txt)}')
            if len(txt) > 100:
                print(f'  Text (first 100): {txt[:100]}')
            else:
                print(f'  Text: {txt}')
            
            print(f'  Paragraphs: {len(shape.text_frame.paragraphs)}')
            
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                print(f'    Para {p_idx}:')
                print(f'      Runs: {len(para.runs)}')
                if len(para.runs) > 0:
                    run = para.runs[0]
                    print(f'      Font size: {run.font.size}')
                    print(f'      Font bold: {run.font.bold}')
                    print(f'      Text: {para.text[:40]}')
        
        print()
