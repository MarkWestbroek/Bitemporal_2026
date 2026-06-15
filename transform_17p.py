#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(pptx_file)
print(f'Totaal slides: {len(prs.slides)}\n')

# Analyseer slide 3 als template (het heeft al het juiste formaat)
slide3 = prs.slides[2]
print('=== TEMPLATE SLIDE 3 ===')

title_shape = None
content_shape = None

for shape in slide3.shapes:
    if hasattr(shape, 'text_frame'):
        text = shape.text.strip()
        if len(text) > 20:
            if 'Europese' in text:
                title_shape = shape
                print(f'Titel shape: "{text[:30]}"')
            elif 'Propriëtaire' in text:
                content_shape = shape
                print(f'Content shape gevonden (7 paragrafen)')

if not title_shape or not content_shape:
    print('ERROR: Could not find template shapes!')
    exit(1)

print(f'  Title pos: ({title_shape.left}, {title_shape.top}), size: {title_shape.width}x{title_shape.height}')
print(f'  Content pos: ({content_shape.left}, {content_shape.top}), size: {content_shape.width}x{content_shape.height}')
print()

# Transformeer slides 5-17
for slide_idx in range(4, 17):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    print(f'Slide {slide_num}:', end=' ')
    
    # Verzamel alle tekst
    title_text = None
    bullet_items = []
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        
        text = shape.text.strip()
        if not text:
            continue
        
        # Bepaal font size
        font_size = None
        if len(shape.text_frame.paragraphs) > 0:
            para = shape.text_frame.paragraphs[0]
            if len(para.runs) > 0:
                font_size = para.runs[0].font.size
        
        # Titel = grootste
        if font_size and font_size > 200000 and not title_text:
            title_text = text
        # Bullet = alles anderen, behalve arrows
        elif text != '→' and text != title_text:
            bullet_items.append(text)
    
    if not title_text:
        print('SKIP')
        continue
    
    # Verwijder alle text shapes
    for shape in list(slide.shapes):
        if hasattr(shape, 'text_frame'):
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Voeg titel toe (copy van template)
    title_box = slide.shapes.add_textbox(
        title_shape.left,
        title_shape.top,
        title_shape.width,
        title_shape.height
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    
    # Voeg content textbox toe met alle bullets
    content_box = slide.shapes.add_textbox(
        content_shape.left,
        content_shape.top,
        content_shape.width,
        content_shape.height
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.clear()
    
    # Eerste introductie line (groter, bold, 24pt)
    p = tf.paragraphs[0]
    p.text = bullet_items[0] if bullet_items else ''
    p.font.size = Pt(24)  # 203200 EMU
    p.font.bold = True
    
    # Rest als bullets met arrows
    for i, bullet in enumerate(bullet_items[1:], 1):
        p = tf.add_paragraph()
        p.text = ''  # Empty line
        
        # Bullet met arrow
        p = tf.add_paragraph()
        p.text = f'→ {bullet}'
        p.font.size = Pt(12)  # 161925 EMU
        p.font.bold = True
    
    print(f'✓ ({len(bullet_items)} items)')

prs.save(pptx_file)
print(f'\n✓ Transformatie compleet!')
