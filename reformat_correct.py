#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt, Inches

# Gebruik niet-17p versie als basis
src_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap.pptx'
dst_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(src_file)
print(f'Totaal slides in bron: {len(prs.slides)}\n')

# We gaan alleen slides 5-17 (index 4-16) aanpassen
# We willen ze hetzelfde formaat als slide 3 (index 2)

# Eerst: pak de template-positioning van slide 3 (index 2)
slide_template = prs.slides[2]
print(f'Slide 3 heeft {len(slide_template.shapes)} shapes')

# Verzamel de positioning van shapes in slide 3
title_shape_template = None
first_bullet_template = None

for shape in slide_template.shapes:
    if hasattr(shape, 'text_frame') and shape.text.strip():
        if len(shape.text_frame.paragraphs) > 0:
            para = shape.text_frame.paragraphs[0]
            if len(para.runs) > 0:
                font_size = para.runs[0].font.size
                # Titel
                if font_size == 361950:
                    title_shape_template = shape
                    print(f'  Titel template: {shape.text}, size={font_size}')
                # Eerste bullet (niet arrow, niet intro)
                elif font_size == 161925 and shape.text.startswith('Propriëtaire'):
                    first_bullet_template = shape
                    print(f'  Bullet template: pos=({shape.left}, {shape.top}), size=({shape.width}, {shape.height})')

if not title_shape_template or not first_bullet_template:
    print('ERROR: Could not find templates!')
    exit(1)

# Nu gaan we slides 5-17 (index 4-16) aanpassen
for slide_idx in range(4, len(prs.slides)):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    print(f'\nProcessing slide {slide_num}...')
    
    # Verzamel titel en bullets uit huidige slide
    title_text = None
    bullet_texts = []
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        
        text = shape.text.strip()
        if not text or text == '→':
            continue
        
        # Bepaal font size
        font_size = None
        if len(shape.text_frame.paragraphs) > 0:
            para = shape.text_frame.paragraphs[0]
            if len(para.runs) > 0:
                font_size = para.runs[0].font.size
        
        # Titel = grootste font
        if font_size and font_size > 200000 and not title_text:
            title_text = text
        elif text and font_size and font_size < 200000:
            bullet_texts.append(text)
    
    if not title_text:
        print(f'  WARNING: Geen titel!')
        continue
    
    # Verwijder alle bestaande TEXT_BOX shapes (behoudt pictures/backgrounds)
    shapes_to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    # Voeg titel toe
    left = title_shape_template.left
    top = title_shape_template.top
    width = title_shape_template.width
    height = title_shape_template.height
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    
    print(f'  ✓ Titel: {title_text}')
    
    # Voeg bullets toe met template-positioning
    bullet_left = first_bullet_template.left
    bullet_top = first_bullet_template.top
    bullet_width = first_bullet_template.width
    bullet_height = first_bullet_template.height
    
    for i, bullet_text in enumerate(bullet_texts):
        current_top = bullet_top + (i * bullet_height) + (i * Inches(0.1))
        
        bullet_box = slide.shapes.add_textbox(bullet_left, current_top, bullet_width, bullet_height)
        bullet_frame = bullet_box.text_frame
        bullet_frame.word_wrap = True
        bullet_frame.clear()
        bullet_p = bullet_frame.paragraphs[0]
        bullet_p.text = bullet_text
        bullet_p.font.size = Pt(12)  # 161925 EMU
        bullet_p.font.bold = True
    
    print(f'  ✓ Aantal bullets: {len(bullet_texts)}')

# Sla op als 17p versie
prs.save(dst_file)
print(f'\n✓ Opgeslagen als: {dst_file}')
