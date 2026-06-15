#!/usr/bin/env python3
from pptx import Presentation

pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(pptx_file)
print(f'Totaal slides: {len(prs.slides)}\n')

# Check slide 3 vs 5
for slide_idx in [2, 4]:
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    print(f'=== SLIDE {slide_num} ===')
    print(f'Aantal shapes: {len(slide.shapes)}')
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, 'text_frame'):
            txt = shape.text[:50] if len(shape.text) > 50 else shape.text
            print(f'  Shape {i}: "{txt}"')
        else:
            print(f'  Shape {i}: {shape.shape_type} (geen text)')
    print()
