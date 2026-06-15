#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt

pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(pptx_file)
print(f'Totaal slides: {len(prs.slides)}\n')

# Transformeer slides 5-17
for slide_idx in range(4, 17):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    # Verzamel alle content
    title_text = None
    all_bullets = []
    title_shape_idx = None
    content_shape_idx = None
    
    for i, shape in enumerate(slide.shapes):
        if not hasattr(shape, 'text_frame'):
            continue
        
        text = shape.text.strip()
        if not text:
            continue
        
        # Font size bepalen
        font_size = None
        if len(shape.text_frame.paragraphs) > 0:
            para = shape.text_frame.paragraphs[0]
            if len(para.runs) > 0:
                font_size = para.runs[0].font.size
        
        # Titel = grote font (360000+)
        if font_size and font_size > 200000 and not title_text:
            title_text = text
            title_shape_idx = i
        # Bullets = rest
        elif text != '→':
            all_bullets.append(text)
            if content_shape_idx is None:
                content_shape_idx = i
    
    if not title_text or content_shape_idx is None:
        print(f'Slide {slide_num}: SKIP')
        continue
    
    # WIJZIG de bestaande shapes in plaats van ze te verwijderen
    # 1. Update title shape
    title_shape = slide.shapes[title_shape_idx]
    title_shape.text_frame.clear()
    p = title_shape.text_frame.paragraphs[0]
    p.text = title_text
    if len(p.runs) > 0:
        p.runs[0].font.size = Pt(44)
        p.runs[0].font.bold = True
    
    # 2. Update content shape
    content_shape = slide.shapes[content_shape_idx]
    content_shape.text_frame.clear()
    
    # Intro line
    if all_bullets:
        p = content_shape.text_frame.paragraphs[0]
        p.text = all_bullets[0]
        p.font.size = Pt(24)
        p.font.bold = True
        
        # Rest as bullets
        for bullet in all_bullets[1:]:
            p = content_shape.text_frame.add_paragraph()
            p.text = ''  # Empty line
            p.font.bold = True
            
            p = content_shape.text_frame.add_paragraph()
            p.text = f'→ {bullet}'
            p.font.size = Pt(12)
            p.font.bold = True
    
    # 3. Verwijder alle ANDERE text shapes
    shapes_to_remove = []
    for i, shape in enumerate(slide.shapes):
        if i != title_shape_idx and i != content_shape_idx:
            if hasattr(shape, 'text_frame'):
                shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    print(f'Slide {slide_num}: ✓ "{title_text[:40]}"')

prs.save(pptx_file)
print(f'\n✓ Transformatie compleet!')
