#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(pptx_file)
print(f'Totaal slides: {len(prs.slides)}\n')

# Voor elk slide van 5 tot 17 (index 4 tot 16)
for slide_idx in range(4, len(prs.slides)):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    print(f'Processing slide {slide_num}...')
    
    # Verzamel alle content
    title_text = None
    bullet_items = []
    shapes_to_remove = []
    title_shape = None
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
            
        text = shape.text.strip()
        if not text:
            continue
        
        # Bepaal font size
        font_size = None
        if len(shape.text_frame.paragraphs) > 0:
            para = shape.text_frame.paragraphs[0]
            if len(para.runs) > 0:
                font_size = para.runs[0].font.size
        
        # Titel = grootte 361950 OF eerste grote font
        if font_size and font_size >= 209550 and not title_text:
            title_text = text
            title_shape = shape
        # Negeer pijltjes
        elif text == '→':
            shapes_to_remove.append(shape)
        # Alles anderen = bullet
        elif text and text != title_text:
            bullet_items.append(text)
            shapes_to_remove.append(shape)
    
    # Nu herstructureer:
    # 1. Zorg dat er een correcte titel is
    if not title_text:
        print(f'  WARNING: Geen titel gevonden!')
        continue
    
    # 2. Verwijder alle kleine shapes
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    # 3. Zorg dat de titel het juiste format heeft (361950 = ~44pt)
    if title_shape:
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(44)  # 361950 EMU = 44pt
                run.font.bold = True
    
    # 4. Voeg één groot tekstvak toe met alle bullets
    if bullet_items:
        # Bepaal positioning (onder de titel)
        left = Inches(0.4)  # 571500 EMU ≈ 0.4 inch
        top = Inches(1.8)   # onder de titel
        width = Inches(6.5)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txBox.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(bullet_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = 0
            p.font.size = Pt(24)  # 203200 EMU ≈ 24pt
            p.font.bold = True
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    print(f'  ✓ Herstructureerd: titel="{title_text}", bullets={len(bullet_items)}')

print(f'\nAlbergeëerde slides: {len(prs.slides) - 4}')
prs.save(pptx_file)
print(f'✓ PPTX opgeslagen!')
