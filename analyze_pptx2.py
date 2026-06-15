#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(pptx_file)
print(f'Totaal slides: {len(prs.slides)}\n')

# Voor elk slide van 5 tot 17 (index 4 tot 16)
for slide_idx in range(4, len(prs.slides)):
    slide = prs.slides[slide_idx]
    print(f'=== SLIDE {slide_idx + 1} ===')
    
    # Vind de titel (grootste font)
    title_text = None
    title_shape = None
    bullet_lines = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text.strip()
            if not text:
                continue
            
            # Bepaal font size
            font_size = None
            if len(shape.text_frame.paragraphs) > 0:
                para = shape.text_frame.paragraphs[0]
                if len(para.runs) > 0:
                    font_size = para.runs[0].font.size
            
            print(f'  Shape: "{text[:50]}" - Font size: {font_size}')
            
            # Titel = grootte 361950
            if font_size and font_size == 361950:
                title_text = text
                title_shape = shape
            # Bullets = alles anders
            elif text and text != '→':
                bullet_lines.append(text)
    
    print(f'  Title: {title_text}')
    print(f'  Bullets ({len(bullet_lines)}): {bullet_lines}')
    print()

print('\nDe volgende stap: voor elke slide 5-17:')
print('1. Titel behouden (font 361950 = ~44pt)')
print('2. Alle andere shapes samensmelten in één groot tekstvak (font 203200 = ~24pt)')
print('3. Pijltjes (→) als bullet-prefix behouden')
