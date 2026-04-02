"""
Script om de PowerPoint-presentatie "Bitemporeel Register" te genereren.
Uitvoeren: python generate_presentatie.py
Output:    bitemporeel_register_presentatie.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Kleuren ──────────────────────────────────────────────────────────
DONKERBLAUW  = RGBColor(0x1E, 0x3A, 0x5F)
MIDDENBLAUW  = RGBColor(0x2C, 0x5F, 0x8A)
LICHTBLAUW   = RGBColor(0x3B, 0x82, 0xF6)
ZACHTBLAUW   = RGBColor(0xBF, 0xDB, 0xFE)
ACHTERGROND  = RGBColor(0xF0, 0xF4, 0xF8)
WIT          = RGBColor(0xFF, 0xFF, 0xFF)
ZWART        = RGBColor(0x20, 0x20, 0x20)
GRIJS        = RGBColor(0x64, 0x74, 0x8B)
DONKERGRIJS  = RGBColor(0x47, 0x55, 0x69)
GROEN        = RGBColor(0x16, 0xA3, 0x4A)
ORANJE       = RGBColor(0xEA, 0x58, 0x0C)
ROOD         = RGBColor(0xDC, 0x26, 0x26)
GEEL_LICHT   = RGBColor(0xFE, 0xF3, 0xC7)
GROEN_LICHT  = RGBColor(0xDC, 0xFC, 0xE7)
BLAUW_LICHT  = RGBColor(0xDB, 0xEA, 0xFE)
ROOD_LICHT   = RGBColor(0xFE, 0xE2, 0xE2)
ORANJE_LICHT = RGBColor(0xFF, 0xED, 0xD5)
PAARS        = RGBColor(0x7C, 0x3A, 0xED)
PAARS_LICHT  = RGBColor(0xED, 0xE9, 0xFE)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def set_text(shape, text, size=18, bold=False, color=ZWART, align=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=18, bold=False, color=ZWART, align=PP_ALIGN.LEFT, font_name="Calibri",
             space_before=Pt(4), space_after=Pt(2), level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p

def add_bullet_text(tf, text, size=18, bold=False, color=ZWART, level=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.level = level
    p.space_before = Pt(4)
    p.space_after = Pt(2)
    return p

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def add_title_bar(slide, title_text, subtitle_text=None):
    """Donkerblauwe titelbalk bovenaan."""
    bar = add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DONKERBLAUW)
    set_text(bar, title_text, size=32, bold=True, color=WIT, align=PP_ALIGN.LEFT)
    bar.text_frame.paragraphs[0].font.name = "Calibri Light"
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.margin_top = Inches(0.15)
    if subtitle_text:
        add_para(bar.text_frame, subtitle_text, size=18, color=ZACHTBLAUW, align=PP_ALIGN.LEFT)
    return bar

def add_footer(slide, slide_num, total=20):
    tb = add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    set_text(tb, f"{slide_num} / {total}", size=11, color=GRIJS, align=PP_ALIGN.RIGHT)

def add_connector_line(slide, x1, y1, x2, y2, color=GRIJS, width=Pt(2)):
    """Voeg een lijn toe als connector."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # type 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_arrow(slide, x1, y1, x2, y2, color=MIDDENBLAUW, width=Pt(2)):
    """Voeg een pijl toe."""
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    # End arrow
    ln = connector.line._ln
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return connector


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: TITEL
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DONKERBLAUW)

# Accent lijn
add_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), LICHTBLAUW)

tb = add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5))
set_text(tb, "Het Bitemporele Register", size=48, bold=True, color=WIT, align=PP_ALIGN.CENTER,
         font_name="Calibri Light")

tb2 = add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(1.2))
set_text(tb2, "Principes, Metamodel & Architectuur", size=28, color=ZACHTBLAUW, align=PP_ALIGN.CENTER)

tb3 = add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8))
set_text(tb3, "Van concept naar implementatie in Go / PostgreSQL", size=18, color=GRIJS, align=PP_ALIGN.CENTER)

add_footer(slide, 1)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Agenda")

items = [
    ("1", "Waarom bitemporeel?", "Het probleem dat we oplossen"),
    ("2", "Twee tijdsdimensies", "Formele tijd & materiële tijd"),
    ("3", "Registratie & Wijzigingen", "Hoe formele tijd wordt vastgelegd"),
    ("4", "UML Metamodel", "Entiteit, Gegevenselement, Relatie"),
    ("5", "Hub + Data Pattern", "Drielagen-architectuur voor correcties"),
    ("6", "Database-ontwerp", "Tabellen, sleutels & versioning"),
    ("7", "Implementatie", "MetaRegistry, Go, Schema-API"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.6) + Inches(0.75) * i
    # Nummer cirkel
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = LICHTBLAUW
    circle.line.fill.background()
    set_text(circle, num, size=18, bold=True, color=WIT, align=PP_ALIGN.CENTER)
    circle.text_frame.margin_left = Inches(0)
    circle.text_frame.margin_right = Inches(0)
    circle.text_frame.margin_top = Inches(0.05)

    tb = add_textbox(slide, Inches(1.8), y, Inches(4), Inches(0.5))
    set_text(tb, title, size=20, bold=True, color=DONKERBLAUW)
    
    tb2 = add_textbox(slide, Inches(6), y + Inches(0.05), Inches(5), Inches(0.5))
    set_text(tb2, desc, size=16, color=GRIJS)

add_footer(slide, 2)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: WAAROM BITEMPOREEL?
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Waarom bitemporeel?", "Het probleem dat we oplossen")

# Probleem-box
box = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2), ROOD_LICHT, ROOD, Pt(2))
tf = set_text(box, "Het probleem", size=22, bold=True, color=ROOD)
box.text_frame.margin_left = Inches(0.3)
box.text_frame.margin_top = Inches(0.2)
add_para(tf, "", size=8)
add_para(tf, "Traditionele databases kennen geen tijdsbesef:", size=17, color=ZWART)
add_para(tf, "▸  UPDATE overschrijft de vorige waarde", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "▸  Geen audittrail: wat was de waarde gisteren?", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "▸  Geen onderscheid: wanneer geregistreerd", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "   vs. wanneer geldig in werkelijkheid", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "▸  Correcties vernietigen de oorspronkelijke data", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "▸  Onmogelijk om historische beslissingen", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "   te reconstrueren", size=16, color=DONKERGRIJS, level=1)

# Oplossing-box
box2 = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), GROEN_LICHT, GROEN, Pt(2))
tf2 = set_text(box2, "De oplossing: bitemporeel", size=22, bold=True, color=GROEN)
box2.text_frame.margin_left = Inches(0.3)
box2.text_frame.margin_top = Inches(0.2)
add_para(tf2, "", size=8)
add_para(tf2, "Twee tijdsassen bewaren alles:", size=17, color=ZWART)
add_para(tf2, "▸  Niets wordt ooit overschreven", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Volledige audittrail van elke wijziging", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Tijdreizen naar elk moment in het verleden", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Correcties behouden het oorspronkelijke pad", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Juridisch verifieerbare registratie", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Toekomstige geldigheid registreerbaar", size=16, color=DONKERGRIJS, level=1)

add_footer(slide, 3)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: TWEE TIJDSDIMENSIES
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Twee tijdsdimensies", "Het fundament van bitemporele registratie")

# Formele tijd box
box_f = add_shape(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
tf = set_text(box_f, "⏱  Formele tijd", size=26, bold=True, color=MIDDENBLAUW)
box_f.text_frame.margin_left = Inches(0.3)
box_f.text_frame.margin_top = Inches(0.2)
add_para(tf, "(registratietijd)", size=16, color=GRIJS)
add_para(tf, "", size=8)
add_para(tf, "Wanneer is iets geregistreerd?", size=18, bold=True, color=ZWART)
add_para(tf, "", size=6)
add_para(tf, "▸  Registratietijdstip (timestamp)", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Opvoer = record komt op de formele tijdlijn", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Afvoer = record gaat van de formele tijdlijn", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Alleen tijdreizen naar het verleden", size=16, color=DONKERGRIJS)
add_para(tf, "", size=8)
add_para(tf, "Vastgelegd via: registratie → wijzigingen", size=15, bold=True, color=MIDDENBLAUW)

# Materiële tijd box
box_m = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8), ORANJE_LICHT, ORANJE, Pt(2))
tf2 = set_text(box_m, "📅  Materiële tijd", size=26, bold=True, color=ORANJE)
box_m.text_frame.margin_left = Inches(0.3)
box_m.text_frame.margin_top = Inches(0.2)
add_para(tf2, "(geldigheidstijd)", size=16, color=GRIJS)
add_para(tf2, "", size=8)
add_para(tf2, "Wanneer geldt iets in de werkelijkheid?", size=18, bold=True, color=ZWART)
add_para(tf2, "", size=6)
add_para(tf2, "▸  Aanvangsdatum (date)", size=16, color=DONKERGRIJS)
add_para(tf2, "▸  Einddatum (date)", size=16, color=DONKERGRIJS)
add_para(tf2, "▸  Tijdreizen naar verleden én toekomst", size=16, color=DONKERGRIJS)
add_para(tf2, "▸  Bijv. voorgenomen verhuizing", size=16, color=DONKERGRIJS)
add_para(tf2, "", size=8)
add_para(tf2, "Vastgelegd via: aanvang / einde records", size=15, bold=True, color=ORANJE)

add_footer(slide, 4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: HOE FORMELE TIJD WERKT - VISUEEL
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Formele tijd: registratie → wijzigingen", "Hoe de formele tijdlijn ontstaat")

# Uitleg links
tb = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(5.5))
tf = set_text(tb, "Kernprincipe", size=22, bold=True, color=DONKERBLAUW)
add_para(tf, "", size=6)
add_para(tf, "Opvoer en afvoer staan niet in de inhoudelijke", size=17, color=ZWART)
add_para(tf, "tabellen, maar in een aparte wijzigingen-tabel.", size=17, color=ZWART)
add_para(tf, "", size=10)
add_para(tf, "Eén registratie = één formeel moment", size=18, bold=True, color=MIDDENBLAUW)
add_para(tf, "", size=6)
add_para(tf, "▸  Heeft een tijdstip (timestamp)", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Bevat 1..n wijzigingen", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Elke wijziging = opvoer of afvoer", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Wijst naar een specifiek record", size=16, color=DONKERGRIJS)
add_para(tf, "", size=10)
add_para(tf, "Soorten registratie:", size=18, bold=True, color=MIDDENBLAUW)
add_para(tf, "▸  Registratie — gewone vastlegging", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Correctie — corrigeert eerdere registratie", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Ongedaanmaking — maakt registratie ongeldig", size=16, color=DONKERGRIJS)

# Diagram rechts: registratie → wijzigingen → records
# Registratie box
reg = add_shape(slide, Inches(7.2), Inches(1.8), Inches(2.8), Inches(1.4), MIDDENBLAUW, line_color=None)
tf_r = set_text(reg, "Registratie", size=20, bold=True, color=WIT, align=PP_ALIGN.CENTER)
reg.text_frame.margin_top = Inches(0.1)
add_para(tf_r, "tijdstip: 2026-03-21 14:30", size=13, color=ZACHTBLAUW, align=PP_ALIGN.CENTER)
add_para(tf_r, "type: Registratie", size=13, color=ZACHTBLAUW, align=PP_ALIGN.CENTER)

# Wijziging boxes
w1 = add_shape(slide, Inches(6.6), Inches(3.8), Inches(2.2), Inches(1.1), BLAUW_LICHT, MIDDENBLAUW)
tf_w1 = set_text(w1, "Wijziging 1", size=15, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
w1.text_frame.margin_top = Inches(0.05)
add_para(tf_w1, "type: Opvoer", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_w1, "record: Persoon #5", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

w2 = add_shape(slide, Inches(9.2), Inches(3.8), Inches(2.2), Inches(1.1), BLAUW_LICHT, MIDDENBLAUW)
tf_w2 = set_text(w2, "Wijziging 2", size=15, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
w2.text_frame.margin_top = Inches(0.05)
add_para(tf_w2, "type: Opvoer", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_w2, "record: Naam #1", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

# Records
r1 = add_shape(slide, Inches(6.6), Inches(5.5), Inches(2.2), Inches(0.9), GROEN_LICHT, GROEN)
set_text(r1, "Persoon #5", size=15, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
r1.text_frame.margin_top = Inches(0.08)
add_para(r1.text_frame, "BSN: 123456782", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

r2 = add_shape(slide, Inches(9.2), Inches(5.5), Inches(2.2), Inches(0.9), GROEN_LICHT, GROEN)
set_text(r2, "Naam v1", size=15, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
r2.text_frame.margin_top = Inches(0.08)
add_para(r2.text_frame, '"De Vries"', size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

# Pijlen (registratie → wijzigingen → records)
add_arrow(slide, Inches(8.6), Inches(3.2), Inches(7.7), Inches(3.8), MIDDENBLAUW)
add_arrow(slide, Inches(8.6), Inches(3.2), Inches(10.3), Inches(3.8), MIDDENBLAUW)
add_arrow(slide, Inches(7.7), Inches(4.9), Inches(7.7), Inches(5.5), GROEN)
add_arrow(slide, Inches(10.3), Inches(4.9), Inches(10.3), Inches(5.5), GROEN)

add_footer(slide, 5)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: OPVOER & AFVOER - AFGELEIDE VELDEN
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Opvoer & afvoer: altijd afgeleid", "De bron van waarheid is de wijzigingen-tabel")

# Kernprincipe
box = add_shape(slide, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.2), GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Pt(2))
tf = set_text(box, "💡  Kernprincipe: opvoer en afvoer in records zijn afgeleide waarden — niet de bron van waarheid.", 
              size=19, bold=True, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)
box.text_frame.margin_top = Inches(0.15)
add_para(tf, "De werkelijke formele tijdlijn wordt bepaald door de registraties en hun wijzigingen.", 
         size=16, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)

# Twee kolommen
# Links: bron van waarheid
lbox = add_shape(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8), WIT, MIDDENBLAUW)
tf_l = set_text(lbox, "Bron van waarheid", size=20, bold=True, color=MIDDENBLAUW)
lbox.text_frame.margin_left = Inches(0.3)
lbox.text_frame.margin_top = Inches(0.15)
add_para(tf_l, "", size=6)
add_para(tf_l, "registratie.tijdstip", size=16, bold=True, color=ZWART)
add_para(tf_l, "  + wijziging.type (opvoer/afvoer)", size=15, color=DONKERGRIJS)
add_para(tf_l, "  + wijziging.representatie_id", size=15, color=DONKERGRIJS)
add_para(tf_l, "", size=8)
add_para(tf_l, "Bij formeel tijdreizen (→ peiltijdstip tf):", size=16, bold=True, color=ZWART)
add_para(tf_l, "▸  Records zelf worden NIET gebruikt", size=15, color=DONKERGRIJS)
add_para(tf_l, "▸  Toestand wordt opnieuw afgeleid", size=15, color=DONKERGRIJS)
add_para(tf_l, "▸  Alle wijzigingen t/m tf verwerken", size=15, color=DONKERGRIJS)

# Rechts: afgeleide velden
rbox = add_shape(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8), WIT, GRIJS)
tf_r = set_text(rbox, "Afgeleide velden in records", size=20, bold=True, color=DONKERGRIJS)
rbox.text_frame.margin_left = Inches(0.3)
rbox.text_frame.margin_top = Inches(0.15)
add_para(tf_r, "", size=6)
add_para(tf_r, "record.opvoer  = afgeleid (tf = nu)", size=16, bold=True, color=ZWART)
add_para(tf_r, "record.afvoer  = afgeleid (tf = nu)", size=16, bold=True, color=ZWART)
add_para(tf_r, "", size=8)
add_para(tf_r, "Geldt voor ALLE lagen:", size=16, bold=True, color=ZWART)
add_para(tf_r, "▸  Entiteit  (A, B)", size=15, color=DONKERGRIJS)
add_para(tf_r, "▸  GE/REL hub  (A_U, Rel_A_B)", size=15, color=DONKERGRIJS)
add_para(tf_r, "▸  _Data  (A_U_Data)", size=15, color=DONKERGRIJS)
add_para(tf_r, "▸  _Aanvang / _Einde  (materieel)", size=15, color=DONKERGRIJS)

add_footer(slide, 6)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: TIJDREIZEN
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Tijdreizen", "De kracht van bitemporele registratie")

# Formeel tijdreizen
box_ft = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.4), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
tf_ft = set_text(box_ft, "⏪  Formeel tijdreizen", size=22, bold=True, color=MIDDENBLAUW)
box_ft.text_frame.margin_left = Inches(0.3)
box_ft.text_frame.margin_top = Inches(0.15)
add_para(tf_ft, "", size=4)
add_para(tf_ft, "Wat was bekend op tijdstip tf?", size=17, bold=True, color=ZWART)
add_para(tf_ft, "▸  Alleen naar het verleden (tf ≤ nu)", size=15, color=DONKERGRIJS)
add_para(tf_ft, "▸  API: ?t=2024-01-01T12:00:00Z", size=15, color=DONKERGRIJS)
add_para(tf_ft, "▸  Verwerk alle wijzigingen t/m tf", size=15, color=DONKERGRIJS)

# Materieel tijdreizen
box_mt = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.4), ORANJE_LICHT, ORANJE, Pt(2))
tf_mt = set_text(box_mt, "⏩  Materieel tijdreizen", size=22, bold=True, color=ORANJE)
box_mt.text_frame.margin_left = Inches(0.3)
box_mt.text_frame.margin_top = Inches(0.15)
add_para(tf_mt, "", size=4)
add_para(tf_mt, "Wat gold op datum tm?", size=17, bold=True, color=ZWART)
add_para(tf_mt, "▸  Verleden én toekomst (tm is vrij)", size=15, color=DONKERGRIJS)
add_para(tf_mt, "▸  Aanvang ≤ tm < einde", size=15, color=DONKERGRIJS)
add_para(tf_mt, "▸  Bijv. toekomstige verhuizing al registreerbaar", size=15, color=DONKERGRIJS)

# Combinatie
box_c = add_shape(slide, Inches(2.5), Inches(4.4), Inches(8.4), Inches(2.6), PAARS_LICHT, PAARS, Pt(2))
tf_c = set_text(box_c, "🔀  Gecombineerd tijdreizen", size=22, bold=True, color=PAARS)
box_c.text_frame.margin_left = Inches(0.3)
box_c.text_frame.margin_top = Inches(0.15)
add_para(tf_c, "", size=4)
add_para(tf_c, "Wat was op registratietijdstip tf bekend over geldigheidsdatum tm?", size=17, bold=True, color=ZWART)
add_para(tf_c, "", size=4)
add_para(tf_c, "▸  Eerst formeel: welke records waren bekend op tf?", size=15, color=DONKERGRIJS)
add_para(tf_c, "▸  Dan materieel: welke daarvan golden op tm?", size=15, color=DONKERGRIJS)
add_para(tf_c, "▸  Maakt volledige audit-trail mogelijk over twee dimensies", size=15, color=DONKERGRIJS)

add_footer(slide, 7)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: UML METAMODEL - REPRESENTATIETYPES
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "UML Metamodel: representatietypes", "De drie bouwstenen van het register")

# Representatie superklasse
rep = add_shape(slide, Inches(4.5), Inches(1.5), Inches(4.4), Inches(1.3), PAARS_LICHT, PAARS, Pt(2))
tf_rep = set_text(rep, "《abstract》Representatie", size=18, bold=True, color=PAURS if False else PAARS, align=PP_ALIGN.CENTER)
rep.text_frame.margin_top = Inches(0.05)
add_para(tf_rep, "opvoer : timestamp  |  afvoer : timestamp", size=13, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_rep, "naam, metatype", size=13, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

# Entiteit
ent = add_shape(slide, Inches(0.8), Inches(3.5), Inches(3.5), Inches(2.5), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
tf_e = set_text(ent, "Entiteit", size=22, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
ent.text_frame.margin_top = Inches(0.1)
add_para(tf_e, "", size=4)
add_para(tf_e, "+ id : integer (PK)", size=14, color=DONKERGRIJS, align=PP_ALIGN.LEFT)
add_para(tf_e, "", size=6)
add_para(tf_e, "▸ Zelfstandig identificeerbaar", size=14, color=DONKERGRIJS)
add_para(tf_e, "▸ Heeft onderliggende GE's", size=14, color=DONKERGRIJS)
add_para(tf_e, "▸ Bijv. Persoon, Locatie", size=14, color=DONKERGRIJS)

# Gegevenselement
ge = add_shape(slide, Inches(4.9), Inches(3.5), Inches(3.5), Inches(2.5), GROEN_LICHT, GROEN, Pt(2))
tf_g = set_text(ge, "Gegevenselement", size=22, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge.text_frame.margin_top = Inches(0.1)
add_para(tf_g, "", size=4)
add_para(tf_g, "+ ent_id : FK  + rel_id : PFK", size=14, color=DONKERGRIJS, align=PP_ALIGN.LEFT)
add_para(tf_g, "", size=6)
add_para(tf_g, "▸ Compositie bij entiteit", size=14, color=DONKERGRIJS)
add_para(tf_g, "▸ Draagt inhoudsvelden", size=14, color=DONKERGRIJS)
add_para(tf_g, "▸ Bijv. Naam, Burgerschap", size=14, color=DONKERGRIJS)

# Relatie
rel = add_shape(slide, Inches(9.0), Inches(3.5), Inches(3.5), Inches(2.5), ORANJE_LICHT, ORANJE, Pt(2))
tf_rl = set_text(rel, "Relatie", size=22, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
rel.text_frame.margin_top = Inches(0.1)
add_para(tf_rl, "", size=4)
add_para(tf_rl, "+ ent_id : FK  + doel_id : FK", size=14, color=DONKERGRIJS, align=PP_ALIGN.LEFT)
add_para(tf_rl, "", size=6)
add_para(tf_rl, "▸ Koppeling twee entiteiten", size=14, color=DONKERGRIJS)
add_para(tf_rl, "▸ Kan eigen inhoudsvelden dragen", size=14, color=DONKERGRIJS)
add_para(tf_rl, "▸ Bijv. Bereikbaarheid (Pers↔Loc)", size=14, color=DONKERGRIJS)

# Specialisatie-pijlen (≈ inheritance)
add_arrow(slide, Inches(2.55), Inches(3.5), Inches(5.7), Inches(2.8), GRIJS, Pt(2))
add_arrow(slide, Inches(6.65), Inches(3.5), Inches(6.65), Inches(2.8), GRIJS, Pt(2))
add_arrow(slide, Inches(10.75), Inches(3.5), Inches(7.7), Inches(2.8), GRIJS, Pt(2))

# Compositie label
tb = add_textbox(slide, Inches(3.6), Inches(6.3), Inches(6), Inches(0.5))
set_text(tb, "Entiteit  ◆───────── 0..*  Gegevenselement          Relatie  erft van  Gegevenselement", 
         size=15, color=GRIJS, align=PP_ALIGN.CENTER)

add_footer(slide, 8)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: CONCREET UML VOORBEELD
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "UML: concreet voorbeeld", "NatuurlijkPersoon en Locatie met hun gegevenselementen")

# Persoon entiteit
p = add_shape(slide, Inches(0.8), Inches(2.0), Inches(3.0), Inches(1.2), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
set_text(p, "《entiteit, materieel》", size=12, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
p.text_frame.margin_top = Inches(0.05)
add_para(p.text_frame, "NatuurlijkPersoon", size=20, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)

# GE's van Persoon
ge1 = add_shape(slide, Inches(0.3), Inches(3.8), Inches(2.5), Inches(1.0), GROEN_LICHT, GROEN)
tf1 = set_text(ge1, "PersoonsIdentificatie", size=14, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge1.text_frame.margin_top = Inches(0.05)
add_para(tf1, "bsn, ingezetene", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

ge2 = add_shape(slide, Inches(0.3), Inches(5.0), Inches(2.5), Inches(1.2), GROEN_LICHT, GROEN)
tf2 = set_text(ge2, "Naam", size=14, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge2.text_frame.margin_top = Inches(0.05)
add_para(tf2, "voorletters, roepnaam", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf2, "tussenvoegsel, achternaam", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

ge3 = add_shape(slide, Inches(3.2), Inches(3.8), Inches(2.5), Inches(1.0), GROEN_LICHT, GROEN)
tf3 = set_text(ge3, "🕐 Burgerschap", size=14, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge3.text_frame.margin_top = Inches(0.05)
add_para(tf3, "landcode, nationaliteit", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

# Relatie
rel_box = add_shape(slide, Inches(5.0), Inches(2.0), Inches(3.2), Inches(1.2), ORANJE_LICHT, ORANJE, Pt(2))
set_text(rel_box, "《relatie, materieel》", size=12, color=ORANJE, align=PP_ALIGN.CENTER)
rel_box.text_frame.margin_top = Inches(0.05)
add_para(rel_box.text_frame, "Bereikbaarheid", size=20, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)

# Locatie entiteit
loc = add_shape(slide, Inches(9.2), Inches(2.0), Inches(3.0), Inches(1.2), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
set_text(loc, "《entiteit, materieel》", size=12, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
loc.text_frame.margin_top = Inches(0.05)
add_para(loc.text_frame, "Locatie", size=20, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)

# GE's van Locatie
ge4 = add_shape(slide, Inches(9.0), Inches(3.8), Inches(2.5), Inches(1.0), GROEN_LICHT, GROEN)
tf4 = set_text(ge4, "Adres", size=14, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge4.text_frame.margin_top = Inches(0.05)
add_para(tf4, "straat, huisnr, postcode", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

# Pijlen
add_arrow(slide, Inches(1.5), Inches(3.2), Inches(1.5), Inches(3.8), GROEN)
add_arrow(slide, Inches(1.5), Inches(3.2), Inches(1.5), Inches(5.0), GROEN)
add_arrow(slide, Inches(3.0), Inches(3.2), Inches(4.4), Inches(3.8), GROEN)
add_arrow(slide, Inches(3.8), Inches(2.6), Inches(5.0), Inches(2.6), ORANJE)
add_arrow(slide, Inches(8.2), Inches(2.6), Inches(9.2), Inches(2.6), ORANJE)
add_arrow(slide, Inches(10.3), Inches(3.2), Inches(10.3), Inches(3.8), GROEN)

# Legenda
leg = add_shape(slide, Inches(3.5), Inches(5.5), Inches(6.0), Inches(1.3), WIT, GRIJS)
tf_leg = set_text(leg, "Legenda", size=14, bold=True, color=DONKERGRIJS)
leg.text_frame.margin_left = Inches(0.2)
leg.text_frame.margin_top = Inches(0.05)
add_para(tf_leg, "🟦  Entiteit (blauw)    🟩  Gegevenselement (groen)    🟧  Relatie (oranje)", size=13, color=DONKERGRIJS)
add_para(tf_leg, "🕐  = materieel (heeft aanvang/einde)    Pijl = compositie (onderdeel van)", size=13, color=DONKERGRIJS)

add_footer(slide, 9)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: HUB + DATA PATTERN - INTRO
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Hub + Data Pattern", "Waarom opsplitsing nodig is")

# Probleem v05
box_p = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.8), ROOD_LICHT, ROOD, Pt(2))
tf_p = set_text(box_p, "Probleem (v05-aanpak)", size=20, bold=True, color=ROOD)
box_p.text_frame.margin_left = Inches(0.3)
box_p.text_frame.margin_top = Inches(0.15)
add_para(tf_p, "", size=4)
add_para(tf_p, "GE-tabel bevat alles in één record:", size=16, color=ZWART)
add_para(tf_p, "  structurele FK's + inhoud + opvoer/afvoer", size=15, color=DONKERGRIJS)
add_para(tf_p, "", size=6)
add_para(tf_p, "▸  Correctie van inhoud → heel record", size=15, color=DONKERGRIJS)
add_para(tf_p, "   vervangen, rel_id gaat op", size=15, color=DONKERGRIJS)
add_para(tf_p, "▸  Geen stabiel ankerpunt voor materiële", size=15, color=DONKERGRIJS)
add_para(tf_p, "   aanvang/einde op GE-niveau", size=15, color=DONKERGRIJS)

# Oplossing v06
box_o = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8), GROEN_LICHT, GROEN, Pt(2))
tf_o = set_text(box_o, "Oplossing (v06 Hub+Data)", size=20, bold=True, color=GROEN)
box_o.text_frame.margin_left = Inches(0.3)
box_o.text_frame.margin_top = Inches(0.15)
add_para(tf_o, "", size=4)
add_para(tf_o, "Splits elk GE/REL-type in lagen:", size=16, color=ZWART)
add_para(tf_o, "", size=6)
add_para(tf_o, "▸  Hub = stabiel identiteitsanker", size=15, color=DONKERGRIJS)
add_para(tf_o, "▸  _Data = geversioned inhoud", size=15, color=DONKERGRIJS)
add_para(tf_o, "▸  _Aanvang/_Einde = materiële tijdlijn", size=15, color=DONKERGRIJS)
add_para(tf_o, "", size=6)
add_para(tf_o, "Elke laag apart corrigeerbaar!", size=16, bold=True, color=GROEN)

# Diagram onderaan: drielagen
hub = add_shape(slide, Inches(2.0), Inches(5.0), Inches(2.8), Inches(1.6), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
tf_h = set_text(hub, "Hub", size=20, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
hub.text_frame.margin_top = Inches(0.1)
add_para(tf_h, "(ent_id, rel_id)", size=14, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_h, "Stabiel ankerpunt", size=13, color=GRIJS, align=PP_ALIGN.CENTER)

data = add_shape(slide, Inches(5.3), Inches(5.0), Inches(2.8), Inches(1.6), GROEN_LICHT, GROEN, Pt(2))
tf_d = set_text(data, "_Data", size=20, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
data.text_frame.margin_top = Inches(0.1)
add_para(tf_d, "(ent_id, rel_id, versie)", size=14, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_d, "Geversioned inhoud", size=13, color=GRIJS, align=PP_ALIGN.CENTER)

ae = add_shape(slide, Inches(8.6), Inches(5.0), Inches(2.8), Inches(1.6), ORANJE_LICHT, ORANJE, Pt(2))
tf_ae = set_text(ae, "_Aanvang / _Einde", size=18, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
ae.text_frame.margin_top = Inches(0.1)
add_para(tf_ae, "(ent_id, rel_id, versie)", size=14, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_ae, "Materiële tijdlijn", size=13, color=GRIJS, align=PP_ALIGN.CENTER)

# Pijlen hub → data, hub → aanvang/einde
add_arrow(slide, Inches(4.8), Inches(5.8), Inches(5.3), Inches(5.8), GROEN, Pt(2))
add_arrow(slide, Inches(4.8), Inches(5.8), Inches(8.6), Inches(5.8), ORANJE, Pt(2))

add_footer(slide, 10)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11: HUB + DATA - HIËRARCHIE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Hub + Data: hiërarchie", "Entiteit → Hub → Data / Aanvang / Einde")

# Entiteit A
ent_a = add_shape(slide, Inches(0.6), Inches(1.7), Inches(2.3), Inches(0.8), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
set_text(ent_a, "NatuurlijkPersoon", size=16, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
ent_a.text_frame.margin_top = Inches(0.1)

# Hub hierarchy for PersoonsIdentificatie (niet-materieel)
h1 = add_shape(slide, Inches(3.5), Inches(1.5), Inches(2.6), Inches(0.6), RGBColor(0xE0, 0xE7, 0xFF), MIDDENBLAUW)
set_text(h1, "PersoonsIdentificatie", size=13, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
h1.text_frame.margin_top = Inches(0.05)

d1 = add_shape(slide, Inches(6.5), Inches(1.5), Inches(2.8), Inches(0.6), GROEN_LICHT, GROEN)
set_text(d1, "PersoonsIdent_Data", size=13, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
d1.text_frame.margin_top = Inches(0.05)

# Hub hierarchy for Naam (niet-materieel)
h2 = add_shape(slide, Inches(3.5), Inches(2.5), Inches(2.6), Inches(0.6), RGBColor(0xE0, 0xE7, 0xFF), MIDDENBLAUW)
set_text(h2, "Naam (hub)", size=13, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
h2.text_frame.margin_top = Inches(0.05)

d2 = add_shape(slide, Inches(6.5), Inches(2.5), Inches(2.8), Inches(0.6), GROEN_LICHT, GROEN)
set_text(d2, "Naam_Data", size=13, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
d2.text_frame.margin_top = Inches(0.05)

# Hub hierarchy for Burgerschap (materieel)
h3 = add_shape(slide, Inches(3.5), Inches(3.5), Inches(2.6), Inches(0.6), RGBColor(0xE0, 0xE7, 0xFF), MIDDENBLAUW)
set_text(h3, "🕐 Burgerschap (hub)", size=13, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
h3.text_frame.margin_top = Inches(0.05)

d3 = add_shape(slide, Inches(6.5), Inches(3.3), Inches(2.8), Inches(0.6), GROEN_LICHT, GROEN)
set_text(d3, "Burgerschap_Data", size=13, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
d3.text_frame.margin_top = Inches(0.05)

a3 = add_shape(slide, Inches(6.5), Inches(4.1), Inches(2.8), Inches(0.6), ORANJE_LICHT, ORANJE)
set_text(a3, "Burgerschap_Aanvang", size=13, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
a3.text_frame.margin_top = Inches(0.05)

e3 = add_shape(slide, Inches(9.7), Inches(4.1), Inches(2.8), Inches(0.6), ROOD_LICHT, ROOD)
set_text(e3, "Burgerschap_Einde", size=13, bold=True, color=ROOD, align=PP_ALIGN.CENTER)
e3.text_frame.margin_top = Inches(0.05)

# Hub hierarchy for Bereikbaarheid (relatie, materieel)
h4 = add_shape(slide, Inches(3.5), Inches(5.1), Inches(2.6), Inches(0.6), ORANJE_LICHT, ORANJE)
set_text(h4, "🕐 Bereikbaarheid (rel)", size=13, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
h4.text_frame.margin_top = Inches(0.05)

d4 = add_shape(slide, Inches(6.5), Inches(4.9), Inches(2.8), Inches(0.6), GROEN_LICHT, GROEN)
set_text(d4, "Bereikbaarheid_Data", size=13, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
d4.text_frame.margin_top = Inches(0.05)

a4 = add_shape(slide, Inches(6.5), Inches(5.7), Inches(2.8), Inches(0.6), ORANJE_LICHT, ORANJE)
set_text(a4, "Bereikbaarheid_Aanvang", size=12, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
a4.text_frame.margin_top = Inches(0.05)

e4 = add_shape(slide, Inches(9.7), Inches(5.7), Inches(2.8), Inches(0.6), ROOD_LICHT, ROOD)
set_text(e4, "Bereikbaarheid_Einde", size=13, bold=True, color=ROOD, align=PP_ALIGN.CENTER)
e4.text_frame.margin_top = Inches(0.05)

# Entiteits-plumbing
ep_a = add_shape(slide, Inches(0.3), Inches(5.1), Inches(2.6), Inches(0.6), ORANJE_LICHT, ORANJE)
set_text(ep_a, "Pers_Aanvang (plumbing)", size=12, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
ep_a.text_frame.margin_top = Inches(0.05)

ep_e = add_shape(slide, Inches(0.3), Inches(5.9), Inches(2.6), Inches(0.6), ROOD_LICHT, ROOD)
set_text(ep_e, "Pers_Einde (plumbing)", size=12, bold=True, color=ROOD, align=PP_ALIGN.CENTER)
ep_e.text_frame.margin_top = Inches(0.05)

# Pijlen
add_arrow(slide, Inches(2.9), Inches(2.1), Inches(3.5), Inches(1.8), MIDDENBLAUW)
add_arrow(slide, Inches(2.9), Inches(2.1), Inches(3.5), Inches(2.8), MIDDENBLAUW)
add_arrow(slide, Inches(2.9), Inches(2.1), Inches(3.5), Inches(3.8), MIDDENBLAUW)
add_arrow(slide, Inches(2.9), Inches(2.1), Inches(3.5), Inches(5.4), ORANJE)

add_arrow(slide, Inches(6.1), Inches(1.8), Inches(6.5), Inches(1.8), GROEN)
add_arrow(slide, Inches(6.1), Inches(2.8), Inches(6.5), Inches(2.8), GROEN)
add_arrow(slide, Inches(6.1), Inches(3.8), Inches(6.5), Inches(3.6), GROEN)
add_arrow(slide, Inches(6.1), Inches(3.8), Inches(6.5), Inches(4.4), ORANJE)

add_arrow(slide, Inches(6.1), Inches(5.4), Inches(6.5), Inches(5.2), GROEN)
add_arrow(slide, Inches(6.1), Inches(5.4), Inches(6.5), Inches(6.0), ORANJE)

# Einde pijlen (van hub naar einde)
add_arrow(slide, Inches(9.3), Inches(4.4), Inches(9.7), Inches(4.4), ROOD)
add_arrow(slide, Inches(9.3), Inches(6.0), Inches(9.7), Inches(6.0), ROOD)

# Entiteit → plumbing
add_arrow(slide, Inches(1.6), Inches(2.5), Inches(1.6), Inches(5.1), ORANJE)

# Legenda
leg = add_shape(slide, Inches(10), Inches(1.5), Inches(2.8), Inches(2.2), WIT, GRIJS)
tf_leg = set_text(leg, "Legenda", size=13, bold=True, color=DONKERGRIJS)
leg.text_frame.margin_left = Inches(0.15)
leg.text_frame.margin_top = Inches(0.05)
add_para(tf_leg, "🟦  Hub", size=12, color=MIDDENBLAUW)
add_para(tf_leg, "🟩  _Data", size=12, color=GROEN)
add_para(tf_leg, "🟧  _Aanvang", size=12, color=ORANJE)
add_para(tf_leg, "🟥  _Einde", size=12, color=ROOD)
add_para(tf_leg, "🕐  Materieel", size=12, color=DONKERGRIJS)

add_footer(slide, 11)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12: VOORBEELD - NAAMSWIJZIGING MET CORRECTIE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Voorbeeld: naamswijziging met correctie", "Hoe hub + data + aanvang/einde samenwerken")

# Stap 1
s1 = add_shape(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(2.6), WIT, MIDDENBLAUW, Pt(2))
tf_s1 = set_text(s1, "Stap 1 — Registratie (t_reg = 21/3/2026)", size=17, bold=True, color=MIDDENBLAUW)
s1.text_frame.margin_left = Inches(0.25)
s1.text_frame.margin_top = Inches(0.1)
add_para(tf_s1, "", size=4)
add_para(tf_s1, "Hub 1: Naam rel_id=1", size=14, bold=True, color=ZWART)
add_para(tf_s1, "  _Data v1: \"Jansen\"  |  Aanvang: geboortedatum  |  Einde: 20/3", size=13, color=DONKERGRIJS)
add_para(tf_s1, "", size=4)
add_para(tf_s1, "Hub 2: Naam rel_id=2", size=14, bold=True, color=ZWART)
add_para(tf_s1, "  _Data v1: \"De Vries\"  |  Aanvang: 21/3/2026  |  Einde: —", size=13, color=DONKERGRIJS)
add_para(tf_s1, "", size=4)
add_para(tf_s1, "⚠  Aanvangsdatum foutief: was 21/3, moet 1/1/2025 zijn", size=14, color=ROOD)

# Stap 2
s2 = add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.6), WIT, ORANJE, Pt(2))
tf_s2 = set_text(s2, "Stap 2 — Correctie (t_corr > t_reg)", size=17, bold=True, color=ORANJE)
s2.text_frame.margin_left = Inches(0.25)
s2.text_frame.margin_top = Inches(0.1)
add_para(tf_s2, "", size=4)
add_para(tf_s2, "Rechtbank: naamswijziging per 1/1/2025", size=14, bold=True, color=ZWART)
add_para(tf_s2, "", size=4)
add_para(tf_s2, "Hub 1: Einde v1 → afgevoerd, Einde v2 = 31/12/2024", size=13, color=DONKERGRIJS)
add_para(tf_s2, "Hub 2: Aanvang v1 → afgevoerd, Aanvang v2 = 1/1/2025", size=13, color=DONKERGRIJS)
add_para(tf_s2, "", size=6)
add_para(tf_s2, "✓  Hub-records ongewijzigd", size=14, color=GROEN)
add_para(tf_s2, "✓  _Data-records ongewijzigd", size=14, color=GROEN)
add_para(tf_s2, "✓  Alleen _Aanvang/_Einde gecorrigeerd!", size=14, bold=True, color=GROEN)

# Key insight
box_i = add_shape(slide, Inches(0.4), Inches(4.5), Inches(12.6), Inches(1.2), GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Pt(2))
tf_i = set_text(box_i, "💡  Elke laag is apart corrigeerbaar — inhoud, identiteit en materiële tijdlijn beïnvloeden elkaar niet.", 
                size=18, bold=True, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)
box_i.text_frame.margin_top = Inches(0.15)
add_para(tf_i, "De drielaagse opsplitsing (Hub + _Data + _Aanvang/_Einde) is hiervoor essentieel.", 
         size=15, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)

# Tabel onderaan
headers = ["Laag", "rel_id", "versie", "waarde", "opvoer", "afvoer"]
data_rows = [
    ["Hub 1", "1", "—", "—", "t_reg", "—"],
    ["Hub 1 _data", "1", "1", "\"Jansen\"", "t_reg", "—"],
    ["Hub 1 _einde", "1", "1", "20/3/2026", "t_reg", "t_corr"],
    ["Hub 1 _einde", "1", "2", "31/12/2024", "t_corr", "—"],
    ["Hub 2 _aanvang", "2", "1", "21/3/2026", "t_reg", "t_corr"],
    ["Hub 2 _aanvang", "2", "2", "1/1/2025", "t_corr", "—"],
]

# Tabel-achtergrond
tbl_y = Inches(5.9)
col_widths = [Inches(2.0), Inches(1.0), Inches(1.0), Inches(2.0), Inches(1.5), Inches(1.5)]
total_w = sum(cw for cw in col_widths)
tbl_x = Inches(1.5)

# Header row
x = tbl_x
for i, h in enumerate(headers):
    hdr = add_rect(slide, x, tbl_y, col_widths[i], Inches(0.35), MIDDENBLAUW)
    set_text(hdr, h, size=12, bold=True, color=WIT, align=PP_ALIGN.CENTER)
    hdr.text_frame.margin_top = Inches(0.02)
    hdr.text_frame.margin_left = Inches(0.05)
    x += col_widths[i]

# Data rows
for r_idx, row in enumerate(data_rows):
    x = tbl_x
    row_y = tbl_y + Inches(0.35) + Inches(0.3) * r_idx
    bg = WIT if r_idx % 2 == 0 else RGBColor(0xF0, 0xF4, 0xF8)
    for c_idx, cell in enumerate(row):
        c = add_rect(slide, x, row_y, col_widths[c_idx], Inches(0.3), bg, GRIJS, Pt(0.5))
        clr = ROOD if cell == "t_corr" else ZWART
        set_text(c, cell, size=11, color=clr, align=PP_ALIGN.CENTER)
        c.text_frame.margin_top = Inches(0.01)
        c.text_frame.margin_left = Inches(0.05)
        x += col_widths[c_idx]

add_footer(slide, 12)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13: DATABASE-ONTWERP: ENTITEIT
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Database-ontwerp: entiteiten", "De basis van het register in PostgreSQL")

# Entiteittabel
ent_box = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(2.4), WIT, MIDDENBLAUW, Pt(2))
tf_e = set_text(ent_box, "natuurlijk_persoon", size=20, bold=True, color=MIDDENBLAUW)
ent_box.text_frame.margin_left = Inches(0.3)
ent_box.text_frame.margin_top = Inches(0.1)
add_para(tf_e, "", size=4)
add_para(tf_e, "id          INTEGER  PK", size=15, color=ZWART, font_name="Consolas")
add_para(tf_e, "opvoer      TIMESTAMPTZ     — afgeleid", size=15, color=GRIJS, font_name="Consolas")
add_para(tf_e, "afvoer      TIMESTAMPTZ     — afgeleid", size=15, color=GRIJS, font_name="Consolas")
add_para(tf_e, "", size=6)
add_para(tf_e, "Geen inhoudsvelden — die zitten in GE's", size=14, color=DONKERGRIJS)

# Registratie tabel
reg_box = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.4), WIT, PAARS, Pt(2))
tf_reg = set_text(reg_box, "registratie", size=20, bold=True, color=PAARS)
reg_box.text_frame.margin_left = Inches(0.3)
reg_box.text_frame.margin_top = Inches(0.1)
add_para(tf_reg, "", size=4)
add_para(tf_reg, "id                          SERIAL  PK", size=14, color=ZWART, font_name="Consolas")
add_para(tf_reg, "tijdstip                    TIMESTAMPTZ", size=14, color=ZWART, font_name="Consolas")
add_para(tf_reg, "registratie_type            TEXT", size=14, color=ZWART, font_name="Consolas")
add_para(tf_reg, "corrigeert_registratie_id   INT FK", size=14, color=GRIJS, font_name="Consolas")
add_para(tf_reg, "maakt_ongedaan_reg_id       INT FK", size=14, color=GRIJS, font_name="Consolas")

# Wijziging tabel
wij_box = add_shape(slide, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.6), WIT, PAARS, Pt(2))
tf_wij = set_text(wij_box, "wijziging", size=20, bold=True, color=PAARS)
wij_box.text_frame.margin_left = Inches(0.3)
wij_box.text_frame.margin_top = Inches(0.1)
add_para(tf_wij, "", size=4)
add_para(tf_wij, "id                  SERIAL  PK", size=14, color=ZWART, font_name="Consolas")
add_para(tf_wij, "registratie_id      INT  FK → registratie", size=14, color=ZWART, font_name="Consolas")
add_para(tf_wij, "wijziging_type      TEXT  (Opvoer/Afvoer)", size=14, color=ZWART, font_name="Consolas")
add_para(tf_wij, "representatienaam   TEXT  (tabelnaam)", size=14, color=ZWART, font_name="Consolas")
add_para(tf_wij, "representatie_id    TEXT  (samengesteld PK)", size=14, color=ZWART, font_name="Consolas")

# Plumbing tabellen
plumb = add_shape(slide, Inches(0.6), Inches(4.3), Inches(5.5), Inches(2.6), WIT, ORANJE, Pt(2))
tf_pl = set_text(plumb, "Materiële plumbing: pers_aanvang / pers_einde", size=16, bold=True, color=ORANJE)
plumb.text_frame.margin_left = Inches(0.3)
plumb.text_frame.margin_top = Inches(0.1)
add_para(tf_pl, "", size=4)
add_para(tf_pl, "id       INT  FK → natuurlijk_persoon(id)", size=14, color=ZWART, font_name="Consolas")
add_para(tf_pl, "versie   SERIAL  (autoincrement per id)", size=14, color=ZWART, font_name="Consolas")
add_para(tf_pl, "datum    DATE", size=14, color=ZWART, font_name="Consolas")
add_para(tf_pl, "opvoer   TIMESTAMPTZ  — afgeleid", size=14, color=GRIJS, font_name="Consolas")
add_para(tf_pl, "afvoer   TIMESTAMPTZ  — afgeleid", size=14, color=GRIJS, font_name="Consolas")
add_para(tf_pl, "", size=4)
add_para(tf_pl, "PK: (id, versie)", size=14, bold=True, color=ORANJE)

add_footer(slide, 13)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14: DATABASE-ONTWERP: HUB + DATA TABELLEN
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Database-ontwerp: Hub + Data tabellen", "Hoe een gegevenselement eruitziet in PostgreSQL")

# Hub tabel
hub_box = add_shape(slide, Inches(0.4), Inches(1.5), Inches(4.0), Inches(2.5), WIT, MIDDENBLAUW, Pt(2))
tf_hub = set_text(hub_box, "naam  (hub)", size=18, bold=True, color=MIDDENBLAUW)
hub_box.text_frame.margin_left = Inches(0.25)
hub_box.text_frame.margin_top = Inches(0.1)
add_para(tf_hub, "", size=4)
add_para(tf_hub, "persoon_id  INT  FK → persoon(id)", size=13, color=ZWART, font_name="Consolas")
add_para(tf_hub, "rel_id      INT  autoincrement/scope", size=13, color=ZWART, font_name="Consolas")
add_para(tf_hub, "opvoer      TIMESTAMPTZ  — afgeleid", size=13, color=GRIJS, font_name="Consolas")
add_para(tf_hub, "afvoer      TIMESTAMPTZ  — afgeleid", size=13, color=GRIJS, font_name="Consolas")
add_para(tf_hub, "", size=4)
add_para(tf_hub, "PK: (persoon_id, rel_id)", size=13, bold=True, color=MIDDENBLAUW)
add_para(tf_hub, "Geen inhoudsvelden!", size=13, bold=True, color=ROOD)

# Data tabel
data_box = add_shape(slide, Inches(4.7), Inches(1.5), Inches(4.0), Inches(2.5), WIT, GROEN, Pt(2))
tf_data = set_text(data_box, "naam_data", size=18, bold=True, color=GROEN)
data_box.text_frame.margin_left = Inches(0.25)
data_box.text_frame.margin_top = Inches(0.1)
add_para(tf_data, "", size=4)
add_para(tf_data, "persoon_id  INT  FK → naam", size=13, color=ZWART, font_name="Consolas")
add_para(tf_data, "rel_id      INT  FK → naam", size=13, color=ZWART, font_name="Consolas")
add_para(tf_data, "versie      INT  autoincrement", size=13, color=ZWART, font_name="Consolas")
add_para(tf_data, "voorletters TEXT", size=13, color=ZWART, font_name="Consolas")
add_para(tf_data, "achternaam  TEXT", size=13, color=ZWART, font_name="Consolas")
add_para(tf_data, "opvoer/afvoer — afgeleid", size=13, color=GRIJS, font_name="Consolas")
add_para(tf_data, "", size=4)
add_para(tf_data, "PK: (persoon_id, rel_id, versie)", size=13, bold=True, color=GROEN)

# Aanvang/Einde tabellen (indien materieel)
ae_box = add_shape(slide, Inches(9.0), Inches(1.5), Inches(4.0), Inches(2.5), WIT, ORANJE, Pt(2))
tf_ae = set_text(ae_box, "naam_aanvang / naam_einde", size=16, bold=True, color=ORANJE)
ae_box.text_frame.margin_left = Inches(0.25)
ae_box.text_frame.margin_top = Inches(0.1)
add_para(tf_ae, "Alleen als IsMaterieel = true", size=12, color=GRIJS)
add_para(tf_ae, "", size=4)
add_para(tf_ae, "persoon_id  INT  FK → naam", size=13, color=ZWART, font_name="Consolas")
add_para(tf_ae, "rel_id      INT  FK → naam", size=13, color=ZWART, font_name="Consolas")
add_para(tf_ae, "versie      INT  autoincrement", size=13, color=ZWART, font_name="Consolas")
add_para(tf_ae, "datum       DATE", size=13, color=ZWART, font_name="Consolas")
add_para(tf_ae, "opvoer/afvoer — afgeleid", size=13, color=GRIJS, font_name="Consolas")
add_para(tf_ae, "", size=4)
add_para(tf_ae, "PK: (persoon_id, rel_id, versie)", size=13, bold=True, color=ORANJE)

# Relatieve autoincrement uitleg
rai = add_shape(slide, Inches(0.4), Inches(4.4), Inches(12.4), Inches(2.6), GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Pt(2))
tf_rai = set_text(rai, "Relatieve autoincrement", size=20, bold=True, color=RGBColor(0x85, 0x4D, 0x0E))
rai.text_frame.margin_left = Inches(0.3)
rai.text_frame.margin_top = Inches(0.1)
add_para(tf_rai, "", size=4)
add_para(tf_rai, "De rel_id en versie worden automatisch opgehoogd binnen de scope van de parent:", size=16, color=ZWART)
add_para(tf_rai, "", size=4)
add_para(tf_rai, "▸  Hub rel_id:  binnen (persoon_id) en type → Persoon #1 heeft Naam rel_id=1, rel_id=2, ...", size=15, color=DONKERGRIJS)
add_para(tf_rai, "▸  _Data versie:  binnen (persoon_id, rel_id) → Naam 1 v1, v2, v3, ...", size=15, color=DONKERGRIJS)
add_para(tf_rai, "▸  _Aanvang versie:  idem, maar vorige versie wordt automatisch afgevoerd (enkelvoudig)", size=15, color=DONKERGRIJS)
add_para(tf_rai, "", size=6)
add_para(tf_rai, "Dit maakt het mogelijk om volledige versiegeschiedenis bij te houden zonder globale sequences.", size=15, bold=True, color=RGBColor(0x85, 0x4D, 0x0E))

add_footer(slide, 14)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15: DATABASE: RELATIE TABEL
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Database-ontwerp: relaties", "Koppeling tussen entiteiten")

# Relatie hub
rel_box = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(2.5), WIT, ORANJE, Pt(2))
tf_rel = set_text(rel_box, "bereikbaarheid  (relatie-hub)", size=18, bold=True, color=ORANJE)
rel_box.text_frame.margin_left = Inches(0.3)
rel_box.text_frame.margin_top = Inches(0.1)
add_para(tf_rel, "", size=4)
add_para(tf_rel, "persoon_id  INT  FK → persoon(id)", size=14, color=ZWART, font_name="Consolas")
add_para(tf_rel, "rel_id      INT  autoincrement/scope", size=14, color=ZWART, font_name="Consolas")
add_para(tf_rel, "locatie_id  INT  FK → locatie(id)", size=14, color=ZWART, font_name="Consolas")
add_para(tf_rel, "opvoer      TIMESTAMPTZ  — afgeleid", size=14, color=GRIJS, font_name="Consolas")
add_para(tf_rel, "afvoer      TIMESTAMPTZ  — afgeleid", size=14, color=GRIJS, font_name="Consolas")
add_para(tf_rel, "", size=4)
add_para(tf_rel, "PK: (persoon_id, rel_id)", size=14, bold=True, color=ORANJE)

# Relatie data
reld = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.5), WIT, GROEN, Pt(2))
tf_reld = set_text(reld, "bereikbaarheid_data", size=18, bold=True, color=GROEN)
reld.text_frame.margin_left = Inches(0.3)
reld.text_frame.margin_top = Inches(0.1)
add_para(tf_reld, "", size=4)
add_para(tf_reld, "persoon_id  INT  FK → bereikbaarheid", size=14, color=ZWART, font_name="Consolas")
add_para(tf_reld, "rel_id      INT  FK → bereikbaarheid", size=14, color=ZWART, font_name="Consolas")
add_para(tf_reld, "versie      INT  autoincrement", size=14, color=ZWART, font_name="Consolas")
add_para(tf_reld, "soort       TEXT  (enum: Woonadres|...)", size=14, color=ZWART, font_name="Consolas")
add_para(tf_reld, "opvoer/afvoer — afgeleid", size=14, color=GRIJS, font_name="Consolas")
add_para(tf_reld, "", size=4)
add_para(tf_reld, "PK: (persoon_id, rel_id, versie)", size=14, bold=True, color=GROEN)

# ER-diagram visueel
# Persoon
p_box = add_shape(slide, Inches(1.0), Inches(5.0), Inches(2.0), Inches(0.8), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
set_text(p_box, "Persoon", size=16, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
p_box.text_frame.margin_top = Inches(0.1)

# Bereikbaarheid hub
b_hub = add_shape(slide, Inches(4.5), Inches(5.0), Inches(2.5), Inches(0.8), ORANJE_LICHT, ORANJE, Pt(2))
set_text(b_hub, "Bereikbaarheid", size=14, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
b_hub.text_frame.margin_top = Inches(0.1)
add_para(b_hub.text_frame, "(hub)", size=11, color=GRIJS, align=PP_ALIGN.CENTER)

# Locatie
l_box = add_shape(slide, Inches(8.5), Inches(5.0), Inches(2.0), Inches(0.8), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
set_text(l_box, "Locatie", size=16, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
l_box.text_frame.margin_top = Inches(0.1)

# Bereikbaarheid_Data
bd = add_shape(slide, Inches(4.5), Inches(6.2), Inches(2.5), Inches(0.7), GROEN_LICHT, GROEN)
set_text(bd, "Bereikbaarh_Data", size=12, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
bd.text_frame.margin_top = Inches(0.08)

# Pijlen
add_arrow(slide, Inches(3.0), Inches(5.4), Inches(4.5), Inches(5.4), ORANJE, Pt(2))
add_arrow(slide, Inches(7.0), Inches(5.4), Inches(8.5), Inches(5.4), ORANJE, Pt(2))
add_arrow(slide, Inches(5.75), Inches(5.8), Inches(5.75), Inches(6.2), GROEN, Pt(2))

# Labels
tb = add_textbox(slide, Inches(3.1), Inches(4.7), Inches(1.5), Inches(0.3))
set_text(tb, "1     *", size=13, color=GRIJS, align=PP_ALIGN.CENTER)

tb2 = add_textbox(slide, Inches(7.0), Inches(4.7), Inches(1.5), Inches(0.3))
set_text(tb2, "*     1", size=13, color=GRIJS, align=PP_ALIGN.CENTER)

add_footer(slide, 15)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16: METAREGISTRY
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "MetaRegistry: single source of truth", "Alle metadata op één plek — dynamisch gedreven")

# Hoofdbox
main = add_shape(slide, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.2), PAURS_LICHT if False else PAARS_LICHT, PAARS, Pt(2))
tf_m = set_text(main, "De MetaRegistry is een Go map van TypeMeta entries — één per representatietype (struct).",
                size=18, bold=True, color=PAARS, align=PP_ALIGN.CENTER)
main.text_frame.margin_top = Inches(0.1)
add_para(tf_m, "Handlers, routes, schema-API en frontend lezen alles dynamisch uit de MetaRegistry.", 
         size=15, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

# TypeMeta velden
tm = add_shape(slide, Inches(0.6), Inches(3.1), Inches(5.8), Inches(4.0), WIT, MIDDENBLAUW, Pt(2))
tf_tm = set_text(tm, "TypeMeta velden", size=18, bold=True, color=MIDDENBLAUW)
tm.text_frame.margin_left = Inches(0.25)
tm.text_frame.margin_top = Inches(0.1)
add_para(tf_tm, "", size=4)
add_para(tf_tm, "▸ Typenaam, Description, Metatype", size=14, color=ZWART)
add_para(tf_tm, "▸ IsMaterieel, GESubtype (hub/data/aanvang/einde)", size=14, color=ZWART)
add_para(tf_tm, "▸ Padnaam (URL), Veldnaam (JSON)", size=14, color=ZWART)
add_para(tf_tm, "▸ Factory / SliceFactory (constructors)", size=14, color=ZWART)
add_para(tf_tm, "▸ Tabelnaam, IDKolom, EntiteitIDKolom", size=14, color=ZWART)
add_para(tf_tm, "▸ HeeftPFK, RelatieveAutoincrement", size=14, color=ZWART)
add_para(tf_tm, "▸ OnderliggendeGegevenselementen []", size=14, color=ZWART)
add_para(tf_tm, "▸ Momentvoorkomen (enkelvoudig/meervoudig)", size=14, color=ZWART)
add_para(tf_tm, "▸ BovenliggendTypenaam", size=14, color=ZWART)
add_para(tf_tm, "▸ Kleur (visualisatie)", size=14, color=ZWART)

# Wat het aanstuurt
drives = add_shape(slide, Inches(6.8), Inches(3.1), Inches(5.8), Inches(4.0), WIT, GROEN, Pt(2))
tf_d = set_text(drives, "Wat de MetaRegistry aanstuurt", size=18, bold=True, color=GROEN)
drives.text_frame.margin_left = Inches(0.25)
drives.text_frame.margin_top = Inches(0.1)
add_para(tf_d, "", size=4)
add_para(tf_d, "🔀  Routes — dynamisch geregistreerd via Padnaam", size=15, color=ZWART)
add_para(tf_d, "", size=4)
add_para(tf_d, "⚙️  Handlers — generiek: MakeGetEntityByMetaHandler", size=15, color=ZWART)
add_para(tf_d, "      (meta), MakeAddByMetaHandler(meta), etc.", size=14, color=DONKERGRIJS)
add_para(tf_d, "", size=4)
add_para(tf_d, "📋  Schema-API — /api/viz/schema retourneert", size=15, color=ZWART)
add_para(tf_d, "      metadata + veldtypes + relaties", size=14, color=DONKERGRIJS)
add_para(tf_d, "", size=4)
add_para(tf_d, "🖥️  Frontend — leest dynamisch het schema", size=15, color=ZWART)
add_para(tf_d, "      Geen hardcoded veldnamen!", size=14, bold=True, color=GROEN)
add_para(tf_d, "", size=4)
add_para(tf_d, "📊  Visualisatie — kleuren, layout, oortjes", size=15, color=ZWART)

add_footer(slide, 16)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17: SCHEMA-API & FRONTEND
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Schema-API & dynamische frontend", "Geen hardcoded structuren — alles uit de MetaRegistry")

# Schema-API flow
tb = add_textbox(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5))
set_text(tb, "Flow:  MetaRegistry  →  Schema-API ( /api/viz/schema )  →  Frontend (React/Vite)", 
         size=18, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)

# Schema-API output voorbeeld
api_box = add_shape(slide, Inches(0.4), Inches(2.4), Inches(6.2), Inches(4.6), WIT, MIDDENBLAUW, Pt(2))
tf_api = set_text(api_box, "Schema-API response (voorbeeld)", size=16, bold=True, color=MIDDENBLAUW)
api_box.text_frame.margin_left = Inches(0.2)
api_box.text_frame.margin_top = Inches(0.1)
add_para(tf_api, "", size=4)
add_para(tf_api, '{', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '  "typenaam": "Naam",', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '  "metatype": "gegevenselement",', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '  "ge_subtype": "hub",', size=12, color=MIDDENBLAUW, font_name="Consolas", bold=True)
add_para(tf_api, '  "is_materieel": false,', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '  "padnaam": "naam",', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '  "velden": [', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '    {"naam":"achternaam","type":"string"},', size=12, color=GROEN, font_name="Consolas")
add_para(tf_api, '    {"naam":"naamgebruik","type":"string",', size=12, color=GROEN, font_name="Consolas")
add_para(tf_api, '     "enum":["EigenNaam","PartnerNaam",...]}', size=12, color=GROEN, font_name="Consolas")
add_para(tf_api, '  ],', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '  "onderliggende": [', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '    {"rolnaam":"data","doeltype":"Naam_Data"}', size=12, color=ORANJE, font_name="Consolas")
add_para(tf_api, '  ]', size=12, color=ZWART, font_name="Consolas")
add_para(tf_api, '}', size=12, color=ZWART, font_name="Consolas")

# Frontend output
fe_box = add_shape(slide, Inches(6.8), Inches(2.4), Inches(6.2), Inches(4.6), WIT, GROEN, Pt(2))
tf_fe = set_text(fe_box, "Frontend: dynamisch gegenereerd", size=16, bold=True, color=GROEN)
fe_box.text_frame.margin_left = Inches(0.25)
fe_box.text_frame.margin_top = Inches(0.1)
add_para(tf_fe, "", size=4)
add_para(tf_fe, "De frontend leest het schema en genereert:", size=15, color=ZWART)
add_para(tf_fe, "", size=4)
add_para(tf_fe, "▸  Formulieren — veldtype bepaalt inputcomponent", size=14, color=DONKERGRIJS)
add_para(tf_fe, "   string → tekstveld, date → datepicker,", size=13, color=GRIJS)
add_para(tf_fe, "   enum → dropdown, bool → checkbox", size=13, color=GRIJS)
add_para(tf_fe, "", size=4)
add_para(tf_fe, "▸  Overzichten — kolommen uit schema-velden", size=14, color=DONKERGRIJS)
add_para(tf_fe, "", size=4)
add_para(tf_fe, "▸  Visualisatie (SVG) — entiteitskaarten,", size=14, color=DONKERGRIJS)
add_para(tf_fe, "   GE-boxes, relatielijnen, aanvang/einde-oortjes", size=13, color=GRIJS)
add_para(tf_fe, "", size=4)
add_para(tf_fe, "▸  Registratiepagina — inzien van audit-trail", size=14, color=DONKERGRIJS)
add_para(tf_fe, "", size=8)
add_para(tf_fe, "Nieuw type toevoegen = alleen MetaRegistry", size=15, bold=True, color=GROEN)
add_para(tf_fe, "aanpassen — frontend past zich automatisch aan", size=15, bold=True, color=GROEN)

add_footer(slide, 17)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18: TECH STACK
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Tech stack & architectuur", "Van MetaRegistry tot browser")

# Layers
layers = [
    ("Frontend", "React + Vite", "Dynamische UI, SVG-visualisatie", GROEN_LICHT, GROEN, Inches(1.5)),
    ("Schema-API", "/api/viz/schema", "Metadata + veldtypes als JSON", BLAUW_LICHT, MIDDENBLAUW, Inches(2.5)),
    ("Handlers", "Generieke Go handlers", "MakeGet…ByMetaHandler(meta)", PAARS_LICHT, PAARS, Inches(3.5)),
    ("MetaRegistry", "Go map[string]TypeMeta", "Single source of truth", GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Inches(4.5)),
    ("ORM", "Bun (Go)", "DB-mapping via struct tags", ORANJE_LICHT, ORANJE, Inches(5.5)),
    ("Database", "PostgreSQL", "Tabellen, FK's, relatieve autoincrement", RGBColor(0xE0, 0xE7, 0xFF), MIDDENBLAUW, Inches(6.5)),
]

for name, tech, desc, bg_clr, line_clr, y in layers:
    box = add_shape(slide, Inches(1.5), y, Inches(10.5), Inches(0.8), bg_clr, line_clr, Pt(2))
    tf = set_text(box, name, size=18, bold=True, color=line_clr)
    box.text_frame.margin_left = Inches(0.3)
    box.text_frame.margin_top = Inches(0.05)
    
    # Tech en desc rechts
    tb_tech = add_textbox(slide, Inches(4.5), y + Inches(0.05), Inches(3.5), Inches(0.7))
    set_text(tb_tech, tech, size=15, bold=True, color=ZWART)
    
    tb_desc = add_textbox(slide, Inches(8.0), y + Inches(0.05), Inches(4.0), Inches(0.7))
    set_text(tb_desc, desc, size=14, color=GRIJS)

# Pijlen tussen lagen
for i in range(len(layers) - 1):
    y1 = layers[i][5] + Inches(0.8)
    y2 = layers[i+1][5]
    add_arrow(slide, Inches(6.75), y1, Inches(6.75), y2, GRIJS, Pt(2))

# HTTP label
tb_gin = add_textbox(slide, Inches(0.2), Inches(3.0), Inches(1.3), Inches(0.5))
set_text(tb_gin, "Gin HTTP", size=13, bold=True, color=GRIJS, align=PP_ALIGN.CENTER)

add_footer(slide, 18)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 19: WERKWIJZE BIJ MODELWIJZIGINGEN
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Werkwijze bij modelwijzigingen", "Wat moet je aanpassen als het model verandert?")

steps = [
    ("1", "Struct definiëren", "Go struct in modellen_entiteiten.go of modellen_ge_rel.go\nmet JSON-tags, Bun-tags en schema-tags"),
    ("2", "Interface-methoden", "GetID(), Metatype(), ClearID(), Get/SetOpvoer(), Get/SetAfvoer()\nimplementeren op de struct"),
    ("3", "MetaRegistry-entry", "TypeMeta entry toevoegen in metaregistry.go\nmet alle metadata (metatype, padnaam, factory, etc.)"),
    ("4", "Onderliggende GE's", "Bij entiteiten: OnderliggendeGegevenselementen aanpassen\nRolnaam + JSONRolnaam matchen met struct-veld"),
    ("5", "DB-tabel aanmaken", "CREATE TABLE in dbsetup/createtables.go\nmet juiste PK, FK's en kolommen"),
    ("6", "Klaar!", "Routes en handlers worden automatisch gegenereerd.\nFrontend past zich aan via de schema-API."),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.5) + Inches(0.95) * i
    
    # Stap nummer
    clr = GROEN if num == "6" else LICHTBLAUW
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = clr
    circle.line.fill.background()
    set_text(circle, num, size=20, bold=True, color=WIT, align=PP_ALIGN.CENTER)
    circle.text_frame.margin_left = Inches(0)
    circle.text_frame.margin_right = Inches(0)
    circle.text_frame.margin_top = Inches(0.05)
    
    # Title
    tb_t = add_textbox(slide, Inches(1.6), y, Inches(3.5), Inches(0.4))
    set_text(tb_t, title, size=18, bold=True, color=DONKERBLAUW)
    
    # Description
    tb_d = add_textbox(slide, Inches(5.2), y, Inches(7.5), Inches(0.8))
    tf = set_text(tb_d, desc.split('\n')[0], size=14, color=DONKERGRIJS)
    if '\n' in desc:
        add_para(tf, desc.split('\n')[1], size=13, color=GRIJS)

add_footer(slide, 19)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 20: SAMENVATTING
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DONKERBLAUW)

# Accent lijn
add_rect(slide, Inches(0), Inches(1.0), Inches(13.333), Inches(0.04), LICHTBLAUW)

tb_title = add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.7))
set_text(tb_title, "Samenvatting", size=36, bold=True, color=WIT, align=PP_ALIGN.CENTER, font_name="Calibri Light")

# Kernpunten in twee kolommen
col1 = add_shape(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4), RGBColor(0x25, 0x4E, 0x7A), line_color=None)
tf1 = set_text(col1, "Principes", size=24, bold=True, color=ZACHTBLAUW)
col1.text_frame.margin_left = Inches(0.3)
col1.text_frame.margin_top = Inches(0.15)
add_para(tf1, "", size=6)
add_para(tf1, "✓  Twee tijdsdimensies:", size=17, color=WIT)
add_para(tf1, "     formeel (wanneer geregistreerd) +", size=16, color=ZACHTBLAUW)
add_para(tf1, "     materieel (wanneer geldig)", size=16, color=ZACHTBLAUW)
add_para(tf1, "", size=6)
add_para(tf1, "✓  Opvoer/afvoer zijn altijd afgeleid", size=17, color=WIT)
add_para(tf1, "     Bron van waarheid: registratie + wijzigingen", size=16, color=ZACHTBLAUW)
add_para(tf1, "", size=6)
add_para(tf1, "✓  Drie metatypes:", size=17, color=WIT)
add_para(tf1, "     Entiteit → Gegevenselement → Relatie", size=16, color=ZACHTBLAUW)
add_para(tf1, "", size=6)
add_para(tf1, "✓  Hub + Data pattern:", size=17, color=WIT)
add_para(tf1, "     Identiteit, inhoud en materialiteit", size=16, color=ZACHTBLAUW)
add_para(tf1, "     apart corrigeerbaar", size=16, color=ZACHTBLAUW)

col2 = add_shape(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.4), RGBColor(0x25, 0x4E, 0x7A), line_color=None)
tf2 = set_text(col2, "Implementatie", size=24, bold=True, color=ZACHTBLAUW)
col2.text_frame.margin_left = Inches(0.3)
col2.text_frame.margin_top = Inches(0.15)
add_para(tf2, "", size=6)
add_para(tf2, "✓  MetaRegistry als single source of truth", size=17, color=WIT)
add_para(tf2, "     Handlers, routes, schema-API dynamisch", size=16, color=ZACHTBLAUW)
add_para(tf2, "", size=6)
add_para(tf2, "✓  Schema-driven frontend", size=17, color=WIT)
add_para(tf2, "     Geen hardcoded veldnamen of structuren", size=16, color=ZACHTBLAUW)
add_para(tf2, "", size=6)
add_para(tf2, "✓  PostgreSQL met relatieve autoincrement", size=17, color=WIT)
add_para(tf2, "     Samengestelde PK's per scope", size=16, color=ZACHTBLAUW)
add_para(tf2, "", size=6)
add_para(tf2, "✓  Tijdreizen over twee dimensies", size=17, color=WIT)
add_para(tf2, "     Foutloze audit-trail, zelfs cross-register", size=16, color=ZACHTBLAUW)
add_para(tf2, "", size=10)
add_para(tf2, "Go  ·  Gin  ·  Bun  ·  PostgreSQL  ·  React/Vite", size=15, color=GRIJS, align=PP_ALIGN.CENTER)

add_footer(slide, 20)


# ── Opslaan ──────────────────────────────────────────────────────────
output_path = "bitemporeel_register_presentatie.pptx"
prs.save(output_path)
print(f"Presentatie opgeslagen als: {output_path}")
print(f"Aantal slides: {len(prs.slides)}")
