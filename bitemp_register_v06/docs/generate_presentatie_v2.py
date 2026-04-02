"""
Script om de PowerPoint-presentatie "Bitemporeel Register" te genereren.
Versie 2 — met Mermaid, roundtrip codegen, import/export, validaties,
            afgeleide velden en referentielijsten.
Uitvoeren: python generate_presentatie_v2.py
Output:    bitemporeel_register_presentatie.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Kleuren ──────────────────────────────────────────────────────────
DONKERBLAUW  = RGBColor(0x1E, 0x3A, 0x5F)
MIDDENBLAUW  = RGBColor(0x2C, 0x5F, 0x8A)
LICHTBLAUW   = RGBColor(0x3B, 0x82, 0xF6)
ZACHTBLAUW   = RGBColor(0xBF, 0xDB, 0xFE)
ACHTERGROND  = RGBColor(0xF0, 0xF4, 0xF8)
WIT          = RGBColor(0xFF, 0xFF, 0xFF)
ZWART        = RGBColor(0x20, 0x20, 0x20)
GRIJS        = RGBColor(0x64, 0x74, 0x8B)
DONKERGRIJS  = RGBColor(0x47, 0x55, 0x69)
GROEN        = RGBColor(0x16, 0xA3, 0x4A)
ORANJE       = RGBColor(0xEA, 0x58, 0x0C)
ROOD         = RGBColor(0xDC, 0x26, 0x26)
GEEL_LICHT   = RGBColor(0xFE, 0xF3, 0xC7)
GROEN_LICHT  = RGBColor(0xDC, 0xFC, 0xE7)
BLAUW_LICHT  = RGBColor(0xDB, 0xEA, 0xFE)
ROOD_LICHT   = RGBColor(0xFE, 0xE2, 0xE2)
ORANJE_LICHT = RGBColor(0xFF, 0xED, 0xD5)
PAARS        = RGBColor(0x7C, 0x3A, 0xED)
PAARS_LICHT  = RGBColor(0xED, 0xE9, 0xFE)
TEAL         = RGBColor(0x0D, 0x94, 0x88)
TEAL_LICHT   = RGBColor(0xCC, 0xFB, 0xF1)
AMBER        = RGBColor(0xD9, 0x77, 0x06)
AMBER_LICHT  = RGBColor(0xFE, 0xF3, 0xC7)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_NUM = [0]  # mutable counter
TOTAL_SLIDES = 26

# ── Helpers ──────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def set_text(shape, text, size=18, bold=False, color=ZWART, align=PP_ALIGN.LEFT, font_name="Calibri"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=18, bold=False, color=ZWART, align=PP_ALIGN.LEFT, font_name="Calibri",
             space_before=Pt(4), space_after=Pt(2), level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def add_title_bar(slide, title_text, subtitle_text=None):
    bar = add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DONKERBLAUW)
    set_text(bar, title_text, size=32, bold=True, color=WIT, align=PP_ALIGN.LEFT)
    bar.text_frame.paragraphs[0].font.name = "Calibri Light"
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.margin_top = Inches(0.15)
    if subtitle_text:
        add_para(bar.text_frame, subtitle_text, size=18, color=ZACHTBLAUW, align=PP_ALIGN.LEFT)
    return bar

def new_slide():
    SLIDE_NUM[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Footer
    tb = add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    set_text(tb, f"{SLIDE_NUM[0]} / {TOTAL_SLIDES}", size=11, color=GRIJS, align=PP_ALIGN.RIGHT)
    return slide

def add_arrow(slide, x1, y1, x2, y2, color=MIDDENBLAUW, width=Pt(2)):
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    ln = connector.line._ln
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return connector

def add_code_box(slide, left, top, width, height, code_lines, title=None, title_color=MIDDENBLAUW):
    """Code-blok met monospace tekst."""
    box = add_shape(slide, left, top, width, height, RGBColor(0x1E, 0x29, 0x3B), line_color=RGBColor(0x33, 0x44, 0x55), line_width=Pt(1))
    box.text_frame.margin_left = Inches(0.2)
    box.text_frame.margin_top = Inches(0.1)
    if title:
        tf = set_text(box, title, size=13, bold=True, color=title_color, font_name="Calibri")
        add_para(tf, "", size=4)
    else:
        tf = set_text(box, "", size=4)
    for line in code_lines:
        add_para(tf, line, size=11, color=RGBColor(0xA5, 0xD6, 0xFF), font_name="Consolas",
                 space_before=Pt(1), space_after=Pt(0))
    return box


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: TITEL
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, DONKERBLAUW)
add_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), LICHTBLAUW)

tb = add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5))
set_text(tb, "Het Bitemporele Register", size=48, bold=True, color=WIT, align=PP_ALIGN.CENTER, font_name="Calibri Light")

tb2 = add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(1.2))
set_text(tb2, "Principes, Metamodel & Architectuur", size=28, color=ZACHTBLAUW, align=PP_ALIGN.CENTER)

tb3 = add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(1.2))
tf3 = set_text(tb3, "Van UML-model tot werkend register", size=18, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf3, "Go  ·  PostgreSQL  ·  React  ·  Mermaid  ·  PlantUML  ·  XMI", size=15, color=GRIJS, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Agenda")

items = [
    ("I",   "Principes", "Waarom bitemporeel? Twee tijdsdimensies, tijdreizen"),
    ("II",  "Metamodel", "UML representatietypes: Entiteit, GE, Relatie"),
    ("III", "Hub + Data", "Drielagen-patroon voor correcties & materialiteit"),
    ("IV",  "Database", "Tabellen, sleutels, relatieve autoincrement"),
    ("V",   "MetaRegistry", "Single source of truth, Schema-API, dynamische frontend"),
    ("VI",  "Roundtrip", "UML Editor ↔ V3 JSON ↔ Code generatie ↔ Register"),
    ("VII", "Import/Export", "Mermaid, PlantUML, XMI (EA), OAS 3.1"),
    ("VIII","Bijzonderheden", "Validaties, afgeleide velden, referentielijsten"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.5) + Inches(0.7) * i
    circle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(0.7), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = LICHTBLAUW
    circle.line.fill.background()
    set_text(circle, num, size=14, bold=True, color=WIT, align=PP_ALIGN.CENTER)
    circle.text_frame.margin_left = Inches(0)
    circle.text_frame.margin_right = Inches(0)
    circle.text_frame.margin_top = Inches(0.05)

    tb = add_textbox(slide, Inches(1.8), y, Inches(3.2), Inches(0.5))
    set_text(tb, title, size=19, bold=True, color=DONKERBLAUW)
    
    tb2 = add_textbox(slide, Inches(5.2), y + Inches(0.03), Inches(7.5), Inches(0.5))
    set_text(tb2, desc, size=15, color=GRIJS)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: WAAROM BITEMPOREEL?
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Waarom bitemporeel?", "Het probleem dat we oplossen")

box = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2), ROOD_LICHT, ROOD, Pt(2))
tf = set_text(box, "Het probleem", size=22, bold=True, color=ROOD)
box.text_frame.margin_left = Inches(0.3)
box.text_frame.margin_top = Inches(0.2)
add_para(tf, "", size=8)
add_para(tf, "Traditionele databases kennen geen tijdsbesef:", size=17, color=ZWART)
add_para(tf, "▸  UPDATE overschrijft de vorige waarde", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "▸  Geen audittrail: wat was de waarde gisteren?", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "▸  Geen onderscheid: wanneer geregistreerd", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "   vs. wanneer geldig in werkelijkheid", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "▸  Correcties vernietigen de oorspronkelijke data", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "▸  Onmogelijk om historische beslissingen", size=16, color=DONKERGRIJS, level=1)
add_para(tf, "   te reconstrueren", size=16, color=DONKERGRIJS, level=1)

box2 = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), GROEN_LICHT, GROEN, Pt(2))
tf2 = set_text(box2, "De oplossing: bitemporeel", size=22, bold=True, color=GROEN)
box2.text_frame.margin_left = Inches(0.3)
box2.text_frame.margin_top = Inches(0.2)
add_para(tf2, "", size=8)
add_para(tf2, "Twee tijdsassen bewaren alles:", size=17, color=ZWART)
add_para(tf2, "▸  Niets wordt ooit overschreven", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Volledige audittrail van elke wijziging", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Tijdreizen naar elk moment in het verleden", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Correcties behouden het oorspronkelijke pad", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Juridisch verifieerbare registratie", size=16, color=DONKERGRIJS, level=1)
add_para(tf2, "▸  Toekomstige geldigheid registreerbaar", size=16, color=DONKERGRIJS, level=1)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: TWEE TIJDSDIMENSIES
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Twee tijdsdimensies", "Het fundament van bitemporele registratie")

box_f = add_shape(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
tf = set_text(box_f, "⏱  Formele tijd", size=26, bold=True, color=MIDDENBLAUW)
box_f.text_frame.margin_left = Inches(0.3)
box_f.text_frame.margin_top = Inches(0.2)
add_para(tf, "(registratietijd)", size=16, color=GRIJS)
add_para(tf, "", size=8)
add_para(tf, "Wanneer is iets geregistreerd?", size=18, bold=True, color=ZWART)
add_para(tf, "", size=6)
add_para(tf, "▸  Registratietijdstip (timestamp)", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Opvoer = record komt op de tijdlijn", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Afvoer = record gaat van de tijdlijn", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Alleen tijdreizen naar het verleden", size=16, color=DONKERGRIJS)
add_para(tf, "", size=8)
add_para(tf, "Vastgelegd via: registratie → wijzigingen", size=15, bold=True, color=MIDDENBLAUW)

box_m = add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8), ORANJE_LICHT, ORANJE, Pt(2))
tf2 = set_text(box_m, "📅  Materiële tijd", size=26, bold=True, color=ORANJE)
box_m.text_frame.margin_left = Inches(0.3)
box_m.text_frame.margin_top = Inches(0.2)
add_para(tf2, "(geldigheidstijd)", size=16, color=GRIJS)
add_para(tf2, "", size=8)
add_para(tf2, "Wanneer geldt iets in de werkelijkheid?", size=18, bold=True, color=ZWART)
add_para(tf2, "", size=6)
add_para(tf2, "▸  Aanvangsdatum (date)", size=16, color=DONKERGRIJS)
add_para(tf2, "▸  Einddatum (date)", size=16, color=DONKERGRIJS)
add_para(tf2, "▸  Tijdreizen naar verleden én toekomst", size=16, color=DONKERGRIJS)
add_para(tf2, "▸  Bijv. voorgenomen verhuizing", size=16, color=DONKERGRIJS)
add_para(tf2, "", size=8)
add_para(tf2, "Vastgelegd via: aanvang / einde records", size=15, bold=True, color=ORANJE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: FORMELE TIJD VISUEEL
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Formele tijd: registratie → wijzigingen", "Hoe de formele tijdlijn ontstaat")

tb = add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(5.5))
tf = set_text(tb, "Kernprincipe", size=22, bold=True, color=DONKERBLAUW)
add_para(tf, "", size=6)
add_para(tf, "Opvoer en afvoer staan niet in de inhoudelijke", size=17, color=ZWART)
add_para(tf, "tabellen, maar in een aparte wijzigingen-tabel.", size=17, color=ZWART)
add_para(tf, "", size=10)
add_para(tf, "Eén registratie = één formeel moment", size=18, bold=True, color=MIDDENBLAUW)
add_para(tf, "", size=6)
add_para(tf, "▸  Heeft een tijdstip (timestamp)", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Bevat 1..n wijzigingen", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Elke wijziging = opvoer of afvoer", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Wijst naar een specifiek record", size=16, color=DONKERGRIJS)
add_para(tf, "", size=10)
add_para(tf, "Soorten registratie:", size=18, bold=True, color=MIDDENBLAUW)
add_para(tf, "▸  Registratie — gewone vastlegging", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Correctie — corrigeert eerdere registratie", size=16, color=DONKERGRIJS)
add_para(tf, "▸  Ongedaanmaking — maakt registratie ongeldig", size=16, color=DONKERGRIJS)

# Diagram rechts
reg = add_shape(slide, Inches(7.2), Inches(1.8), Inches(2.8), Inches(1.4), MIDDENBLAUW)
tf_r = set_text(reg, "Registratie", size=20, bold=True, color=WIT, align=PP_ALIGN.CENTER)
reg.text_frame.margin_top = Inches(0.1)
add_para(tf_r, "tijdstip: 2026-03-21 14:30", size=13, color=ZACHTBLAUW, align=PP_ALIGN.CENTER)
add_para(tf_r, "type: Registratie", size=13, color=ZACHTBLAUW, align=PP_ALIGN.CENTER)

w1 = add_shape(slide, Inches(6.6), Inches(3.8), Inches(2.2), Inches(1.1), BLAUW_LICHT, MIDDENBLAUW)
tf_w1 = set_text(w1, "Wijziging 1", size=15, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
w1.text_frame.margin_top = Inches(0.05)
add_para(tf_w1, "type: Opvoer", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_w1, "record: Persoon #5", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

w2 = add_shape(slide, Inches(9.2), Inches(3.8), Inches(2.2), Inches(1.1), BLAUW_LICHT, MIDDENBLAUW)
tf_w2 = set_text(w2, "Wijziging 2", size=15, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
w2.text_frame.margin_top = Inches(0.05)
add_para(tf_w2, "type: Opvoer", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_w2, "record: Naam #1", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

r1 = add_shape(slide, Inches(6.6), Inches(5.5), Inches(2.2), Inches(0.9), GROEN_LICHT, GROEN)
set_text(r1, "Persoon #5", size=15, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
r1.text_frame.margin_top = Inches(0.08)
add_para(r1.text_frame, "BSN: 123456782", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

r2 = add_shape(slide, Inches(9.2), Inches(5.5), Inches(2.2), Inches(0.9), GROEN_LICHT, GROEN)
set_text(r2, "Naam v1", size=15, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
r2.text_frame.margin_top = Inches(0.08)
add_para(r2.text_frame, '"De Vries"', size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

add_arrow(slide, Inches(8.6), Inches(3.2), Inches(7.7), Inches(3.8), MIDDENBLAUW)
add_arrow(slide, Inches(8.6), Inches(3.2), Inches(10.3), Inches(3.8), MIDDENBLAUW)
add_arrow(slide, Inches(7.7), Inches(4.9), Inches(7.7), Inches(5.5), GROEN)
add_arrow(slide, Inches(10.3), Inches(4.9), Inches(10.3), Inches(5.5), GROEN)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: OPVOER & AFVOER
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Opvoer & afvoer: altijd afgeleid", "De bron van waarheid is de wijzigingen-tabel")

box = add_shape(slide, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.2), GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Pt(2))
tf = set_text(box, "💡  Kernprincipe: opvoer en afvoer in records zijn afgeleide waarden — niet de bron van waarheid.", 
              size=19, bold=True, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)
box.text_frame.margin_top = Inches(0.15)
add_para(tf, "De werkelijke formele tijdlijn wordt bepaald door de registraties en hun wijzigingen.", 
         size=16, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)

lbox = add_shape(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8), WIT, MIDDENBLAUW)
tf_l = set_text(lbox, "Bron van waarheid", size=20, bold=True, color=MIDDENBLAUW)
lbox.text_frame.margin_left = Inches(0.3)
lbox.text_frame.margin_top = Inches(0.15)
add_para(tf_l, "", size=6)
add_para(tf_l, "registratie.tijdstip", size=16, bold=True, color=ZWART)
add_para(tf_l, "  + wijziging.type (opvoer/afvoer)", size=15, color=DONKERGRIJS)
add_para(tf_l, "  + wijziging.representatie_id", size=15, color=DONKERGRIJS)
add_para(tf_l, "", size=8)
add_para(tf_l, "Bij formeel tijdreizen (→ peiltijdstip tf):", size=16, bold=True, color=ZWART)
add_para(tf_l, "▸  Records zelf worden NIET gebruikt", size=15, color=DONKERGRIJS)
add_para(tf_l, "▸  Toestand wordt opnieuw afgeleid", size=15, color=DONKERGRIJS)
add_para(tf_l, "▸  Alle wijzigingen t/m tf verwerken", size=15, color=DONKERGRIJS)

rbox = add_shape(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8), WIT, GRIJS)
tf_r = set_text(rbox, "Afgeleide velden in records", size=20, bold=True, color=DONKERGRIJS)
rbox.text_frame.margin_left = Inches(0.3)
rbox.text_frame.margin_top = Inches(0.15)
add_para(tf_r, "", size=6)
add_para(tf_r, "record.opvoer  = afgeleid (tf = nu)", size=16, bold=True, color=ZWART)
add_para(tf_r, "record.afvoer  = afgeleid (tf = nu)", size=16, bold=True, color=ZWART)
add_para(tf_r, "", size=8)
add_para(tf_r, "Geldt voor ALLE lagen:", size=16, bold=True, color=ZWART)
add_para(tf_r, "▸  Entiteit  (Persoon, Locatie)", size=15, color=DONKERGRIJS)
add_para(tf_r, "▸  GE/REL hub  (Naam, Bereikbaarheid)", size=15, color=DONKERGRIJS)
add_para(tf_r, "▸  _Data  (Naam_Data)", size=15, color=DONKERGRIJS)
add_para(tf_r, "▸  _Aanvang / _Einde  (materieel)", size=15, color=DONKERGRIJS)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: TIJDREIZEN
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Tijdreizen", "De kracht van bitemporele registratie")

box_ft = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.4), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
tf_ft = set_text(box_ft, "⏪  Formeel tijdreizen", size=22, bold=True, color=MIDDENBLAUW)
box_ft.text_frame.margin_left = Inches(0.3)
box_ft.text_frame.margin_top = Inches(0.15)
add_para(tf_ft, "", size=4)
add_para(tf_ft, "Wat was bekend op tijdstip tf?", size=17, bold=True, color=ZWART)
add_para(tf_ft, "▸  Alleen naar het verleden (tf ≤ nu)", size=15, color=DONKERGRIJS)
add_para(tf_ft, "▸  API: ?t=2024-01-01T12:00:00Z", size=15, color=DONKERGRIJS)
add_para(tf_ft, "▸  Verwerk alle wijzigingen t/m tf", size=15, color=DONKERGRIJS)

box_mt = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.4), ORANJE_LICHT, ORANJE, Pt(2))
tf_mt = set_text(box_mt, "⏩  Materieel tijdreizen", size=22, bold=True, color=ORANJE)
box_mt.text_frame.margin_left = Inches(0.3)
box_mt.text_frame.margin_top = Inches(0.15)
add_para(tf_mt, "", size=4)
add_para(tf_mt, "Wat gold op datum tm?", size=17, bold=True, color=ZWART)
add_para(tf_mt, "▸  Verleden én toekomst (tm is vrij)", size=15, color=DONKERGRIJS)
add_para(tf_mt, "▸  Aanvang ≤ tm < einde", size=15, color=DONKERGRIJS)
add_para(tf_mt, "▸  Bijv. toekomstige verhuizing al registreerbaar", size=15, color=DONKERGRIJS)

box_c = add_shape(slide, Inches(2.5), Inches(4.4), Inches(8.4), Inches(2.6), PAARS_LICHT, PAARS, Pt(2))
tf_c = set_text(box_c, "🔀  Gecombineerd tijdreizen", size=22, bold=True, color=PAARS)
box_c.text_frame.margin_left = Inches(0.3)
box_c.text_frame.margin_top = Inches(0.15)
add_para(tf_c, "", size=4)
add_para(tf_c, "Wat was op registratietijdstip tf bekend over geldigheidsdatum tm?", size=17, bold=True, color=ZWART)
add_para(tf_c, "", size=4)
add_para(tf_c, "▸  Eerst formeel: welke records waren bekend op tf?", size=15, color=DONKERGRIJS)
add_para(tf_c, "▸  Dan materieel: welke daarvan golden op tm?", size=15, color=DONKERGRIJS)
add_para(tf_c, "▸  Maakt volledige audit-trail mogelijk over twee dimensies", size=15, color=DONKERGRIJS)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: UML METAMODEL + MERMAID
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "UML Metamodel: representatietypes", "De drie bouwstenen van het register — met Mermaid-broncode")

# Visueel diagram links
rep = add_shape(slide, Inches(0.6), Inches(1.5), Inches(3.5), Inches(1.1), PAARS_LICHT, PAARS, Pt(2))
tf_rep = set_text(rep, "《abstract》Representatie", size=16, bold=True, color=PAARS, align=PP_ALIGN.CENTER)
rep.text_frame.margin_top = Inches(0.05)
add_para(tf_rep, "opvoer | afvoer | naam | metatype", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

ent = add_shape(slide, Inches(0.3), Inches(3.0), Inches(2.0), Inches(1.8), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
tf_e = set_text(ent, "Entiteit", size=18, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
ent.text_frame.margin_top = Inches(0.05)
add_para(tf_e, "+ id : integer", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_e, "", size=4)
add_para(tf_e, "Zelfstandig", size=12, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf_e, "identificeerbaar", size=12, color=GRIJS, align=PP_ALIGN.CENTER)

ge = add_shape(slide, Inches(2.5), Inches(3.0), Inches(2.0), Inches(1.8), GROEN_LICHT, GROEN, Pt(2))
tf_g = set_text(ge, "Gegevens-", size=18, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge.text_frame.margin_top = Inches(0.05)
add_para(tf_g, "element", size=18, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
add_para(tf_g, "+ ent_id + rel_id", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_g, "Compositie bij ENT", size=12, color=GRIJS, align=PP_ALIGN.CENTER)

rel_b = add_shape(slide, Inches(2.5), Inches(5.2), Inches(2.0), Inches(1.5), ORANJE_LICHT, ORANJE, Pt(2))
tf_r = set_text(rel_b, "Relatie", size=18, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
rel_b.text_frame.margin_top = Inches(0.05)
add_para(tf_r, "+ doel_id : FK", size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_r, "Koppelt 2 ENTs", size=12, color=GRIJS, align=PP_ALIGN.CENTER)

# Inheritance arrows
add_arrow(slide, Inches(1.3), Inches(3.0), Inches(2.0), Inches(2.6), GRIJS, Pt(2))
add_arrow(slide, Inches(3.5), Inches(3.0), Inches(3.0), Inches(2.6), GRIJS, Pt(2))
# Relatie erft van GE
add_arrow(slide, Inches(3.5), Inches(5.2), Inches(3.5), Inches(4.8), GRIJS, Pt(2))
# Compositie label
tb_comp = add_textbox(slide, Inches(0.3), Inches(5.0), Inches(2.0), Inches(0.3))
set_text(tb_comp, "ENT ◆── 0..* GE", size=12, color=GRIJS, align=PP_ALIGN.CENTER)

# Mermaid broncode rechts
mermaid_lines = [
    "classDiagram",
    "  direction TB",
    "",
    '  class Representatie {',
    '    <<abstract>>',
    '    opvoer : timestamp',
    '    afvoer : timestamp',
    '  }',
    "",
    '  class Entiteit {',
    '    +id : integer',
    '  }',
    "",
    '  class Gegevenselement {',
    '    +ent_id : FK',
    '    +rel_id : PFK',
    '  }',
    "",
    '  class Relatie {',
    '    +doel_id : FK',
    '  }',
    "",
    '  Representatie <|-- Entiteit',
    '  Representatie <|-- Gegevenselement',
    '  Gegevenselement <|-- Relatie',
    '  Entiteit "1" *-- "0..*" Gegevenselement',
]
add_code_box(slide, Inches(5.2), Inches(1.4), Inches(7.6), Inches(5.5),
             mermaid_lines, title="Mermaid classDiagram — exporteerbaar vanuit de UML-editor", title_color=LICHTBLAUW)

# Label
tb_note = add_textbox(slide, Inches(5.2), Inches(6.9), Inches(7.6), Inches(0.4))
set_text(tb_note, "↑ De editor exporteert naar Mermaid, PlantUML en XMI — deze code is direct bruikbaar in documentatie", 
         size=12, color=GRIJS, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: CONCREET UML + MERMAID VOORBEELD
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "UML concreet voorbeeld", "NatuurlijkPersoon ↔ Bereikbaarheid ↔ Locatie — in Mermaid-formaat")

# Visueel links
p = add_shape(slide, Inches(0.4), Inches(1.6), Inches(2.4), Inches(0.9), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
set_text(p, "《ent, materieel》", size=10, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
p.text_frame.margin_top = Inches(0.02)
add_para(p.text_frame, "NatuurlijkPersoon", size=16, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)

ge1 = add_shape(slide, Inches(0.2), Inches(2.8), Inches(2.0), Inches(0.7), GROEN_LICHT, GROEN)
set_text(ge1, "PersoonsIdentificatie", size=11, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge1.text_frame.margin_top = Inches(0.02)
add_para(ge1.text_frame, "bsn, ingezetene", size=10, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

ge2 = add_shape(slide, Inches(0.2), Inches(3.7), Inches(2.0), Inches(0.7), GROEN_LICHT, GROEN)
set_text(ge2, "Naam", size=11, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge2.text_frame.margin_top = Inches(0.02)
add_para(ge2.text_frame, "voorletters, achternaam", size=10, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

ge3 = add_shape(slide, Inches(2.5), Inches(2.8), Inches(2.0), Inches(0.7), GROEN_LICHT, GROEN)
set_text(ge3, "🕐 Burgerschap", size=11, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
ge3.text_frame.margin_top = Inches(0.02)
add_para(ge3.text_frame, "landcode, nationaliteit", size=10, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

rel_box = add_shape(slide, Inches(2.5), Inches(3.7), Inches(2.0), Inches(0.7), ORANJE_LICHT, ORANJE)
set_text(rel_box, "🕐 Bereikbaarheid", size=11, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
rel_box.text_frame.margin_top = Inches(0.02)
add_para(rel_box.text_frame, "soort (enum)", size=10, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

loc = add_shape(slide, Inches(2.5), Inches(4.6), Inches(2.0), Inches(0.7), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
set_text(loc, "《ent》Locatie", size=11, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
loc.text_frame.margin_top = Inches(0.08)

# Pijlen
add_arrow(slide, Inches(1.2), Inches(2.5), Inches(1.2), Inches(2.8), GROEN)
add_arrow(slide, Inches(1.2), Inches(2.5), Inches(1.2), Inches(3.7), GROEN)
add_arrow(slide, Inches(2.6), Inches(2.5), Inches(3.5), Inches(2.8), GROEN)
add_arrow(slide, Inches(2.6), Inches(2.5), Inches(3.5), Inches(3.7), ORANJE)
add_arrow(slide, Inches(3.5), Inches(4.4), Inches(3.5), Inches(4.6), ORANJE)

# Mermaid code rechts
mermaid_concrete = [
    "classDiagram",
    '  class NatuurlijkPersoon {',
    '    <<entiteit, materieel>>',
    '    +id : integer',
    '    /weergavenaam : string',
    '  }',
    '  class PersoonsIdentificatie {',
    '    <<GE>> bsn : string(BSN)',
    '    ingezetene : bool',
    '  }',
    '  class Naam {',
    '    <<GE, meervoudig>>',
    '    voorletters, roepnaam : string',
    '    achternaam : string',
    '    naamgebruik : Naamgebruiksoort',
    '  }',
    '  class Burgerschap {',
    '    <<GE, materieel>>',
    '    landcode : string',
    '    nationaliteit : string',
    '  }',
    '  class Bereikbaarheid {',
    '    <<relatie, materieel>>',
    '    soort : Bereikbaarheidsoort',
    '  }',
    '  class Locatie {',
    '    <<entiteit, materieel>>',
    '    +id : integer',
    '  }',
    '  NatuurlijkPersoon "1" *-- "0..1" PersoonsId',
    '  NatuurlijkPersoon "1" *-- "0..*" Naam',
    '  NatuurlijkPersoon "1" *-- "0..*" Burgerschap',
    '  NatuurlijkPersoon "1" --> "0..*" Bereikbaarheid',
    '  Bereikbaarheid --> "1" Locatie',
]
add_code_box(slide, Inches(5.0), Inches(1.4), Inches(8.0), Inches(5.5),
             mermaid_concrete, title="Mermaid classDiagram — volledig voorbeeld", title_color=LICHTBLAUW)

# Legenda
leg = add_shape(slide, Inches(0.2), Inches(5.6), Inches(4.5), Inches(1.0), WIT, GRIJS)
tf_leg = set_text(leg, "🟦 Entiteit  🟩 GE  🟧 Relatie  🕐 Materieel  / = afgeleid veld", 
                  size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
leg.text_frame.margin_top = Inches(0.1)
add_para(tf_leg, "Velden met / prefix zijn afgeleide velden (UML-conventie)", size=11, color=GRIJS, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: HUB + DATA PATTERN
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Hub + Data Pattern", "Waarom opsplitsing nodig is")

box_p = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.8), ROOD_LICHT, ROOD, Pt(2))
tf_p = set_text(box_p, "Probleem (alles in één tabel)", size=20, bold=True, color=ROOD)
box_p.text_frame.margin_left = Inches(0.3)
box_p.text_frame.margin_top = Inches(0.15)
add_para(tf_p, "", size=4)
add_para(tf_p, "GE-tabel bevat alles in één record:", size=16, color=ZWART)
add_para(tf_p, "  structurele FK's + inhoud + opvoer/afvoer", size=15, color=DONKERGRIJS)
add_para(tf_p, "", size=6)
add_para(tf_p, "▸  Correctie van inhoud → heel record vervangen", size=15, color=DONKERGRIJS)
add_para(tf_p, "▸  Geen stabiel ankerpunt voor materiële tijd", size=15, color=DONKERGRIJS)
add_para(tf_p, "▸  rel_id gaat op bij elke inhoudscorrectie", size=15, color=DONKERGRIJS)

box_o = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8), GROEN_LICHT, GROEN, Pt(2))
tf_o = set_text(box_o, "Oplossing: Hub + Data", size=20, bold=True, color=GROEN)
box_o.text_frame.margin_left = Inches(0.3)
box_o.text_frame.margin_top = Inches(0.15)
add_para(tf_o, "", size=4)
add_para(tf_o, "Splits elk GE/REL-type in lagen:", size=16, color=ZWART)
add_para(tf_o, "", size=6)
add_para(tf_o, "▸  Hub = stabiel identiteitsanker", size=15, color=DONKERGRIJS)
add_para(tf_o, "▸  _Data = geversioned inhoud", size=15, color=DONKERGRIJS)
add_para(tf_o, "▸  _Aanvang/_Einde = materiële tijdlijn", size=15, color=DONKERGRIJS)
add_para(tf_o, "", size=6)
add_para(tf_o, "Elke laag apart corrigeerbaar!", size=16, bold=True, color=GROEN)

# Drie lagen diagram
hub = add_shape(slide, Inches(2.0), Inches(5.0), Inches(2.8), Inches(1.6), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
tf_h = set_text(hub, "Hub", size=20, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
hub.text_frame.margin_top = Inches(0.1)
add_para(tf_h, "(ent_id, rel_id)", size=14, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_h, "Stabiel ankerpunt", size=13, color=GRIJS, align=PP_ALIGN.CENTER)

data = add_shape(slide, Inches(5.3), Inches(5.0), Inches(2.8), Inches(1.6), GROEN_LICHT, GROEN, Pt(2))
tf_d = set_text(data, "_Data", size=20, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
data.text_frame.margin_top = Inches(0.1)
add_para(tf_d, "(ent_id, rel_id, versie)", size=14, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_d, "Geversioned inhoud", size=13, color=GRIJS, align=PP_ALIGN.CENTER)

ae = add_shape(slide, Inches(8.6), Inches(5.0), Inches(2.8), Inches(1.6), ORANJE_LICHT, ORANJE, Pt(2))
tf_ae = set_text(ae, "_Aanvang / _Einde", size=18, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
ae.text_frame.margin_top = Inches(0.1)
add_para(tf_ae, "(ent_id, rel_id, versie)", size=14, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
add_para(tf_ae, "Materiële tijdlijn", size=13, color=GRIJS, align=PP_ALIGN.CENTER)

add_arrow(slide, Inches(4.8), Inches(5.8), Inches(5.3), Inches(5.8), GROEN, Pt(2))
add_arrow(slide, Inches(4.8), Inches(5.8), Inches(8.6), Inches(5.8), ORANJE, Pt(2))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11: HUB + DATA HIËRARCHIE
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Hub + Data: hiërarchie", "Entiteit → Hub → Data / Aanvang / Einde")

ent_a = add_shape(slide, Inches(0.4), Inches(1.7), Inches(2.3), Inches(0.8), BLAUW_LICHT, MIDDENBLAUW, Pt(2))
set_text(ent_a, "NatuurlijkPersoon", size=15, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)
ent_a.text_frame.margin_top = Inches(0.1)

rows = [
    ("PersoonsIdentificatie", False, Inches(1.5)),
    ("Naam (hub)", False, Inches(2.5)),
    ("🕐 Burgerschap (hub)", True, Inches(3.5)),
    ("🕐 Bereikbaarheid (rel)", True, Inches(5.1)),
]

for name, materieel, y in rows:
    clr = ORANJE if "rel" in name else MIDDENBLAUW
    bg = ORANJE_LICHT if "rel" in name else RGBColor(0xE0, 0xE7, 0xFF)
    h = add_shape(slide, Inches(3.3), y, Inches(2.6), Inches(0.6), bg, clr)
    set_text(h, name, size=13, bold=True, color=clr, align=PP_ALIGN.CENTER)
    h.text_frame.margin_top = Inches(0.05)

    d = add_shape(slide, Inches(6.3), y - Inches(0.2), Inches(2.5), Inches(0.55), GROEN_LICHT, GROEN)
    data_name = name.replace("🕐 ", "").replace(" (hub)", "").replace(" (rel)", "") + "_Data"
    set_text(d, data_name, size=12, bold=True, color=GROEN, align=PP_ALIGN.CENTER)
    d.text_frame.margin_top = Inches(0.05)

    add_arrow(slide, Inches(5.9), y + Inches(0.15), Inches(6.3), y, GROEN)

    if materieel:
        a = add_shape(slide, Inches(6.3), y + Inches(0.5), Inches(2.5), Inches(0.45), ORANJE_LICHT, ORANJE)
        a_name = name.replace("🕐 ", "").replace(" (hub)", "").replace(" (rel)", "") + "_Aanvang"
        set_text(a, a_name, size=11, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
        a.text_frame.margin_top = Inches(0.03)

        e = add_shape(slide, Inches(9.2), y + Inches(0.5), Inches(2.5), Inches(0.45), ROOD_LICHT, ROOD)
        e_name = name.replace("🕐 ", "").replace(" (hub)", "").replace(" (rel)", "") + "_Einde"
        set_text(e, e_name, size=11, bold=True, color=ROOD, align=PP_ALIGN.CENTER)
        e.text_frame.margin_top = Inches(0.03)

        add_arrow(slide, Inches(5.9), y + Inches(0.3), Inches(6.3), y + Inches(0.7), ORANJE)
        add_arrow(slide, Inches(8.8), y + Inches(0.7), Inches(9.2), y + Inches(0.7), ROOD)

# Ent → hubs
for y in [Inches(1.8), Inches(2.8), Inches(3.8), Inches(5.4)]:
    add_arrow(slide, Inches(2.7), Inches(2.1), Inches(3.3), y, MIDDENBLAUW)

# Plumbing
ep_a = add_shape(slide, Inches(0.2), Inches(5.1), Inches(2.5), Inches(0.5), ORANJE_LICHT, ORANJE)
set_text(ep_a, "Pers_Aanvang (plumbing)", size=11, bold=True, color=ORANJE, align=PP_ALIGN.CENTER)
ep_a.text_frame.margin_top = Inches(0.05)
ep_e = add_shape(slide, Inches(0.2), Inches(5.8), Inches(2.5), Inches(0.5), ROOD_LICHT, ROOD)
set_text(ep_e, "Pers_Einde (plumbing)", size=11, bold=True, color=ROOD, align=PP_ALIGN.CENTER)
ep_e.text_frame.margin_top = Inches(0.05)
add_arrow(slide, Inches(1.4), Inches(2.5), Inches(1.4), Inches(5.1), ORANJE)

# Legenda
leg = add_shape(slide, Inches(9.5), Inches(1.5), Inches(3.3), Inches(2.2), WIT, GRIJS)
tf_leg = set_text(leg, "Legenda", size=13, bold=True, color=DONKERGRIJS)
leg.text_frame.margin_left = Inches(0.15)
leg.text_frame.margin_top = Inches(0.05)
add_para(tf_leg, "🟦 Hub (identiteitsanker)", size=12, color=MIDDENBLAUW)
add_para(tf_leg, "🟩 _Data (inhoud, versioned)", size=12, color=GROEN)
add_para(tf_leg, "🟧 _Aanvang", size=12, color=ORANJE)
add_para(tf_leg, "🟥 _Einde", size=12, color=ROOD)
add_para(tf_leg, "🕐 Materieel type", size=12, color=DONKERGRIJS)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12: VOORBEELD NAAMSWIJZIGING MET CORRECTIE
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Voorbeeld: naamswijziging met correctie", "Hoe hub + data + aanvang/einde samenwerken")

s1 = add_shape(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(2.4), WIT, MIDDENBLAUW, Pt(2))
tf_s1 = set_text(s1, "Stap 1 — Registratie (t_reg = 21/3/2026)", size=16, bold=True, color=MIDDENBLAUW)
s1.text_frame.margin_left = Inches(0.25)
s1.text_frame.margin_top = Inches(0.1)
add_para(tf_s1, "", size=3)
add_para(tf_s1, "Hub 1: Naam rel_id=1", size=13, bold=True, color=ZWART)
add_para(tf_s1, "  _Data v1: \"Jansen\"  |  Aanvang: geboortedatum  |  Einde: 20/3", size=12, color=DONKERGRIJS)
add_para(tf_s1, "Hub 2: Naam rel_id=2", size=13, bold=True, color=ZWART)
add_para(tf_s1, "  _Data v1: \"De Vries\"  |  Aanvang: 21/3/2026  |  Einde: —", size=12, color=DONKERGRIJS)
add_para(tf_s1, "⚠  Aanvangsdatum foutief: was 21/3, moet 1/1/2025 zijn", size=13, color=ROOD)

s2 = add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.4), WIT, ORANJE, Pt(2))
tf_s2 = set_text(s2, "Stap 2 — Correctie (t_corr > t_reg)", size=16, bold=True, color=ORANJE)
s2.text_frame.margin_left = Inches(0.25)
s2.text_frame.margin_top = Inches(0.1)
add_para(tf_s2, "", size=3)
add_para(tf_s2, "Rechtbank: naamswijziging per 1/1/2025", size=13, bold=True, color=ZWART)
add_para(tf_s2, "Hub 1: Einde v1 → afg, Einde v2 = 31/12/2024", size=12, color=DONKERGRIJS)
add_para(tf_s2, "Hub 2: Aanvang v1 → afg, Aanvang v2 = 1/1/2025", size=12, color=DONKERGRIJS)
add_para(tf_s2, "✓ Hub-records ongewijzigd  ✓ _Data ongewijzigd", size=13, color=GROEN)
add_para(tf_s2, "✓ Alleen _Aanvang/_Einde gecorrigeerd!", size=13, bold=True, color=GROEN)

box_i = add_shape(slide, Inches(0.4), Inches(4.2), Inches(12.6), Inches(0.8), GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Pt(2))
set_text(box_i, "💡  Elke laag apart corrigeerbaar — identiteit, inhoud en materiële tijdlijn beïnvloeden elkaar niet.", 
         size=17, bold=True, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)
box_i.text_frame.margin_top = Inches(0.1)

# Mini-tabel
headers = ["Laag", "rel_id", "versie", "waarde", "opvoer", "afvoer"]
data_rows = [
    ["Hub 1", "1", "—", "—", "t_reg", "—"],
    ["Hub 1 _data", "1", "1", '"Jansen"', "t_reg", "—"],
    ["Hub 1 _einde", "1", "1", "20/3/2026", "t_reg", "t_corr"],
    ["Hub 1 _einde", "1", "2", "31/12/2024", "t_corr", "—"],
    ["Hub 2 _aanvang", "2", "1", "21/3/2026", "t_reg", "t_corr"],
    ["Hub 2 _aanvang", "2", "2", "1/1/2025", "t_corr", "—"],
]
col_widths = [Inches(2.0), Inches(1.0), Inches(1.0), Inches(2.0), Inches(1.5), Inches(1.5)]
tbl_x = Inches(1.5)
tbl_y = Inches(5.2)

x = tbl_x
for i, h in enumerate(headers):
    hdr = add_rect(slide, x, tbl_y, col_widths[i], Inches(0.3), MIDDENBLAUW)
    set_text(hdr, h, size=11, bold=True, color=WIT, align=PP_ALIGN.CENTER)
    hdr.text_frame.margin_top = Inches(0.02)
    x += col_widths[i]

for r_idx, row in enumerate(data_rows):
    x = tbl_x
    row_y = tbl_y + Inches(0.3) + Inches(0.26) * r_idx
    bg = WIT if r_idx % 2 == 0 else ACHTERGROND
    for c_idx, cell in enumerate(row):
        c = add_rect(slide, x, row_y, col_widths[c_idx], Inches(0.26), bg, GRIJS, Pt(0.5))
        clr = ROOD if cell == "t_corr" else ZWART
        set_text(c, cell, size=10, color=clr, align=PP_ALIGN.CENTER)
        c.text_frame.margin_top = Inches(0.01)
        x += col_widths[c_idx]


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13: DATABASE-ONTWERP ENTITEIT
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Database-ontwerp: entiteiten", "PostgreSQL tabellen voor entiteiten en registratie")

ent_box = add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(2.2), WIT, MIDDENBLAUW, Pt(2))
tf_e = set_text(ent_box, "natuurlijk_persoon", size=20, bold=True, color=MIDDENBLAUW)
ent_box.text_frame.margin_left = Inches(0.3)
ent_box.text_frame.margin_top = Inches(0.1)
add_para(tf_e, "", size=4)
add_para(tf_e, "id          INTEGER  PK", size=14, color=ZWART, font_name="Consolas")
add_para(tf_e, "opvoer      TIMESTAMPTZ     — afgeleid", size=14, color=GRIJS, font_name="Consolas")
add_para(tf_e, "afvoer      TIMESTAMPTZ     — afgeleid", size=14, color=GRIJS, font_name="Consolas")
add_para(tf_e, "", size=4)
add_para(tf_e, "Geen inhoudsvelden — die zitten in GE's", size=13, color=DONKERGRIJS)

reg_box = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.2), WIT, PAARS, Pt(2))
tf_reg = set_text(reg_box, "registratie", size=20, bold=True, color=PAARS)
reg_box.text_frame.margin_left = Inches(0.3)
reg_box.text_frame.margin_top = Inches(0.1)
add_para(tf_reg, "", size=4)
add_para(tf_reg, "id                          SERIAL  PK", size=13, color=ZWART, font_name="Consolas")
add_para(tf_reg, "tijdstip                    TIMESTAMPTZ", size=13, color=ZWART, font_name="Consolas")
add_para(tf_reg, "registratie_type            TEXT", size=13, color=ZWART, font_name="Consolas")
add_para(tf_reg, "corrigeert_registratie_id   INT FK", size=13, color=GRIJS, font_name="Consolas")

wij_box = add_shape(slide, Inches(6.8), Inches(4.1), Inches(5.8), Inches(2.4), WIT, PAARS, Pt(2))
tf_wij = set_text(wij_box, "wijziging", size=20, bold=True, color=PAARS)
wij_box.text_frame.margin_left = Inches(0.3)
wij_box.text_frame.margin_top = Inches(0.1)
add_para(tf_wij, "", size=4)
add_para(tf_wij, "id                  SERIAL  PK", size=13, color=ZWART, font_name="Consolas")
add_para(tf_wij, "registratie_id      INT  FK → registratie", size=13, color=ZWART, font_name="Consolas")
add_para(tf_wij, "wijziging_type      TEXT  (Opvoer/Afvoer)", size=13, color=ZWART, font_name="Consolas")
add_para(tf_wij, "representatienaam   TEXT  (tabelnaam)", size=13, color=ZWART, font_name="Consolas")
add_para(tf_wij, "representatie_id    TEXT  (samengesteld PK)", size=13, color=ZWART, font_name="Consolas")

plumb = add_shape(slide, Inches(0.6), Inches(4.1), Inches(5.5), Inches(2.4), WIT, ORANJE, Pt(2))
tf_pl = set_text(plumb, "Materieel: pers_aanvang / pers_einde", size=16, bold=True, color=ORANJE)
plumb.text_frame.margin_left = Inches(0.3)
plumb.text_frame.margin_top = Inches(0.1)
add_para(tf_pl, "", size=4)
add_para(tf_pl, "id       INT  FK → persoon(id)", size=13, color=ZWART, font_name="Consolas")
add_para(tf_pl, "versie   SERIAL  (autoincrement per id)", size=13, color=ZWART, font_name="Consolas")
add_para(tf_pl, "datum    DATE", size=13, color=ZWART, font_name="Consolas")
add_para(tf_pl, "opvoer   TIMESTAMPTZ  — afgeleid", size=13, color=GRIJS, font_name="Consolas")
add_para(tf_pl, "afvoer   TIMESTAMPTZ  — afgeleid", size=13, color=GRIJS, font_name="Consolas")
add_para(tf_pl, "PK: (id, versie)", size=13, bold=True, color=ORANJE)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14: DATABASE HUB + DATA TABELLEN
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Database: Hub + Data + Aanvang/Einde", "Hoe een gegevenselement eruitziet in PostgreSQL")

hub_box = add_shape(slide, Inches(0.3), Inches(1.5), Inches(4.2), Inches(2.2), WIT, MIDDENBLAUW, Pt(2))
tf_hub = set_text(hub_box, "naam  (hub)", size=17, bold=True, color=MIDDENBLAUW)
hub_box.text_frame.margin_left = Inches(0.2)
hub_box.text_frame.margin_top = Inches(0.1)
add_para(tf_hub, "", size=3)
add_para(tf_hub, "persoon_id  INT  FK → persoon(id)", size=12, color=ZWART, font_name="Consolas")
add_para(tf_hub, "rel_id      INT  autoincrement/scope", size=12, color=ZWART, font_name="Consolas")
add_para(tf_hub, "opvoer/afvoer  — afgeleid", size=12, color=GRIJS, font_name="Consolas")
add_para(tf_hub, "PK: (persoon_id, rel_id)", size=12, bold=True, color=MIDDENBLAUW)
add_para(tf_hub, "Geen inhoudsvelden!", size=12, bold=True, color=ROOD)

data_box = add_shape(slide, Inches(4.7), Inches(1.5), Inches(4.2), Inches(2.2), WIT, GROEN, Pt(2))
tf_data = set_text(data_box, "naam_data", size=17, bold=True, color=GROEN)
data_box.text_frame.margin_left = Inches(0.2)
data_box.text_frame.margin_top = Inches(0.1)
add_para(tf_data, "", size=3)
add_para(tf_data, "persoon_id + rel_id  FK → naam", size=12, color=ZWART, font_name="Consolas")
add_para(tf_data, "versie      INT  autoincrement", size=12, color=ZWART, font_name="Consolas")
add_para(tf_data, "voorletters TEXT, achternaam TEXT", size=12, color=ZWART, font_name="Consolas")
add_para(tf_data, "opvoer/afvoer  — afgeleid", size=12, color=GRIJS, font_name="Consolas")
add_para(tf_data, "PK: (persoon_id, rel_id, versie)", size=12, bold=True, color=GROEN)

ae_box = add_shape(slide, Inches(9.1), Inches(1.5), Inches(4.0), Inches(2.2), WIT, ORANJE, Pt(2))
tf_ae = set_text(ae_box, "naam_aanvang / _einde", size=16, bold=True, color=ORANJE)
ae_box.text_frame.margin_left = Inches(0.2)
ae_box.text_frame.margin_top = Inches(0.1)
add_para(tf_ae, "Als IsMaterieel = true", size=11, color=GRIJS)
add_para(tf_ae, "", size=3)
add_para(tf_ae, "persoon_id + rel_id  FK → naam", size=12, color=ZWART, font_name="Consolas")
add_para(tf_ae, "versie  INT  autoincrement", size=12, color=ZWART, font_name="Consolas")
add_para(tf_ae, "datum   DATE", size=12, color=ZWART, font_name="Consolas")
add_para(tf_ae, "PK: (persoon_id, rel_id, versie)", size=12, bold=True, color=ORANJE)

# Relatieve autoincrement
rai = add_shape(slide, Inches(0.3), Inches(4.0), Inches(12.7), Inches(2.8), GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Pt(2))
tf_rai = set_text(rai, "Relatieve autoincrement", size=20, bold=True, color=RGBColor(0x85, 0x4D, 0x0E))
rai.text_frame.margin_left = Inches(0.3)
rai.text_frame.margin_top = Inches(0.1)
add_para(tf_rai, "", size=4)
add_para(tf_rai, "Werkt binnen de scope van de parent — geen globale sequences:", size=16, color=ZWART)
add_para(tf_rai, "", size=4)
add_para(tf_rai, "▸  Hub rel_id:  binnen (persoon_id) en type → Persoon #1 heeft Naam rel_id=1, rel_id=2, ...", size=14, color=DONKERGRIJS)
add_para(tf_rai, "▸  _Data versie:  binnen (persoon_id, rel_id) → Naam 1 v1, v2, v3, ...", size=14, color=DONKERGRIJS)
add_para(tf_rai, "▸  _Aanvang versie:  idem, maar vorige versie wordt automatisch afgevoerd (enkelvoudig)", size=14, color=DONKERGRIJS)
add_para(tf_rai, "", size=6)
add_para(tf_rai, "Dit maakt volledige versiegeschiedenis mogelijk zonder globale sequences.", size=14, bold=True, color=RGBColor(0x85, 0x4D, 0x0E))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15: METAREGISTRY
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "MetaRegistry: single source of truth", "Alle metadata op één plek — dynamisch gedreven")

main = add_shape(slide, Inches(0.6), Inches(1.6), Inches(12.2), Inches(1.0), PAARS_LICHT, PAARS, Pt(2))
tf_m = set_text(main, "De MetaRegistry is een Go map van TypeMeta entries — één per representatietype.",
                size=17, bold=True, color=PAARS, align=PP_ALIGN.CENTER)
main.text_frame.margin_top = Inches(0.05)
add_para(tf_m, "Routes, handlers, schema-API en frontend worden allemaal dynamisch gedreven door de MetaRegistry.", 
         size=14, color=DONKERGRIJS, align=PP_ALIGN.CENTER)

tm = add_shape(slide, Inches(0.6), Inches(2.9), Inches(5.8), Inches(4.0), WIT, MIDDENBLAUW, Pt(2))
tf_tm = set_text(tm, "TypeMeta velden", size=18, bold=True, color=MIDDENBLAUW)
tm.text_frame.margin_left = Inches(0.25)
tm.text_frame.margin_top = Inches(0.1)
add_para(tf_tm, "", size=3)
add_para(tf_tm, "▸ Typenaam, Description, Metatype", size=13, color=ZWART)
add_para(tf_tm, "▸ IsMaterieel, GESubtype (hub/data/aanvang/einde)", size=13, color=ZWART)
add_para(tf_tm, "▸ Padnaam (URL), Veldnaam (JSON)", size=13, color=ZWART)
add_para(tf_tm, "▸ Factory / SliceFactory (constructors)", size=13, color=ZWART)
add_para(tf_tm, "▸ Tabelnaam, IDKolom, EntiteitIDKolom", size=13, color=ZWART)
add_para(tf_tm, "▸ HeeftPFK, RelatieveAutoincrement", size=13, color=ZWART)
add_para(tf_tm, "▸ OnderliggendeGegevenselementen []", size=13, color=ZWART)
add_para(tf_tm, "▸ Momentvoorkomen (enkelvoudig/meervoudig)", size=13, color=ZWART)
add_para(tf_tm, "▸ BovenliggendTypenaam", size=13, color=ZWART)
add_para(tf_tm, "▸ Kleur (visualisatie)", size=13, color=ZWART)

drives = add_shape(slide, Inches(6.8), Inches(2.9), Inches(5.8), Inches(4.0), WIT, GROEN, Pt(2))
tf_d = set_text(drives, "Wat de MetaRegistry aanstuurt", size=18, bold=True, color=GROEN)
drives.text_frame.margin_left = Inches(0.25)
drives.text_frame.margin_top = Inches(0.1)
add_para(tf_d, "", size=3)
add_para(tf_d, "🔀  Routes — dynamisch geregistreerd via Padnaam", size=14, color=ZWART)
add_para(tf_d, "", size=3)
add_para(tf_d, "⚙️  Handlers — generiek per MetaRegistry-entry", size=14, color=ZWART)
add_para(tf_d, "", size=3)
add_para(tf_d, "📋  Schema-API — /api/viz/schema retourneert", size=14, color=ZWART)
add_para(tf_d, "      metadata + veldtypes + relaties", size=13, color=DONKERGRIJS)
add_para(tf_d, "", size=3)
add_para(tf_d, "🖥️  Frontend — leest dynamisch het schema", size=14, color=ZWART)
add_para(tf_d, "      Geen hardcoded veldnamen!", size=13, bold=True, color=GROEN)
add_para(tf_d, "", size=3)
add_para(tf_d, "📊  Visualisatie — kleuren, layout, oortjes", size=14, color=ZWART)
add_para(tf_d, "", size=3)
add_para(tf_d, "🔄  Code generatie — MetaRegistry wordt gegenereerd", size=14, color=ZWART)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16: SCHEMA-API & FRONTEND
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Schema-API & dynamische frontend", "Geen hardcoded structuren — alles uit de MetaRegistry")

tb = add_textbox(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5))
set_text(tb, "Flow:  MetaRegistry  →  Schema-API ( /api/viz/schema )  →  Frontend (React/Vite)", 
         size=18, bold=True, color=MIDDENBLAUW, align=PP_ALIGN.CENTER)

api_box = add_shape(slide, Inches(0.4), Inches(2.3), Inches(6.2), Inches(4.8), WIT, MIDDENBLAUW, Pt(2))
tf_api = set_text(api_box, "Schema-API response (voorbeeld)", size=15, bold=True, color=MIDDENBLAUW)
api_box.text_frame.margin_left = Inches(0.2)
api_box.text_frame.margin_top = Inches(0.1)
add_para(tf_api, "", size=3)
for line in [
    '{', '  "typenaam": "Naam",', '  "metatype": "gegevenselement",',
    '  "ge_subtype": "hub",', '  "is_materieel": false,',
    '  "padnaam": "naam",', '  "velden": [',
    '    {"naam":"achternaam","type":"string"},',
    '    {"naam":"naamgebruik","type":"string",',
    '     "enum":["EigenNaam","PartnerNaam",...]}',
    '  ],', '  "onderliggende": [',
    '    {"rolnaam":"data","doeltype":"Naam_Data"}',
    '  ]', '}'
]:
    bold = "ge_subtype" in line
    clr = LICHTBLAUW if bold else RGBColor(0xA5, 0xD6, 0xFF) if line.startswith('    ') else ZWART
    add_para(tf_api, line, size=11, color=clr, font_name="Consolas", bold=bold, space_before=Pt(1), space_after=Pt(0))

fe_box = add_shape(slide, Inches(6.8), Inches(2.3), Inches(6.2), Inches(4.8), WIT, GROEN, Pt(2))
tf_fe = set_text(fe_box, "Frontend: dynamisch gegenereerd", size=15, bold=True, color=GROEN)
fe_box.text_frame.margin_left = Inches(0.25)
fe_box.text_frame.margin_top = Inches(0.1)
add_para(tf_fe, "", size=4)
add_para(tf_fe, "De frontend leest het schema en genereert:", size=14, color=ZWART)
add_para(tf_fe, "", size=3)
add_para(tf_fe, "▸  Formulieren — veldtype bepaalt inputcomponent", size=14, color=DONKERGRIJS)
add_para(tf_fe, "   string → tekstveld, date → datepicker,", size=12, color=GRIJS)
add_para(tf_fe, "   enum → dropdown, bool → checkbox", size=12, color=GRIJS)
add_para(tf_fe, "", size=3)
add_para(tf_fe, "▸  Overzichten — kolommen uit schema-velden", size=14, color=DONKERGRIJS)
add_para(tf_fe, "", size=3)
add_para(tf_fe, "▸  Visualisatie (SVG) — entiteitskaarten,", size=14, color=DONKERGRIJS)
add_para(tf_fe, "   GE-boxes, relatielijnen, aanvang/einde-oortjes", size=12, color=GRIJS)
add_para(tf_fe, "", size=3)
add_para(tf_fe, "▸  Registratiepagina — inzien van audit-trail", size=14, color=DONKERGRIJS)
add_para(tf_fe, "", size=6)
add_para(tf_fe, "Nieuw type toevoegen = MetaRegistry aanpassen", size=14, bold=True, color=GROEN)
add_para(tf_fe, "→ frontend past zich automatisch aan", size=14, bold=True, color=GROEN)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17: ROUNDTRIP UML ↔ CODE (OVERZICHT)
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Roundtrip: UML ↔ Code ↔ Register", "Van visueel model tot werkend register en terug")

# Flow diagram met grote pijlen
steps = [
    ("UML Editor", "React + ReactFlow\nvisueel modelleren", BLAUW_LICHT, MIDDENBLAUW, Inches(0.4)),
    ("V3 JSON", "Uitwisselformaat\nmetamodel_v3.json", GEEL_LICHT, AMBER, Inches(2.6)),
    ("Codegen", "cmd/codegen/\n7 Go-bestanden", PAARS_LICHT, PAARS, Inches(4.8)),
    ("Go API", "MetaRegistry\n+ handlers + routes", GROEN_LICHT, GROEN, Inches(7.0)),
    ("PostgreSQL", "Tabellen + FK's\nrelatieve autoincrement", RGBColor(0xE0, 0xE7, 0xFF), MIDDENBLAUW, Inches(9.2)),
    ("Register", "Werkend bitemporeel\nregister + frontend", ORANJE_LICHT, ORANJE, Inches(11.4)),
]

for name, desc, bg, line, x in steps:
    box = add_shape(slide, x, Inches(2.0), Inches(1.7), Inches(2.2), bg, line, Pt(2))
    tf = set_text(box, name, size=16, bold=True, color=line, align=PP_ALIGN.CENTER)
    box.text_frame.margin_top = Inches(0.1)
    box.text_frame.margin_left = Inches(0.05)
    add_para(tf, "", size=4)
    for dl in desc.split('\n'):
        add_para(tf, dl, size=11, color=DONKERGRIJS, align=PP_ALIGN.CENTER, space_before=Pt(1), space_after=Pt(0))

# Forward arrows
for i in range(len(steps) - 1):
    x1 = steps[i][4] + Inches(1.7)
    x2 = steps[i+1][4]
    add_arrow(slide, x1, Inches(3.1), x2, Inches(3.1), GRIJS, Pt(2))

# Reverse arrow (grote terugpijl)
from pptx.oxml.ns import qn
rev = slide.shapes.add_connector(1, Inches(11.4), Inches(4.5), Inches(2.1), Inches(4.5))
rev.line.color.rgb = LICHTBLAUW
rev.line.width = Pt(3)
rev.line.dash_style = 4  # dash-dot
ln = rev.line._ln
head = ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
ln.append(head)

tb_rev = add_textbox(slide, Inches(3.5), Inches(4.6), Inches(7.0), Inches(0.4))
set_text(tb_rev, "⬅  Reverse: GET /api/schema/model → V3 JSON → Editor (bidirectionele roundtrip)", 
         size=14, bold=True, color=LICHTBLAUW, align=PP_ALIGN.CENTER)

# Wat codegen oplevert
gen_box = add_shape(slide, Inches(0.4), Inches(5.2), Inches(12.7), Inches(1.8), WIT, PAARS, Pt(2))
tf_gen = set_text(gen_box, "Code generator levert 7 bestanden:", size=16, bold=True, color=PAARS)
gen_box.text_frame.margin_left = Inches(0.3)
gen_box.text_frame.margin_top = Inches(0.1)
add_para(tf_gen, "", size=3)
add_para(tf_gen, "modellen_entiteiten.go    Entiteit-structs + plumbing         modellen_input.go     Platte API-input-structs", size=12, color=DONKERGRIJS, font_name="Consolas")
add_para(tf_gen, "modellen_ge_rel.go        Hub, _Data, _Aanvang, _Einde        metaregistry.go       Volledige MetaRegistry", size=12, color=DONKERGRIJS, font_name="Consolas")
add_para(tf_gen, "modellen_methods.go       Interface-methoden (GetID, etc.)     datatype_registry.go  Custom datatypes (BSN, etc.)", size=12, color=DONKERGRIJS, font_name="Consolas")
add_para(tf_gen, "enum_registry.go          Enum-waarden + editor-layouts", size=12, color=DONKERGRIJS, font_name="Consolas")


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18: ROUNDTRIP DETAIL - V3 JSON
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "V3 JSON: het uitwisselformaat", "Single exchange format tussen editor, codegen en API")

# Structuur links
struct_box = add_shape(slide, Inches(0.4), Inches(1.5), Inches(5.5), Inches(5.5), WIT, AMBER, Pt(2))
tf_s = set_text(struct_box, "V3 JSON structuur", size=18, bold=True, color=AMBER)
struct_box.text_frame.margin_left = Inches(0.25)
struct_box.text_frame.margin_top = Inches(0.1)
add_para(tf_s, "", size=4)
add_para(tf_s, "versie, naam, beschrijving", size=14, bold=True, color=ZWART)
add_para(tf_s, "", size=4)
add_para(tf_s, "datatypes []", size=14, bold=True, color=TEAL)
add_para(tf_s, "  naam, basistype, format, validatie, normalisatie", size=12, color=DONKERGRIJS)
add_para(tf_s, "", size=4)
add_para(tf_s, "enums []", size=14, bold=True, color=PAARS)
add_para(tf_s, "  goType, baseType, waarden [{constNaam, waarde}]", size=12, color=DONKERGRIJS)
add_para(tf_s, "", size=4)
add_para(tf_s, "referentielijstInstanties []", size=14, bold=True, color=ORANJE)
add_para(tf_s, "  systeemnaam, naam, omschrijving", size=12, color=DONKERGRIJS)
add_para(tf_s, "", size=4)
add_para(tf_s, "entiteiten []", size=14, bold=True, color=MIDDENBLAUW)
add_para(tf_s, "  typenaam, isMaterieel, kleur, meervoud", size=12, color=DONKERGRIJS)
add_para(tf_s, "  afgeleideVelden []", size=12, color=DONKERGRIJS)
add_para(tf_s, "  gegevenselementen []", size=12, color=GROEN)
add_para(tf_s, "    naam, velden [], isMaterieel, momentvoorkomen", size=12, color=DONKERGRIJS)
add_para(tf_s, "  relaties []", size=12, color=ORANJE)
add_para(tf_s, "    naam, doelEntiteit, velden [], isMaterieel", size=12, color=DONKERGRIJS)
add_para(tf_s, "", size=4)
add_para(tf_s, "Editor-posities (x, y) worden meebewaard", size=13, color=GRIJS)
add_para(tf_s, "maar genegeerd door codegen", size=13, color=GRIJS)

# Bronnen / bestemmingen
dest_box = add_shape(slide, Inches(6.2), Inches(1.5), Inches(6.7), Inches(5.5), WIT, MIDDENBLAUW, Pt(2))
tf_d = set_text(dest_box, "Wie leest/schrijft V3 JSON?", size=18, bold=True, color=MIDDENBLAUW)
dest_box.text_frame.margin_left = Inches(0.25)
dest_box.text_frame.margin_top = Inches(0.1)
add_para(tf_d, "", size=6)
add_para(tf_d, "📝  UML Editor → V3 JSON", size=16, bold=True, color=ZWART)
add_para(tf_d, "     Exporteert het visuele model naar JSON", size=13, color=DONKERGRIJS)
add_para(tf_d, "     POST /api/schema/model — upload naar backend", size=13, color=DONKERGRIJS)
add_para(tf_d, "", size=6)
add_para(tf_d, "⚙️  Codegen ← V3 JSON", size=16, bold=True, color=ZWART)
add_para(tf_d, "     --input file.json  of  --from-url http://...", size=13, color=DONKERGRIJS)
add_para(tf_d, "     Genereert 7 Go-bestanden", size=13, color=DONKERGRIJS)
add_para(tf_d, "", size=6)
add_para(tf_d, "🔄  API → V3 JSON (reverse)", size=16, bold=True, color=ZWART)
add_para(tf_d, "     GET /api/schema/model — export van runtime", size=13, color=DONKERGRIJS)
add_para(tf_d, "     MetaRegistry → V3 JSON via v3_exporter.go", size=13, color=DONKERGRIJS)
add_para(tf_d, "", size=6)
add_para(tf_d, "📥  Editor ← V3 JSON (load from API)", size=16, bold=True, color=ZWART)
add_para(tf_d, "     Laadt model bij opstart vanuit de API", size=13, color=DONKERGRIJS)
add_para(tf_d, "     Volledige bidirectionele roundtrip!", size=14, bold=True, color=GROEN)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 19: IMPORT/EXPORT STANDAARDEN
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Import / Export standaarden", "Interoperabiliteit met de buitenwereld")

# Export
exp_box = add_shape(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.5), WIT, GROEN, Pt(2))
tf_exp = set_text(exp_box, "Export vanuit UML Editor", size=20, bold=True, color=GROEN)
exp_box.text_frame.margin_left = Inches(0.3)
exp_box.text_frame.margin_top = Inches(0.1)
add_para(tf_exp, "", size=6)

formats_exp = [
    ("Mermaid", "classDiagram syntax — direct bruikbaar in\nMarkdown, GitHub, documentatie", "exportMermaid.js"),
    ("PlantUML", "UML class diagram syntax — PNG/SVG via\nPlantUML server of CLI", "exportPlantUML.js"),
    ("XMI 1.1", "UML 1.4 XMI — Sparx Enterprise Architect\ncompatibel inclusief diagramposities", "exportXMI.js"),
    ("V3 JSON", "Eigen metamodel-formaat — voor codegen,\nAPI import/export, roundtrip", "/api/schema/model"),
]
for name, desc, file in formats_exp:
    add_para(tf_exp, f"▸  {name}", size=16, bold=True, color=ZWART)
    for line in desc.split('\n'):
        add_para(tf_exp, f"   {line}", size=13, color=DONKERGRIJS)
    add_para(tf_exp, f"   → {file}", size=11, color=GRIJS, font_name="Consolas")
    add_para(tf_exp, "", size=3)

# Import
imp_box = add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.5), WIT, MIDDENBLAUW, Pt(2))
tf_imp = set_text(imp_box, "Import naar UML Editor", size=20, bold=True, color=MIDDENBLAUW)
imp_box.text_frame.margin_left = Inches(0.3)
imp_box.text_frame.margin_top = Inches(0.1)
add_para(tf_imp, "", size=6)

formats_imp = [
    ("XMI", "Sparx EA XMI met EA-extensies,\ndiagramposities, stereotypes", "importXMI.js"),
    ("Mermaid", "classDiagram met stereotypes,\nvelden en cardinaliteiten", "importMermaid.js"),
    ("PlantUML", "Class diagrams met stereotypes", "importPlantUML.js"),
    ("V3 JSON", "Van API of bestand —\nbidirectionele roundtrip", "GET /api/schema/model"),
]
for name, desc, file in formats_imp:
    add_para(tf_imp, f"▸  {name}", size=16, bold=True, color=ZWART)
    for line in desc.split('\n'):
        add_para(tf_imp, f"   {line}", size=13, color=DONKERGRIJS)
    add_para(tf_imp, f"   → {file}", size=11, color=GRIJS, font_name="Consolas")
    add_para(tf_imp, "", size=3)

# OAS note
add_para(tf_imp, "", size=4)
add_para(tf_imp, "Schema-API gebruikt OAS 3.1 type-systeem", size=14, bold=True, color=MIDDENBLAUW)
add_para(tf_imp, "(type + format) voor veldtypering", size=13, color=DONKERGRIJS)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 20: MERMAID IN ACTIE
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Mermaid & PlantUML: echte output", "Gegenereerde UML-code vanuit de editor")

# Mermaid voorbeeld
mermaid_full = [
    "classDiagram",
    "  direction TB",
    "",
    '  class NatuurlijkPersoon {',
    '    <<entiteit, materieel>>',
    '    +id : integer',
    '    /weergavenaam : string',
    '  }',
    '  class PersoonsIdentificatie {',
    '    <<gegevenselement>>',
    '    +bsn : string {BSN}',
    '    -ingezetene : boolean',
    '  }',
    '  class Naam {',
    '    <<gegevenselement>>',
    '    +voorletters : string',
    '    +achternaam : string',
    '    +naamgebruik : Naamgebruiksoort',
    '  }',
    '  NatuurlijkPersoon "1" --> "0..1" PersoonsId : Us',
    '  NatuurlijkPersoon "1" --> "0..*" Naam : Vs',
]
add_code_box(slide, Inches(0.3), Inches(1.4), Inches(6.3), Inches(5.6),
             mermaid_full, title="Mermaid classDiagram (geëxporteerd)", title_color=LICHTBLAUW)

# PlantUML voorbeeld
plantuml_lines = [
    "@startuml",
    "skinparam classAttributeIconSize 0",
    "hide empty methods",
    "",
    "class NatuurlijkPersoon <<entiteit, materieel>> {",
    "  + id : integer",
    "  / weergavenaam : string",
    "}",
    "",
    "class PersoonsIdentificatie <<gegevenselement>> {",
    "  + bsn : string {BSN}",
    "  - ingezetene : boolean",
    "}",
    "",
    'class Bereikbaarheid <<relatie, materieel>> {',
    '  + soort : Bereikbaarheidsoort',
    '}',
    "",
    "enum Bereikbaarheidsoort {",
    "  Woonadres",
    "  Briefadres",
    "  Correspondentieadres",
    "}",
    "",
    'NatuurlijkPersoon "1" --> "0..1" PersoonsId : Us',
    '@enduml',
]
add_code_box(slide, Inches(6.8), Inches(1.4), Inches(6.3), Inches(5.6),
             plantuml_lines, title="PlantUML (geëxporteerd)", title_color=GROEN)

# Note
tb_note = add_textbox(slide, Inches(0.3), Inches(7.0), Inches(13), Inches(0.4))
set_text(tb_note, "De editor exporteert ook XMI 1.1 (Sparx EA-compatibel) met diagramposities — voor enterprise tooling interop", 
         size=13, color=GRIJS, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 21: VALIDATIES
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Validaties: datatypes met regels", "BSN, NLPostcode, IBAN — gedefinieerd in het metamodel")

# Architectuur
arch = add_shape(slide, Inches(0.4), Inches(1.5), Inches(4.5), Inches(2.5), WIT, TEAL, Pt(2))
tf_a = set_text(arch, "Validatie-architectuur", size=18, bold=True, color=TEAL)
arch.text_frame.margin_left = Inches(0.25)
arch.text_frame.margin_top = Inches(0.1)
add_para(tf_a, "", size=3)
add_para(tf_a, "Validatie zit op datatype-niveau, niet per veld:", size=14, color=ZWART)
add_para(tf_a, "", size=3)
add_para(tf_a, "▸  Metamodel (V3 JSON) definieert datatypes", size=13, color=DONKERGRIJS)
add_para(tf_a, "▸  Codegen → DatatypeRegistry (Go)", size=13, color=DONKERGRIJS)
add_para(tf_a, "▸  Frontend validation lib (framework-agnostisch JS)", size=13, color=DONKERGRIJS)
add_para(tf_a, "▸  Velden verwijzen naar datatype bij naam", size=13, color=DONKERGRIJS)
add_para(tf_a, "▸  3 regeltypen: checksum, formula, function", size=13, color=DONKERGRIJS)

# Pipeline
pipe = add_shape(slide, Inches(0.4), Inches(4.3), Inches(4.5), Inches(2.8), WIT, TEAL, Pt(2))
tf_p = set_text(pipe, "Validatie-pipeline", size=18, bold=True, color=TEAL)
pipe.text_frame.margin_left = Inches(0.25)
pipe.text_frame.margin_top = Inches(0.1)
add_para(tf_p, "", size=3)
add_para(tf_p, "1. Normalisatie — trim, uppercase, strip spaces", size=13, color=DONKERGRIJS)
add_para(tf_p, "2. Type-check — string/integer/number/boolean", size=13, color=DONKERGRIJS)
add_para(tf_p, "3. Lengte/bereik — min, max, minLength, maxLength", size=13, color=DONKERGRIJS)
add_para(tf_p, "4. Pattern — regex (bijv. postcode-formaat)", size=13, color=DONKERGRIJS)
add_para(tf_p, "5. Regels — checksum (BSN), formula, function", size=13, color=DONKERGRIJS)
add_para(tf_p, "", size=4)
add_para(tf_p, "Return: { geldig, genormaliseerd, fouten[] }", size=13, bold=True, color=TEAL, font_name="Consolas")

# BSN voorbeeld
bsn_lines = [
    '{',
    '  "naam": "BSN",',
    '  "basistype": "string",',
    '  "format": "bsn",',
    '  "validatie": {',
    '    "pattern": "^[0-9]{9}$",',
    '    "minLength": 9, "maxLength": 9,',
    '    "regels": [{',
    '      "naam": "11-proef",',
    '      "type": "checksum",',
    '      "expressie":',
    '        "(9*d1+8*d2+7*d3+6*d4',
    '         +5*d5+4*d6+3*d7+2*d8',
    '         -1*d9) % 11 == 0"',
    '    }]',
    '  }',
    '}',
]
add_code_box(slide, Inches(5.2), Inches(1.5), Inches(3.8), Inches(4.0),
             bsn_lines, title="BSN datatype (metamodel)", title_color=TEAL)

# NLPostcode voorbeeld
pc_lines = [
    '{',
    '  "naam": "NLPostcode",',
    '  "basistype": "string",',
    '  "format": "nl-postcode",',
    '  "validatie": {',
    '    "pattern":',
    '      "^[1-9][0-9]{3}\\\\s?[A-Za-z]{2}$",',
    '    "minLength": 6, "maxLength": 7,',
    '    "foutmelding":',
    '      "Voer geldige postcode in"',
    '  },',
    '  "normalisatie": "uppercase_letters",',
    '  "weergave": {',
    '    "placeholder": "1234 AB",',
    '    "inputMask": "0000 AA"',
    '  }',
    '}',
]
add_code_box(slide, Inches(9.2), Inches(1.5), Inches(3.8), Inches(4.0),
             pc_lines, title="NLPostcode (metamodel)", title_color=TEAL)

# IBAN note
iban = add_shape(slide, Inches(5.2), Inches(5.8), Inches(7.8), Inches(1.0), TEAL_LICHT, TEAL, Pt(1))
tf_ib = set_text(iban, "IBAN — MOD-97 check (ISO 13616): geïmplementeerd als named function in FUNCTION_REGISTRY", 
                 size=14, color=TEAL, align=PP_ALIGN.CENTER)
iban.text_frame.margin_top = Inches(0.05)
add_para(tf_ib, "Validatie-library is framework-agnostisch JS — herbruikbaar in editor én data-frontend", 
         size=12, color=DONKERGRIJS, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 22: AFGELEIDE VELDEN
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Afgeleide velden", "Berekende waarden uit het model — met CEL-expressies")

# Twee niveaus
n_box = add_shape(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(2.5), WIT, ORANJE, Pt(2))
tf_n = set_text(n_box, "Twee niveaus van afleiding", size=18, bold=True, color=ORANJE)
n_box.text_frame.margin_left = Inches(0.25)
n_box.text_frame.margin_top = Inches(0.1)
add_para(tf_n, "", size=4)
add_para(tf_n, "Veld-niveau (binnen één GE/relatie)", size=15, bold=True, color=ZWART)
add_para(tf_n, "  volledig_adres = straat + \" \" + huisnr + \", \" + postcode", size=12, color=DONKERGRIJS, font_name="Consolas")
add_para(tf_n, "", size=6)
add_para(tf_n, "Representatie-niveau (over child GE's heen)", size=15, bold=True, color=ZWART)
add_para(tf_n, "  weergavenaam = Naam.roepnaam ?: Naam.voorletters", size=12, color=DONKERGRIJS, font_name="Consolas")
add_para(tf_n, "                 + \" \" + Naam.achternaam", size=12, color=DONKERGRIJS, font_name="Consolas")

# Expressietalen
lang = add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.5), WIT, PAARS, Pt(2))
tf_l = set_text(lang, "Expressietalen", size=18, bold=True, color=PAARS)
lang.text_frame.margin_left = Inches(0.25)
lang.text_frame.margin_top = Inches(0.1)
add_para(tf_l, "", size=4)
add_para(tf_l, "▸  CEL (Common Expression Language) — default", size=14, color=ZWART)
add_para(tf_l, "   Google's type-safe, Go-native (K8s, Firebase)", size=12, color=DONKERGRIJS)
add_para(tf_l, "▸  Expr — lightweight Go-based", size=14, color=ZWART)
add_para(tf_l, "▸  JsonLogic — JSON-serialiseerbaar, cross-platform", size=14, color=ZWART)
add_para(tf_l, "▸  Pseudo — documentatie, niet geëvalueerd", size=14, color=ZWART)

# Weergavevelden
wv = add_shape(slide, Inches(0.4), Inches(4.3), Inches(6.2), Inches(2.8), WIT, GROEN, Pt(2))
tf_wv = set_text(wv, "Weergavevelden (isWeergaveVeld)", size=16, bold=True, color=GROEN)
wv.text_frame.margin_left = Inches(0.25)
wv.text_frame.margin_top = Inches(0.1)
add_para(tf_wv, "", size=3)
add_para(tf_wv, "Velden met isWeergaveVeld: true verschijnen", size=14, color=ZWART)
add_para(tf_wv, "op visuele kaarten in Index en Tijdlijn:", size=14, color=ZWART)
add_para(tf_wv, "", size=4)
add_para(tf_wv, "▸  Entiteitskaart: \"J. de Vries\"", size=14, color=DONKERGRIJS)
add_para(tf_wv, "▸  GE-box: samenvatting van de inhoud", size=14, color=DONKERGRIJS)
add_para(tf_wv, "▸  Meerdere weergavevelden: \"Jansen | EigenNaam\"", size=14, color=DONKERGRIJS)
add_para(tf_wv, "", size=4)
add_para(tf_wv, "Frontend evalueert CEL via celEvaluator.js", size=13, color=GRIJS)
add_para(tf_wv, "(hub-aware: kent de pad-notatie GE.veld)", size=13, color=GRIJS)

# UML-conventie
uml = add_shape(slide, Inches(6.8), Inches(4.3), Inches(6.2), Inches(2.8), WIT, ORANJE, Pt(2))
tf_uml = set_text(uml, "UML-conventie voor afgeleide velden", size=16, bold=True, color=ORANJE)
uml.text_frame.margin_left = Inches(0.25)
uml.text_frame.margin_top = Inches(0.1)
add_para(tf_uml, "", size=4)
add_para(tf_uml, "In UML: /prefix (schuine streep = afgeleid)", size=14, color=ZWART)
add_para(tf_uml, "", size=4)
add_para(tf_uml, "  /weergavenaam : string", size=15, color=ORANJE, font_name="Consolas")
add_para(tf_uml, "", size=4)
add_para(tf_uml, "▸  Oranje / prefix in de editor canvas", size=14, color=DONKERGRIJS)
add_para(tf_uml, "▸  Italic weergave op diagram-nodes", size=14, color=DONKERGRIJS)
add_para(tf_uml, "▸  Edit-panel met oranje accent", size=14, color=DONKERGRIJS)
add_para(tf_uml, "", size=4)
add_para(tf_uml, "Bewaard in V3 JSON → roundtrip behouden", size=13, color=GRIJS)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 23: REFERENTIELIJSTEN
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Referentielijsten", "Landenlijst, EU-Lidstaten — als subtypes van bestaande metatypes")

# Kernidee
ki = add_shape(slide, Inches(0.4), Inches(1.5), Inches(12.6), Inches(1.0), GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Pt(2))
set_text(ki, "💡  Referentielijsten zijn geen apart mechanisme — ze hergebruiken hub+data, formele/materiële tijd en de registratie-API.", 
         size=16, bold=True, color=RGBColor(0x85, 0x4D, 0x0E), align=PP_ALIGN.CENTER)
ki.text_frame.margin_top = Inches(0.1)

# Drie subtypes
st = add_shape(slide, Inches(0.4), Inches(2.8), Inches(6.2), Inches(4.2), WIT, MIDDENBLAUW, Pt(2))
tf_st = set_text(st, "Drie subtypes van bestaande metatypes", size=17, bold=True, color=MIDDENBLAUW)
st.text_frame.margin_left = Inches(0.25)
st.text_frame.margin_top = Inches(0.1)
add_para(tf_st, "", size=4)

subtypes = [
    ("referentielijst", "subtype van Entiteit", "Eén Go struct; elke instantie (bijv.\n\"Landenlijst\") is een record", MIDDENBLAUW),
    ("referentielijst_item", "subtype van Entiteit", "Bijv. Land — gewone entiteit met\nvrij modelleerbare GE's", GROEN),
    ("referentielijst_items", "subtype van Relatie", "Koppeltabel: FK naar lijst + FK naar item\nBijv. LandenlijstLand", ORANJE),
]
for name, parent, desc, clr in subtypes:
    add_para(tf_st, f"▸  {name}", size=15, bold=True, color=clr)
    add_para(tf_st, f"   {parent}", size=12, color=GRIJS)
    for line in desc.split('\n'):
        add_para(tf_st, f"   {line}", size=12, color=DONKERGRIJS)
    add_para(tf_st, "", size=4)

add_para(tf_st, "Onderscheid zit puur in metamodel-metadata,", size=13, bold=True, color=MIDDENBLAUW)
add_para(tf_st, "niet in runtime-gedrag!", size=13, bold=True, color=MIDDENBLAUW)

# Mermaid rechts
ref_mermaid = [
    "classDiagram",
    "  direction LR",
    "",
    "  class Referentielijst {",
    "    <<entiteit>>",
    "    +id : integer",
    "  }",
    "  class ReferentielijstNaam {",
    "    <<GE hub>>",
    "    naam : string",
    "  }",
    "  class Land {",
    "    <<referentielijst_item>>",
    "    +id : integer",
    "  }",
    "  class LandenlijstLand {",
    "    <<referentielijst_items>>",
    "    +landenlijst_id : FK",
    "    +land_id : FK",
    "  }",
    "  class Landcode {",
    "    <<GE hub>>",
    "    code : string",
    "  }",
    "",
    '  Referentielijst *-- LandenlijstLand',
    '  LandenlijstLand *-- Land',
    '  Land *-- Landcode',
    '  ReferentielijstNaam --* Referentielijst',
]
add_code_box(slide, Inches(6.8), Inches(2.8), Inches(6.2), Inches(4.2),
             ref_mermaid, title="Mermaid: referentielijsten voorbeeld", title_color=LICHTBLAUW)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 24: TECH STACK
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Tech stack & architectuur", "Van MetaRegistry tot browser")

layers = [
    ("Frontend", "React + Vite", "Dynamische UI, SVG-visualisatie", GROEN_LICHT, GROEN, Inches(1.5)),
    ("Schema-API", "/api/viz/schema", "Metadata + veldtypes als JSON (OAS 3.1)", BLAUW_LICHT, MIDDENBLAUW, Inches(2.4)),
    ("Handlers", "Generieke Go handlers", "MakeGet…ByMetaHandler(meta)", PAARS_LICHT, PAARS, Inches(3.3)),
    ("MetaRegistry", "Go map[string]TypeMeta", "Single source of truth (gegenereerd)", GEEL_LICHT, RGBColor(0xCA, 0x8A, 0x04), Inches(4.2)),
    ("ORM", "Bun (Go)", "DB-mapping via struct tags", ORANJE_LICHT, ORANJE, Inches(5.1)),
    ("Database", "PostgreSQL", "Tabellen, FK's, relatieve autoincrement", RGBColor(0xE0, 0xE7, 0xFF), MIDDENBLAUW, Inches(6.0)),
]

for name, tech, desc, bg_clr, line_clr, y in layers:
    box = add_shape(slide, Inches(1.5), y, Inches(10.5), Inches(0.72), bg_clr, line_clr, Pt(2))
    tf = set_text(box, name, size=16, bold=True, color=line_clr)
    box.text_frame.margin_left = Inches(0.3)
    box.text_frame.margin_top = Inches(0.05)
    
    tb_tech = add_textbox(slide, Inches(4.5), y + Inches(0.03), Inches(3.5), Inches(0.6))
    set_text(tb_tech, tech, size=14, bold=True, color=ZWART)
    
    tb_desc = add_textbox(slide, Inches(8.0), y + Inches(0.03), Inches(4.0), Inches(0.6))
    set_text(tb_desc, desc, size=13, color=GRIJS)

for i in range(len(layers) - 1):
    y1 = layers[i][5] + Inches(0.72)
    y2 = layers[i+1][5]
    add_arrow(slide, Inches(6.75), y1, Inches(6.75), y2, GRIJS, Pt(2))

# Sidebar: extra tools
side = add_shape(slide, Inches(0.2), Inches(1.5), Inches(1.2), Inches(5.3), WIT, GRIJS)
tf_s = set_text(side, "Tools", size=13, bold=True, color=DONKERGRIJS, align=PP_ALIGN.CENTER)
side.text_frame.margin_top = Inches(0.05)
add_para(tf_s, "", size=3)
add_para(tf_s, "Gin HTTP", size=11, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf_s, "Codegen", size=11, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf_s, "UML Editor", size=11, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf_s, "Mermaid", size=11, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf_s, "PlantUML", size=11, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf_s, "XMI/EA", size=11, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf_s, "CEL", size=11, color=GRIJS, align=PP_ALIGN.CENTER)
add_para(tf_s, "GraphQL*", size=11, color=GRIJS, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 25: WERKWIJZE BIJ MODELWIJZIGINGEN
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, ACHTERGROND)
add_title_bar(slide, "Werkwijze bij modelwijzigingen", "Twee paden: handmatig of via codegen roundtrip")

# Pad A: Codegen (aanbevolen)
path_a = add_shape(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.5), GROEN_LICHT, GROEN, Pt(2))
tf_a = set_text(path_a, "Pad A: via UML Editor + Codegen", size=18, bold=True, color=GROEN)
path_a.text_frame.margin_left = Inches(0.25)
path_a.text_frame.margin_top = Inches(0.1)
add_para(tf_a, "(aanbevolen)", size=14, color=GRIJS)
add_para(tf_a, "", size=6)

steps_a = [
    ("1", "Model aanpassen in UML Editor"),
    ("2", "Exporteer V3 JSON (of save naar API)"),
    ("3", "Run codegen: go run ./cmd/codegen --input model.json"),
    ("4", "7 bestanden worden gegenereerd"),
    ("5", "DB-tabel aanmaken (createtables.go)"),
    ("6", "go build && go test"),
    ("7", "Frontend past zich automatisch aan!"),
]
for num, text in steps_a:
    clr = GROEN if num == "7" else ZWART
    add_para(tf_a, f"  {num}.  {text}", size=14, color=clr, bold=(num == "7"))
    add_para(tf_a, "", size=2)

# Pad B: Handmatig
path_b = add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.5), WIT, MIDDENBLAUW, Pt(2))
tf_b = set_text(path_b, "Pad B: handmatig", size=18, bold=True, color=MIDDENBLAUW)
path_b.text_frame.margin_left = Inches(0.25)
path_b.text_frame.margin_top = Inches(0.1)
add_para(tf_b, "(voor edge cases / fijntuning)", size=14, color=GRIJS)
add_para(tf_b, "", size=6)

steps_b = [
    ("1", "Struct definiëren (modellen_*.go)", "Go struct met JSON/Bun/schema tags"),
    ("2", "Interface-methoden implementeren", "GetID, Metatype, ClearID, Opvoer/Afvoer"),
    ("3", "MetaRegistry-entry toevoegen", "TypeMeta met alle metadata"),
    ("4", "OnderliggendeGE's aanpassen", "Bij entiteiten: Rolnaam + JSONRolnaam"),
    ("5", "DB-tabel aanmaken", "CREATE TABLE in createtables.go"),
    ("6", "Klaar!", "Routes en handlers automatisch gegenereerd"),
]
for num, title, desc in steps_b:
    clr = GROEN if num == "6" else MIDDENBLAUW
    add_para(tf_b, f"  {num}.  {title}", size=14, bold=True, color=clr)
    add_para(tf_b, f"       {desc}", size=12, color=DONKERGRIJS)
    add_para(tf_b, "", size=2)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 26: SAMENVATTING
# ═══════════════════════════════════════════════════════════════════════
slide = new_slide()
set_slide_bg(slide, DONKERBLAUW)

add_rect(slide, Inches(0), Inches(1.0), Inches(13.333), Inches(0.04), LICHTBLAUW)

tb_title = add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.7))
set_text(tb_title, "Samenvatting", size=36, bold=True, color=WIT, align=PP_ALIGN.CENTER, font_name="Calibri Light")

col1 = add_shape(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(5.6), RGBColor(0x25, 0x4E, 0x7A))
tf1 = set_text(col1, "Principes", size=22, bold=True, color=ZACHTBLAUW)
col1.text_frame.margin_left = Inches(0.25)
col1.text_frame.margin_top = Inches(0.1)
add_para(tf1, "", size=4)
add_para(tf1, "✓  Twee tijdsdimensies", size=15, color=WIT)
add_para(tf1, "     formeel + materieel", size=13, color=ZACHTBLAUW)
add_para(tf1, "✓  Opvoer/afvoer altijd afgeleid", size=15, color=WIT)
add_para(tf1, "✓  Drie metatypes: ENT → GE → REL", size=15, color=WIT)
add_para(tf1, "✓  Hub + Data pattern", size=15, color=WIT)
add_para(tf1, "     Apart corrigeerbaar", size=13, color=ZACHTBLAUW)
add_para(tf1, "✓  Tijdreizen over 2 dimensies", size=15, color=WIT)

col2 = add_shape(slide, Inches(4.7), Inches(1.3), Inches(4.0), Inches(5.6), RGBColor(0x25, 0x4E, 0x7A))
tf2 = set_text(col2, "Tooling", size=22, bold=True, color=ZACHTBLAUW)
col2.text_frame.margin_left = Inches(0.25)
col2.text_frame.margin_top = Inches(0.1)
add_para(tf2, "", size=4)
add_para(tf2, "✓  UML Editor met roundtrip", size=15, color=WIT)
add_para(tf2, "     V3 JSON als uitwisselformaat", size=13, color=ZACHTBLAUW)
add_para(tf2, "✓  Code generatie (7 bestanden)", size=15, color=WIT)
add_para(tf2, "✓  Import/Export:", size=15, color=WIT)
add_para(tf2, "     Mermaid, PlantUML, XMI/EA", size=13, color=ZACHTBLAUW)
add_para(tf2, "✓  Validatie-framework", size=15, color=WIT)
add_para(tf2, "     BSN, NLPostcode, IBAN", size=13, color=ZACHTBLAUW)
add_para(tf2, "✓  Afgeleide velden (CEL)", size=15, color=WIT)
add_para(tf2, "✓  Referentielijsten", size=15, color=WIT)

col3 = add_shape(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.6), RGBColor(0x25, 0x4E, 0x7A))
tf3 = set_text(col3, "Implementatie", size=22, bold=True, color=ZACHTBLAUW)
col3.text_frame.margin_left = Inches(0.25)
col3.text_frame.margin_top = Inches(0.1)
add_para(tf3, "", size=4)
add_para(tf3, "✓  MetaRegistry = single source", size=15, color=WIT)
add_para(tf3, "     Dynamische routes, handlers", size=13, color=ZACHTBLAUW)
add_para(tf3, "✓  Schema-driven frontend", size=15, color=WIT)
add_para(tf3, "     Geen hardcoded structuren", size=13, color=ZACHTBLAUW)
add_para(tf3, "✓  PostgreSQL + relatieve AI", size=15, color=WIT)
add_para(tf3, "     Samengestelde PK's/scope", size=13, color=ZACHTBLAUW)
add_para(tf3, "✓  Go + Gin + Bun", size=15, color=WIT)
add_para(tf3, "✓  React + Vite frontend", size=15, color=WIT)
add_para(tf3, "✓  OAS 3.1 type-systeem", size=15, color=WIT)


# ── Opslaan ──────────────────────────────────────────────────────────
output_path = "bitemporeel_register_presentatie_v2.pptx"
prs.save(output_path)
print(f"Presentatie opgeslagen als: {output_path}")
print(f"Aantal slides: {len(prs.slides)}")
