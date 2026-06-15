#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt, Inches

pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'

prs = Presentation(pptx_file)
print(f'Totaal slides: {len(prs.slides)}\n')

# Analyseer slide 3 en 4 als template
print('=== TEMPLATE ANALYSE ===')
slide3 = prs.slides[2]
print(f'Slide 3: {len(slide3.shapes)} shapes')

title_template = None
bullet_template = None

for shape in slide3.shapes:
    if not hasattr(shape, 'text_frame'):
        continue
    
    text = shape.text.strip()
    if not text:
        continue
    
    font_size = None
    if len(shape.text_frame.paragraphs) > 0:
        para = shape.text_frame.paragraphs[0]
        if len(para.runs) > 0:
            font_size = para.runs[0].font.size
    
    if font_size == 361950:
        title_template = shape
        print(f'  Titel template: "{text}" @ ({shape.left}, {shape.top})')
    elif font_size == 161925 and shape.text.startswith('Propriëtaire'):
        bullet_template = shape
        print(f'  Bullet template @ ({shape.left}, {shape.top})')
        print(f'    Size: {shape.width} x {shape.height}')

if not title_template or not bullet_template:
    print('ERROR: Could not find templates!')
    exit(1)

print()

# Nu transformeer slides 5-17 (index 4-16)
for slide_idx in range(4, 17):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    print(f'Processing slide {slide_num}...', end=' ')
    
    # Verzamel alle content
    title_text = None
    bullet_texts = []
    
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        
        text = shape.text.strip()
        if not text or text == '→':
            continue
        
        font_size = None
        if len(shape.text_frame.paragraphs) > 0:
            para = shape.text_frame.paragraphs[0]
            if len(para.runs) > 0:
                font_size = para.runs[0].font.size
        
        # Titel = grootste font (360000+)
        if font_size and font_size > 200000 and not title_text:
            title_text = text
        # Bullet items
        elif text and font_size and font_size < 200000:
            bullet_texts.append(text)
    
    if not title_text:
        print('SKIP (no title)')
        continue
    
    # Verwijder alle TEXT_BOX shapes (behoudt PICTURE/backgrounds)
    shapes_to_remove = []
    for shape in list(slide.shapes):
        if hasattr(shape, 'text_frame'):
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    
    # Voeg titel toe
    title_box = slide.shapes.add_textbox(
        title_template.left, 
        title_template.top, 
        title_template.width, 
        title_template.height
    )
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    
    # Voeg bullets toe (één per tekstbox, template-style)
    y_offset = 0
    for bullet_text in bullet_texts:
        bullet_top = bullet_template.top + (y_offset * Inches(0.35))
        
        bullet_box = slide.shapes.add_textbox(
            bullet_template.left,
            bullet_top,
            bullet_template.width,
            bullet_template.height
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = bullet_text
        p.font.size = Pt(12)  # 161925 EMU
        p.font.bold = True
        
        y_offset += 1
    
    print(f'✓ ({len(bullet_texts)} bullets)')

prs.save(pptx_file)
print(f'\n✓ Opgeslagen!')
