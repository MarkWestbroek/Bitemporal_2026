#!/usr/bin/env python3
from pptx import Presentation

pptx_file = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'
prs = Presentation(pptx_file)

slide3 = prs.slides[2]
print('=== SLIDE 3 SHAPES ===')

for i, shape in enumerate(slide3.shapes):
    if hasattr(shape, 'text_frame'):
        text = shape.text[:50].replace('\n', ' ')
        length = len(shape.text)
        print(f'{i}: TEXT "{text}" ({length} chars)')
        
        if len(shape.text_frame.paragraphs) > 0:
            p = shape.text_frame.paragraphs[0]
            if len(p.runs) > 0:
                print(f'   Font: {p.runs[0].font.size}')
    else:
        print(f'{i}: {shape.shape_type}')
