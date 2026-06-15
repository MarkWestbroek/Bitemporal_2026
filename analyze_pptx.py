#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt
import os

# Zoek het bestand
search_paths = [
    'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx',
    'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap.pptx'
]

pptx_file = None
for path in search_paths:
    if os.path.exists(path):
        pptx_file = path
        break

if pptx_file:
    print(f'Gevonden: {pptx_file}')
    prs = Presentation(pptx_file)
    print(f'Totaal slides: {len(prs.slides)}\n')
    
    # Analyseer slide 3 (index 2)
    print('=== SLIDE 3 ===')
    slide3 = prs.slides[2]
    for i, shape in enumerate(slide3.shapes):
        if hasattr(shape, 'text_frame'):
            txt = shape.text[:40] if len(shape.text) > 40 else shape.text
            print(f'Shape {i}: "{txt}"')
            if len(shape.text_frame.paragraphs) > 0:
                para = shape.text_frame.paragraphs[0]
                if len(para.runs) > 0:
                    run = para.runs[0]
                    print(f'  Font size: {run.font.size}, Bold: {run.font.bold}')
                    print(f'  Left: {shape.left}, Top: {shape.top}, Width: {shape.width}, Height: {shape.height}')
    
    # Analyseer slide 5 (index 4)
    print('\n=== SLIDE 5 ===')
    slide5 = prs.slides[4]
    for i, shape in enumerate(slide5.shapes):
        if hasattr(shape, 'text_frame'):
            txt = shape.text[:40] if len(shape.text) > 40 else shape.text
            print(f'Shape {i}: "{txt}"')
            if len(shape.text_frame.paragraphs) > 0:
                para = shape.text_frame.paragraphs[0]
                if len(para.runs) > 0:
                    run = para.runs[0]
                    print(f'  Font size: {run.font.size}, Bold: {run.font.bold}')
                    print(f'  Left: {shape.left}, Top: {shape.top}, Width: {shape.width}, Height: {shape.height}')
else:
    print('Bestand niet gevonden')
