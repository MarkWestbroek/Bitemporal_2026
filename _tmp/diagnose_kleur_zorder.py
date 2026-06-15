#!/usr/bin/env python3
"""Diagnose: tekstkleur en z-order in slide 3 (template) vs slide 5 (getransformeerd)."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

# Gebruik de getransformeerde versie
SRC = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\Deepseek\\FOST & apidays Amsterdam 2026 v0.5.pptx'

prs = Presentation(SRC)

for slide_idx in [2, 4]:  # slide 3 (idx 2) en slide 5 (idx 4)
    slide = prs.slides[slide_idx]
    print(f'=== SLIDE {slide_idx+1} ({len(slide.shapes)} shapes) ===')
    
    for i, shape in enumerate(slide.shapes):
        stype = str(shape.shape_type)
        has_text = hasattr(shape, 'text_frame')
        
        # Positie en afmeting
        L = shape.left if shape.left else 0
        T = shape.top if shape.top else 0
        W = shape.width if shape.width else 0
        H = shape.height if shape.height else 0
        
        info = f'  [{i}] {stype} L={L} T={T} W={W} H={H}'
        
        if has_text and shape.text_frame:
            txt = shape.text_frame.text[:60].replace('\n',' | ')
            
            # Font info: size, bold, color
            font_info = ''
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fs = run.font.size
                    fb = run.font.bold
                    fc = run.font.color
                    if fc and fc.rgb:
                        font_info = f' sz={fs/12700:.0f}pt' if fs else ''
                        font_info += ' BOLD' if fb else ''
                        font_info += f' color=#{fc.rgb}'
                    elif fs:
                        font_info = f' sz={fs/12700:.0f}pt' if fs else ''
                        font_info += ' BOLD' if fb else ''
                    break
                if font_info:
                    break
            
            # Check paragraph-level font
            if not font_info:
                for para in shape.text_frame.paragraphs:
                    pf = para.font
                    if pf.size or pf.bold or (pf.color and pf.color.rgb):
                        font_info = f' sz={pf.size/12700:.0f}pt' if pf.size else ''
                        font_info += ' BOLD' if pf.bold else ''
                        if pf.color and pf.color.rgb:
                            font_info += f' color=#{pf.color.rgb}'
                        break
            
            info += font_info
            info += f'\n       "{txt}"'
        
        print(info)
    print()

# Ook checken: slide 5 background fill / image
print('=== SLIDE 5 BACKGROUND CHECK ===')
slide5 = prs.slides[4]
bg = slide5.background
print(f'  Background fill: {bg.fill.type if bg.fill else "none"}')

# Check picture shape details
for i, shape in enumerate(slide5.shapes):
    if shape.shape_type == 13:  # PICTURE
        print(f'  Picture [{i}]: L={shape.left} T={shape.top} W={shape.width} H={shape.height}')
