#!/usr/bin/env python3
"""
Fix voor NIEUW.pptx: tekst in slides 5-17 heeft de verkeerde kleur (zwart).
Kopieer de kleur uit het template (slide 3) naar alle tekst in slides 5-17.
En zorg dat tekst-shapes bovenaan in de z-order staan.
"""
from pptx import Presentation
from lxml import etree
from copy import deepcopy

NIEUW = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p - NIEUW.pptx'
ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

prs = Presentation(NIEUW)

# Haal de exacte kleur-XML op uit slide 3
slide3 = prs.slides[2]
title_color_xml = None
content_color_xml = None

for shape in slide3.shapes:
    if not hasattr(shape, 'text_frame') or not shape.text.strip():
        continue
    paras = shape.text_frame.paragraphs
    if not paras or not paras[0].runs:
        continue
    run = paras[0].runs[0]
    rPr = run._r.find(f'{{{ns}}}rPr')
    if rPr is None:
        continue
    fill = rPr.find(f'{{{ns}}}solidFill')
    if fill is None:
        continue
    font_size = run.font.size
    if font_size and font_size > 250000:
        title_color_xml = deepcopy(fill)
        print(f'Titel kleur: {etree.tostring(title_color_xml, encoding="unicode")}')
    else:
        content_color_xml = deepcopy(fill)
        print(f'Content kleur: {etree.tostring(content_color_xml, encoding="unicode")}')

if not title_color_xml or not content_color_xml:
    print('ERROR: kleuren niet gevonden in slide 3')
    exit(1)

def apply_color_to_shape(shape, color_xml):
    """Zet de kleur op alle runs in een shape."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(f'{{{ns}}}rPr')
            if rPr is None:
                rPr = etree.SubElement(run._r, f'{{{ns}}}rPr')
                rPr.set('lang', 'nl-NL')
                rPr.set('b', '1')
                run._r.insert(0, rPr)
            # Verwijder bestaande solidFill
            for old_fill in rPr.findall(f'{{{ns}}}solidFill'):
                rPr.remove(old_fill)
            # Voeg nieuwe kleur in
            rPr.insert(0, deepcopy(color_xml))

def bring_text_to_front(slide):
    """Verplaats alle tekst-shapes naar het einde van spTree (bovenaan in z-order)."""
    spTree = slide.shapes._spTree
    text_sps = []
    for child in list(spTree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'sp':
            # Check of dit een text shape is
            txBody = child.find(f'.//{{{ns}}}txBody')
            if txBody is not None:
                t_elems = txBody.findall(f'.//{{{ns}}}t')
                has_text = any(t.text and t.text.strip() for t in t_elems)
                if has_text:
                    text_sps.append(child)
    # Verplaats naar einde
    for sp in text_sps:
        spTree.remove(sp)
        spTree.append(sp)

# Pas kleur en z-order toe op slides 5-17
for slide_idx in range(4, 17):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    fixed_shapes = 0
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame') or not shape.text.strip():
            continue
        # Bepaal of het een titel is
        font_size = None
        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
            font_size = shape.text_frame.paragraphs[0].runs[0].font.size
        
        if font_size and font_size > 250000:
            apply_color_to_shape(shape, title_color_xml)
        else:
            apply_color_to_shape(shape, content_color_xml)
        fixed_shapes += 1
    
    bring_text_to_front(slide)
    print(f'Slide {slide_num}: ✓ {fixed_shapes} shapes gekleurd + z-order gefixed')

prs.save(NIEUW)
print(f'\n✓ Opgeslagen!')
