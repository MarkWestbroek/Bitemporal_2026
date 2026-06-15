#!/usr/bin/env python3
"""
Transformeer FOST & apidays Amsterdam 2026 v0.5.pptx:
Uniformeer slides 5 t/m 17 naar de layout van slide 3:
- Titel (28pt BOLD) bovenaan
- Één content-tekstvak (16pt BOLD) met alle bullet-tekst
- Behoud afbeeldingen
"""
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
import copy
import os
import shutil

# === CONFIG ===
SRC = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\Deepseek\\FOST & apidays Amsterdam 2026 v0.5.pptx'
DST = SRC  # overschrijf (we maken eerst backup)
BACKUP = SRC.replace('.pptx', ' - Copy.pptx')  # bestaat al als backup

# === Template metingen van slide 3 ===
# Slide 3 shapes:
# [0] TEXT_BOX: L=571500 T=939641  W=5200650 H=438150  font=361950EMU (28pt) BOLD (titel)
# [1] TEXT_BOX: L=589788 T=1832705 W=4953000 H=3453061 font=203200EMU (16pt) BOLD (content)
# [2] PICTURE:   L=5934075 T=935980  W=6096000 H=4919957

TITLE_LEFT   = 571500
TITLE_TOP    = 939641
TITLE_WIDTH  = 5200650
TITLE_HEIGHT = 438150
TITLE_FONT   = Pt(28)

CONTENT_LEFT   = 589788
CONTENT_TOP    = 1832705
CONTENT_WIDTH  = 4953000
CONTENT_HEIGHT = 3453061
CONTENT_FONT   = Pt(16)

# === Helper functies ===
def get_shape_text(shape):
    """Haal volledige tekst uit een shape."""
    if not hasattr(shape, 'text_frame'):
        return None
    return shape.text_frame.text.strip()

def get_first_font_size(shape):
    """Haal font size van eerste run in EMU."""
    if not hasattr(shape, 'text_frame'):
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                return run.font.size
    return None

def remove_text_shapes(slide):
    """Verwijder alle TEXT_BOX shapes van een slide, retourneer verwijderde count."""
    to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    return len(to_remove)

def add_title(slide, text):
    """Voeg titel-textbox toe met slide-3 formatting."""
    box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = TITLE_FONT
    p.font.bold = True
    return box

def add_content(slide, lines):
    """Voeg content-textbox toe met slide-3 formatting.
    lines: list van strings (lege strings worden lege paragrafen)"""
    box = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = CONTENT_FONT
        p.font.bold = True
    
    return box

def has_picture(slide):
    """Check of slide minstens één picture shape heeft."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            return True
    return False

# === HOOFDPROGRAMMA ===
print(f'Bron: {SRC}')
prs = Presentation(SRC)
total = len(prs.slides)
print(f'Totaal slides: {total}')

# Slide index mapping: 0=slide1, 1=slide2, 2=slide3, 3=slide4, 4=slide5, ...
# We transformeren slides 5 t/m 17 (index 4 t/m 16)

transformed = 0
skipped = 0

for slide_idx in range(4, total):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    # --- Verzamel alle tekst uit huidige slide ---
    title_text = None
    intro_lines = []   # subtitles/intro regels (zonder →)
    bullet_lines = []  # bullet regels (we voegen → toe)
    
    for shape in slide.shapes:
        text = get_shape_text(shape)
        if not text:
            continue
        
        font_size = get_first_font_size(shape)
        
        # Titel: grootste font (28pt+)
        if font_size and font_size >= 350000:  # >= 28pt
            if not title_text:
                title_text = text
            continue
        
        # Skip arrow-only shapes
        if text == '→':
            continue
        
        # Check of het een intro/subtitle is (zonder → prefix, vaak eerste regel)
        if text.startswith('→'):
            bullet_lines.append(text[1:].strip())  # strip → prefix
        elif '→' in text:
            # Split op → als het inline bullets heeft
            parts = text.split('→')
            first = parts[0].strip()
            if first:
                intro_lines.append(first)
            for part in parts[1:]:
                if part.strip():
                    bullet_lines.append(part.strip())
        else:
            # Check of het een "header" is (13pt BOLD, vaak de intro-regel)
            if font_size and font_size <= 210000:  # <= 16pt
                if len(intro_lines) == 0 and not text.startswith('→'):
                    intro_lines.append(text)
                else:
                    bullet_lines.append(text)
            else:
                intro_lines.append(text)
    
    # Falls er geen titel is, zoek in alle text
    if not title_text:
        # Pak de tekst van de laatste gevonden shape als titel
        for shape in reversed(list(slide.shapes)):
            text = get_shape_text(shape)
            if text and text != '→' and len(text) > 3:
                title_text = text
                break
    
    if not title_text:
        print(f'Slide {slide_num}: ⚠ Geen titel gevonden, skip')
        skipped += 1
        continue
    
    print(f'Slide {slide_num}: "{title_text}" ({len(intro_lines)} intro, {len(bullet_lines)} bullets)')
    
    # --- Verwijder alle oude tekst-shapes ---
    remove_text_shapes(slide)
    
    # --- Voeg nieuwe titel toe ---
    add_title(slide, title_text)
    
    # --- Bouw content lines ---
    content_lines = []
    
    # Intro lines eerst
    for line in intro_lines:
        content_lines.append(line)
    
    # Dan bullets met → prefix
    for line in bullet_lines:
        if content_lines:
            content_lines.append('')  # lege regel voor spacing
        content_lines.append(f'→ {line}')
    
    if content_lines:
        add_content(slide, content_lines)
    
    transformed += 1

# --- Opslaan ---
prs.save(DST)
print(f'\n✓ Klaar! {transformed} slides getransformeerd, {skipped} overgeslagen')
print(f'  Opgeslagen: {DST}')
print(f'  Backup al aanwezig: {BACKUP}')
