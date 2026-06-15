#!/usr/bin/env python3
"""
Uniformeer slides 5-17 naar het formaat van slides 3 & 4.
Strategie: GEEN nieuwe shapes aanmaken, maar bestaande shapes aanpassen.
- Intro-textbox krijgt alle bullet-content als paragrafen (slide-3-stijl)
- Losse bullet-textboxen en pijl-textboxen worden verwijderd
"""
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
from lxml import etree

SRC = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'
DST = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p - NIEUW.pptx'

prs = Presentation(SRC)

# Haal de XML-template paragrafen op uit slide 3 (index 2)
# para[0] = intro (203200pt), para[1] = leeg, para[2] = "→ bullet" patroon
slide3_content = None
for shape in prs.slides[2].shapes:
    if hasattr(shape, 'text_frame') and len(shape.text_frame.paragraphs) >= 3:
        slide3_content = shape
        break

if not slide3_content:
    print('ERROR: Template content shape niet gevonden in slide 3')
    exit(1)

print('Template paragrafen uit slide 3:')
for i, p in enumerate(slide3_content.text_frame.paragraphs[:4]):
    print(f'  para[{i}]: {len(p.runs)} runs, text="{p.text[:40]}"')

print()

def get_para_xml_templates(content_shape):
    """Haal XML-template-elementen op voor intro-para en bullet-para."""
    paras = content_shape.text_frame.paragraphs
    # para 0 = intro (24pt), para 1 = leeg, para 2 = bullet
    intro_para_xml = deepcopy(paras[0]._p)
    empty_para_xml = deepcopy(paras[1]._p)
    bullet_para_xml = deepcopy(paras[2]._p)
    return intro_para_xml, empty_para_xml, bullet_para_xml

intro_tmpl, empty_tmpl, bullet_tmpl = get_para_xml_templates(slide3_content)

def set_para_text(para_xml, text):
    """Zet alle run-teksten samen tot één tekst, behoud opmaak van eerste run."""
    nsmap = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    # Verwijder alle <a:r> runs
    for r in para_xml.findall(f'{{{nsmap}}}r'):
        para_xml.remove(r)
    # Maak één nieuwe run aan met de tekst
    # Kopieer rPr van origineel als die er is
    first_r = None
    r_elem = etree.SubElement(para_xml, f'{{{nsmap}}}r')
    # Voeg rPr toe (bold, etc.) - leeg = inherit
    rPr = etree.SubElement(r_elem, f'{{{nsmap}}}rPr')
    rPr.set('lang', 'nl-NL')
    rPr.set('b', '1')
    # Voeg tekst toe
    t_elem = etree.SubElement(r_elem, f'{{{nsmap}}}t')
    t_elem.text = text
    return para_xml

# Transformeer slides 5-17
for slide_idx in range(4, 17):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1

    # Identificeer alle text shapes
    title_shape = None
    intro_shape = None  # eerste niet-pijl bullet (wordt de content-host)
    bullet_shapes = []  # overige bullets
    arrow_shapes = []   # pijl-shapes (→)

    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        text = shape.text.strip()
        if not text:
            continue

        # Font size van eerste run
        font_size = None
        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
            font_size = shape.text_frame.paragraphs[0].runs[0].font.size

        if text == '→':
            arrow_shapes.append(shape)
        elif font_size and font_size > 250000:
            title_shape = shape
        elif intro_shape is None:
            intro_shape = shape  # eerste content = intro
        else:
            bullet_shapes.append(shape)

    if not title_shape or not intro_shape:
        print(f'Slide {slide_num}: SKIP (title={title_shape is not None}, intro={intro_shape is not None})')
        continue

    intro_text = intro_shape.text.strip()
    bullet_texts = [s.text.strip() for s in bullet_shapes]

    print(f'Slide {slide_num}: intro="{intro_text[:40]}", {len(bullet_texts)} bullets')

    # Herbouw de intro_shape met alle content als paragrafen
    tf = intro_shape.text_frame
    txBody = tf._txBody
    nsmap = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Verwijder alle bestaande <a:p> paragrafen
    for p_elem in txBody.findall(f'{{{nsmap}}}p'):
        txBody.remove(p_elem)

    # Voeg intro-paragraaf toe
    intro_p = deepcopy(intro_tmpl)
    set_para_text(intro_p, intro_text)
    txBody.append(intro_p)

    # Voeg per bullet: lege regel + bullet-paragraaf
    for bt in bullet_texts:
        txBody.append(deepcopy(empty_tmpl))
        bp = deepcopy(bullet_tmpl)
        set_para_text(bp, f'→ {bt}')
        txBody.append(bp)

    # Verwijder losse bullet en pijl shapes
    for shape in bullet_shapes + arrow_shapes:
        shape.element.getparent().remove(shape.element)

prs.save(DST)
print(f'\n✓ Opgeslagen als NIEUW bestand: {DST}')
