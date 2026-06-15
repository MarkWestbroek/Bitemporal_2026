#!/usr/bin/env python3
"""
Transformeer FOST & apidays Amsterdam 2026 v0.5.pptx:
Uniformeer slides 5 t/m 17 naar de layout van slide 3:
- Titel (28pt BOLD) bovenaan
- Eén content-tekstvak (16pt BOLD) met alle bullet-tekst
- Behoud afbeeldingen

Verbeterde v2: slimme titel-detectie via positie + font-size.
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

# === CONFIG ===
SRC = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\Deepseek\\FOST & apidays Amsterdam 2026 v0.5 - Copy.pptx'
DST = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\Deepseek\\FOST & apidays Amsterdam 2026 v0.5.pptx'

# === Template metingen van slide 3 ===
TITLE_LEFT   = 571500
TITLE_TOP    = 939641
TITLE_WIDTH  = 5200650
TITLE_HEIGHT = 438150
TITLE_FONT   = Pt(28)
TITLE_COLOR  = RGBColor(0xF5, 0xF5, 0xF5)  # bijna wit (matcht slide 3)

CONTENT_LEFT   = 589788
CONTENT_TOP    = 1832705
CONTENT_WIDTH  = 4953000
CONTENT_HEIGHT = 3453061
CONTENT_FONT   = Pt(16)
CONTENT_COLOR  = RGBColor(0xDA, 0xFF, 0xDE)  # licht mintgroen (matcht slide 3)

# De standaard "bottom title" positie in originele slides 5-15:
BOTTOM_TITLE_TOP = 571500  # T-waarde van de 28pt titel in slides 5+

# === Helper functies ===
def get_shape_text(shape):
    if not hasattr(shape, 'text_frame'):
        return None
    return shape.text_frame.text.strip()

def get_first_font_size(shape):
    if not hasattr(shape, 'text_frame'):
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                return run.font.size
    return None

def remove_text_shapes(slide):
    """Verwijder alle shapes met text_frame (TEXT_BOX + AUTO_SHAPE met tekst)."""
    to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape.element
        sp.getparent().remove(sp)
    return len(to_remove)

def send_pictures_to_back(slide):
    """Verplaats alle PICTURE shapes naar het begin van de z-order (achtergrond)."""
    sp_tree = slide.shapes[0].element.getparent()  # het <p:spTree> element
    pictures = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            pictures.append(shape.element)
    # Verplaats elk picture element naar positie 0 (begin)
    for i, pic_el in enumerate(pictures):
        sp_tree.remove(pic_el)
        sp_tree.insert(i, pic_el)

def add_title(slide, text):
    box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = TITLE_FONT
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    return box

def add_content(slide, lines):
    """lines: list van strings (lege strings worden lege paragrafen)"""
    box = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = CONTENT_FONT
        run.font.bold = True
        run.font.color.rgb = CONTENT_COLOR
    return box

# === HOOFDPROGRAMMA ===
print(f'Bron (backup): {SRC}')
print(f'Doel: {DST}')
prs = Presentation(SRC)
total = len(prs.slides)

transformed = 0

for slide_idx in range(4, total):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    
    # --- Stap 1: Vind de titel ---
    # Strategie: zoek eerst naar een textbox op de standaard "bottom title" positie
    # (T=571500 ± 100000, 28pt BOLD). Als die er niet is, pak de grootste font.
    
    title_text = None
    title_candidates = []  # (text, font_size, is_bold, top)
    
    for shape in slide.shapes:
        text = get_shape_text(shape)
        if not text or text == '→':
            continue
        font_size = get_first_font_size(shape)
        if not font_size:
            continue
        
        # Check bold status
        is_bold = False
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.bold:
                    is_bold = True
                break
            break
        
        top = shape.top if shape.top is not None else 0
        title_candidates.append((text, font_size, is_bold, top))
    
    # Sorteer: eerst op "is standaard titelpositie" (T ≈ 571500), dan op font-size
    def title_score(c):
        text, fs, bold, top = c
        # Skip decoratieve tekst (alleen quotes, streepjes, etc.)
        if len(text) <= 3 and all(ch in '"\u201c\u201d\u2018\u2019\u0027\u2013\u2014\u2022\u25cf \n' for ch in text):
            return -1
        score = 0
        # Prefereer standaard titelpositie (T ≈ 571500)
        if abs(top - BOTTOM_TITLE_TOP) < 200000:
            score += 500000  # hoog gewicht zodat positie belangrijker is dan font-size
        # Prefereer BOLD
        if bold:
            score += 100000
        # Groter font = hogere score
        score += fs
        return score
    
    title_candidates.sort(key=title_score, reverse=True)
    
    if title_candidates:
        title_text = title_candidates[0][0]
    else:
        print(f'Slide {slide_num}: ⚠ Geen titel gevonden, skip')
        continue
    
    # --- Stap 2: Verzamel alle overige tekst als content ---
    intro_lines = []
    bullet_lines = []
    
    for shape in slide.shapes:
        text = get_shape_text(shape)
        if not text or text == '→':
            continue
        
        font_size = get_first_font_size(shape)
        
        # Skip de titel zelf
        if text == title_text:
            continue
        
        # Verwerk tekst met → bullets
        if '→' in text:
            parts = text.split('→')
            first = parts[0].strip()
            if first and first not in intro_lines:
                intro_lines.append(first)
            for part in parts[1:]:
                stripped = part.strip()
                if stripped and stripped not in bullet_lines:
                    bullet_lines.append(stripped)
        elif font_size and font_size >= 350000:
            # Andere large-font tekst die geen titel is → als intro
            if text not in intro_lines:
                intro_lines.append(text)
        else:
            # Kleine tekst: check of het een intro-line of bullet is
            if not intro_lines:
                intro_lines.append(text)
            elif text not in bullet_lines:
                bullet_lines.append(text)
    
    print(f'Slide {slide_num}: "{title_text}" ({len(intro_lines)} intro, {len(bullet_lines)} bullets)')
    
    # --- Stap 3: Verwijder oude text shapes en bouw nieuwe op ---
    remove_text_shapes(slide)
    add_title(slide, title_text)
    
    # --- Stap 4: Bouw content ---
    content_lines = []
    for line in intro_lines:
        content_lines.append(line)
    
    for line in bullet_lines:
        if content_lines:
            content_lines.append('')
        content_lines.append(f'→ {line}')
    
    if content_lines:
        add_content(slide, content_lines)
    
    # Zet alle afbeeldingen naar de achtergrond zodat tekst altijd zichtbaar is
    send_pictures_to_back(slide)
    
    transformed += 1

prs.save(DST)
print(f'\n✓ Klaar! {transformed} slides getransformeerd')
print(f'  Opgeslagen: {DST}')
