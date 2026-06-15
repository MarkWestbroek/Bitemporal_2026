from pptx import Presentation
prs = Presentation('d:/Git/Bitemporal_2026/bitemp_register_v06/docs/presentaties/Deepseek/FOST & apidays Amsterdam 2026 v0.5.pptx')
for slide_idx in [4, 5, 15, 16]:  # slides 5,6,16,17
    s = prs.slides[slide_idx]
    print(f'=== GETRANSFORMEERDE SLIDE {slide_idx+1} ({len(s.shapes)} shapes) ===')
    for i, sh in enumerate(s.shapes):
        st = str(sh.shape_type)
        if hasattr(sh, 'text_frame') and sh.text_frame.text.strip():
            txt = sh.text_frame.text.strip()[:70]
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    c = r.font.color
                    try:
                        crgb = str(c.rgb) if c and c.rgb else 'inherited'
                    except:
                        crgb = 'no_rgb'
                    print(f'  [{i}] {st} fs={r.font.size} bold={r.font.bold} color=#{crgb}')
                    print(f'       "{txt}"')
                    break
                break
        else:
            L = sh.left if sh.left else 0
            T = sh.top if sh.top else 0
            print(f'  [{i}] {st} L={L} T={T}')
    print()
