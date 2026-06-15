#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
import os

# Find the v0.5 file - it was attached so let's try the root
search = [
    'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\Deepseek\\FOST & apidays Amsterdam 2026 v0.5.pptx',
]
pptx_file = None
for p in search:
    if os.path.exists(p):
        pptx_file = p
        break

if not pptx_file:
    print("FILE NOT FOUND in:")
    for p in search:
        print(f"  {p}")
    # Try listing
    root = 'd:\\Git\\Bitemporal_2026'
    for f in os.listdir(root):
        if 'FOST' in f or 'apidays' in f:
            print(f"  Found similar: {f}")
    exit(1)

prs = Presentation(pptx_file)
print(f'Bestand: {pptx_file}')
print(f'Totaal slides: {len(prs.slides)}')
print(f'Slide width: {prs.slide_width}, height: {prs.slide_height}')
print()

for slide_idx in range(len(prs.slides)):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    print(f'=== SLIDE {slide_num} ({len(slide.shapes)} shapes) ===')
    for i, shape in enumerate(slide.shapes):
        stype = str(shape.shape_type)
        has_text = hasattr(shape, 'text_frame')
        txt_preview = ''
        font_info = ''
        if has_text and shape.text_frame:
            txt = shape.text_frame.text.strip()
            txt_preview = txt[:100].replace('\n',' | ')
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fs = run.font.size
                    fb = run.font.bold
                    if fs:
                        font_info = f' font={fs}EMU ({fs/12700:.0f}pt)'
                        if fb: font_info += ' BOLD'
                    break
                if font_info: break
        left_emu = shape.left if shape.left is not None else 0
        top_emu = shape.top if shape.top is not None else 0
        w_emu = shape.width if shape.width is not None else 0
        h_emu = shape.height if shape.height is not None else 0
        print(f'  [{i}] {stype} L={left_emu} T={top_emu} W={w_emu} H={h_emu}{font_info}')
        if txt_preview:
            print(f'       "{txt_preview}"')
    print()
