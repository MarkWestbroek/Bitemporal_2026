#!/usr/bin/env python3
from pptx import Presentation

src = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'
prs = Presentation(src)
print(f'Totaal slides: {len(prs.slides)}\n')

for slide_idx in range(min(7, len(prs.slides))):
    slide = prs.slides[slide_idx]
    print(f'=== SLIDE {slide_idx+1} ({len(slide.shapes)} shapes) ===')
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, 'text_frame'):
            txt = shape.text
            print(f'  Shape[{i}] type={shape.shape_type} name="{shape.name}"')
            print(f'    text (first 80): "{txt[:80].replace(chr(10)," | ")}"')
            print(f'    paragraphs: {len(shape.text_frame.paragraphs)}')
            for pi, para in enumerate(shape.text_frame.paragraphs[:4]):
                runs_info = [(r.text[:30], r.font.size, r.font.bold) for r in para.runs[:2]]
                print(f'      para[{pi}]: runs={len(para.runs)} - {runs_info}')
        else:
            print(f'  Shape[{i}] type={shape.shape_type} name="{shape.name}" (geen tekst)')
    print()
