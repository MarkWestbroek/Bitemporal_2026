#!/usr/bin/env python3
from pptx import Presentation
from lxml import etree

# Check het NIEUWE bestand (al getransformeerd)
FILES = [
    ('ORIGINEEL', 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'),
    ('NIEUW', 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p - NIEUW.pptx'),
]

nsmap = 'http://schemas.openxmlformats.org/drawingml/2006/main'

for label, path in FILES:
    prs = Presentation(path)
    print(f'\n=== {label} - SLIDE 3 ===')
    slide = prs.slides[2]
    spTree = slide.shapes._spTree
    children = list(spTree)
    print(f'spTree children volgorde ({len(children)} items):')
    for i, child in enumerate(children):
        name = child.get('id','?') + ':' + child.get('{http://schemas.openxmlformats.org/drawingml/2006/main}nvSpPr', '')
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        # get name
        nvSpPr = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
        if nvSpPr is None:
            nvSpPr = child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}nvPr')
        n = nvSpPr.get('name','?') if nvSpPr is not None else '?'
        print(f'  [{i}] {tag}: {n}')
    
    print(f'\n=== {label} - SLIDE 3 tekst kleuren ===')
    for shape in prs.slides[2].shapes:
        if hasattr(shape, 'text_frame') and shape.text.strip():
            txt = shape.text[:30].replace('\n',' ')
            # Zoek solidFill recursief
            colors = shape.element.findall(f'.//{{{nsmap}}}solidFill')
            srgb = [(c.find(f'{{{nsmap}}}srgbClr'), c.find(f'{{{nsmap}}}schemeClr')) for c in colors[:1]]
            color_str = ''
            if colors:
                srgb_elem = colors[0].find(f'{{{nsmap}}}srgbClr')
                scheme_elem = colors[0].find(f'{{{nsmap}}}schemeClr')
                if srgb_elem is not None:
                    color_str = f'srgb: #{srgb_elem.get("val","?")}'
                elif scheme_elem is not None:
                    color_str = f'scheme: {scheme_elem.get("val","?")}'
            print(f'  "{txt}" → {color_str}')
