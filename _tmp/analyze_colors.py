#!/usr/bin/env python3
from pptx import Presentation
from lxml import etree

SRC = 'd:\\Git\\Bitemporal_2026\\bitemp_register_v06\\docs\\presentaties\\apidays\\FOST & apidays Amsterdam 2026 - Volledige Recap 17p.pptx'
prs = Presentation(SRC)

nsmap = 'http://schemas.openxmlformats.org/drawingml/2006/main'

for slide_idx in [2, 4]:  # slide 3 en 5
    slide = prs.slides[slide_idx]
    print(f'=== SLIDE {slide_idx+1} ===')
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, 'text_frame') and shape.text.strip():
            txt = shape.text[:30].replace('\n',' ')
            # Kleur van eerste run
            color_xml = ''
            try:
                p = shape.text_frame.paragraphs[0]
                if p.runs:
                    rPr = p.runs[0]._r.find(f'{{{nsmap}}}rPr')
                    if rPr is not None:
                        solidFill = rPr.find(f'{{{nsmap}}}solidFill')
                        if solidFill is not None:
                            color_xml = etree.tostring(solidFill, encoding='unicode')[:80]
                        else:
                            color_xml = '(geen solidFill in rPr)'
                    else:
                        color_xml = '(geen rPr)'
                else:
                    color_xml = '(geen runs)'
            except Exception as e:
                color_xml = str(e)
            print(f'  [{i}] "{txt}" → kleur: {color_xml}')
    print()
